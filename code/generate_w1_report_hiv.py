#!/usr/bin/env python3
"""
Gerador de Relatório Delphi W1 HIV/SIDA
================================
Uso:
  python generate_w1_report_hiv.py <resultados.xlsx> [dicionario.xlsx]
      [--output-dir DIR]
      [--exclude-experts ESP1,ESP2,...]
      [--exclude-file ficheiro.txt]
      [--config config.env]
  [--simple-table]
      [--simple-sections | --compact-report]

Argumentos:
    resultados.xlsx   Ficheiro com os resultados agregados (obrigatório).
                      Folha usada: 'Responses' (tall-skinny, uma linha por
                      expert_code × intervention).
    dicionario.xlsx   Ficheiro dicionário com os metadados das intervenções
                      (opcional mas recomendado).
                      Folha usada: primeira cujo nome começa por 'Catalogo'.
    --output-dir DIR  Directório de saída (opcional, padrão: ../reports/).
    --exclude-experts Lista de códigos de especialistas a excluir (separados por vírgula,
              ponto-e-vírgula ou espaço). Exemplo: --exclude-experts 001PM,001XX
    --exclude-file    Ficheiro com códigos de especialistas a excluir (um por linha,
              com suporte a comentários iniciados por '#').
    --config          Ficheiro de configuração opcional para exclusões e metadados.
              Chaves suportadas: REPORT_EXCLUDE_EXPERTS, EXCLUDE_EXPERTS,
              REPORT_EXCLUDE_FILE, EXCLUDE_FILE, TOPIC_NAME.
    --simple-sections Omite secções visuais extensas: XY plot (2g),
              "Pontuação e Metodologia", diagrama aluvial e tabela de pontuação.
    --compact-report  Alias para --simple-sections.
    --simple-table    Substitui a tabela detalhada de pontuação por uma versão
          simples (Intervenção, S_pond, Índice de eficiência, Grupo 1/2/3),
          ordenável e com ordenação inicial por S_pond.

Metadados das intervenções (código, nome, componente, URL):
    1. Dicionário externo (se fornecido) — folha Catalogo_*
    2. Folha Catalogo_* dentro do próprio ficheiro de resultados
    3. Colunas url_* na folha Submissions do ficheiro de resultados
       (nome e componente ficam em branco neste último caso)
    Se nenhuma fonte estiver disponível, o script usa os códigos brutos.

Colunas obrigatórias na folha Responses:
    expert_code, intervention, gate
Colunas opcionais:
    dup, which_dup, which_dup_other, intg, which_intg, which_intg_other,
    res, oth, oth_reason, impact, cmt, exp, modality, group

Saída:
  delphi_w1_hiv_relatorio_<timestamp>.html  no directório de trabalho actual
"""

import sys
import os
import math
import json
from collections import defaultdict, OrderedDict
from datetime import datetime
import statistics
from pathlib import Path

try:
    import pandas as pd
    import openpyxl
except ImportError:
    sys.exit("Erro: dependências em falta. Execute: pip install pandas openpyxl")

try:
    from pptx import Presentation as _Prs
    from pptx.util import Inches as _In, Pt as _Pt
    from pptx.dml.color import RGBColor as _RGB
    from pptx.enum.text import PP_ALIGN as _PA
    import io as _io
    _PPTX_OK = True
except ImportError:
    _PPTX_OK = False

try:
    import cairosvg as _cairosvg
    _CAIRO_OK = True
except ImportError:
    _CAIRO_OK = False

try:
    from playwright.sync_api import sync_playwright as _sync_playwright
    _PLAYWRIGHT_OK = True
except ImportError:
    _PLAYWRIGHT_OK = False

# ─────────────────────────────────────────────────────────────────────────────
# EXTRACÇÃO DE METADADOS
# ─────────────────────────────────────────────────────────────────────────────

def _find_catalogo_sheet(wb):
    """Return the first sheet whose name starts with 'Catalogo' (case-insensitive)."""
    for name in wb.sheetnames:
        if name.lower().startswith("catalogo"):
            return wb[name]
    return None

def _parse_catalogo_sheet(ws):
    """
    Parse a Catalogo_* sheet.
    Expects row 1 as section headers (ignored), row 2 as column names,
    data from row 3 onward.
    Returns list of dicts with keys: code, label, component, url
    """
    # Find header row: look for a row containing 'Código'
    header_row = None
    for i, row in enumerate(ws.iter_rows(values_only=True), 1):
        if any(str(c).strip() == "Código" for c in row if c):
            header_row = i
            headers = [str(c).strip() if c else "" for c in row]
            break
    if header_row is None:
        return []

    col = {h: i for i, h in enumerate(headers)}
    interventions = []
    for row in ws.iter_rows(min_row=header_row + 1, values_only=True):
        code = row[col.get("Código", -1)] if col.get("Código") is not None else None
        if not code or not str(code).strip():
            continue
        code = str(code).strip()
        # Look for either "Actividade" (HIV) or "Intervenção" (other programs)
        label = None
        if "Actividade" in col:
            label = row[col["Actividade"]]
        elif "Intervenção" in col:
            label = row[col["Intervenção"]]
        programa  = row[col["Programa"]]        if "Programa"     in col else None
        component = row[col["Componente"]]      if "Componente"   in col else None
        url       = row[col["URL da Ficha"]]    if "URL da Ficha" in col else None
        team      = row[col["Team"]]            if "Team"         in col else None
        component_txt = str(component).strip() if component else ""
        interventions.append({
            "code":      code,
            "label":     str(label).strip()     if label     else code,
          "programa":  str(programa).strip()  if programa  else "",
          "component": component_txt,
          "comp_macro": comp_macro_from_raw(component_txt),
            "url":       str(url).strip()        if url       else "",
            "team":      str(team).strip()       if team and str(team).strip() not in ("", "None", "nan") else "",
        })
    return interventions

def _parse_submissions_urls(wb):
    """
    Fallback: extract URLs from url_* columns in the Submissions sheet.
    Returns list of dicts with code and url only (label/component empty).
    """
    if "Submissions" not in wb.sheetnames:
        return []
    ws = wb["Submissions"]
    headers = [c.value for c in ws[1]]
    url_cols = [(i, h) for i, h in enumerate(headers)
                if h and str(h).lower().startswith("url_")]
    seen = {}
    for row in ws.iter_rows(min_row=2, values_only=True):
        for i, h in url_cols:
            code = h[4:]  # strip "url_"
            if code not in seen and row[i]:
                seen[code] = str(row[i]).strip()
    return [{"code": c, "label": c, "programa": "", "component": "", "comp_macro": "Outros", "url": u}
            for c, u in sorted(seen.items())]

def _parse_responses_codes(df):
  """Last-resort fallback: just use the codes found in the data."""
  codes = sorted(df["intervention"].dropna().unique())
  return [
    {"code": c, "label": c, "programa": "", "component": "", "comp_macro": "Outros", "url": ""}
    for c in codes
  ]

def load_metadata(results_path, dict_path=None):
    """
    Load intervention metadata with the following priority:
      1. External dictionary file (dict_path), Catalogo_* sheet
      2. Catalogo_* sheet inside results file
      3. url_* columns in Submissions sheet of results file
      4. Raw codes from Responses sheet
    Returns ordered list of dicts: code, label, component, url
    """
    interventions = []

    # 1. External dictionary
    if dict_path and os.path.exists(dict_path):
        wb = openpyxl.load_workbook(dict_path, data_only=True)
        ws = _find_catalogo_sheet(wb)
        if ws:
            interventions = _parse_catalogo_sheet(ws)
            if interventions:
                print(f"  Metadados: dicionário externo '{os.path.basename(dict_path)}' "
                      f"({len(interventions)} intervenções)")
                return interventions

    # 2. Catalogo_* in results file
    wb_res = openpyxl.load_workbook(results_path, data_only=True)
    ws = _find_catalogo_sheet(wb_res)
    if ws:
        interventions = _parse_catalogo_sheet(ws)
        if interventions:
            print(f"  Metadados: folha Catalogo_* no ficheiro de resultados "
                  f"({len(interventions)} intervenções)")
            return interventions

    # 3. url_* columns in Submissions
    interventions = _parse_submissions_urls(wb_res)
    if interventions:
        print(f"  Metadados: colunas url_* na folha Submissions "
              f"({len(interventions)} intervenções, sem nomes)")
        return interventions

    # 4. Raw codes
    return None  # signal to caller to build from df after loading

# ─────────────────────────────────────────────────────────────────────────────
# NORMALIZAÇÃO
# ─────────────────────────────────────────────────────────────────────────────

GATE_NORM = {
    "sim_def": "sim_def", "sim definitivamente": "sim_def",
    "sim,definitivamente": "sim_def", "sim": "sim_def",
    "possivelmente": "possivelmente",
    "nao": "nao", "não": "nao", "no": "nao",
}
YN_NORM = {"sim": "sim", "yes": "sim", "nao": "nao", "não": "nao", "no": "nao"}

def norm_gate(v):
    if pd.isna(v): return None
    return GATE_NORM.get(str(v).strip().lower(), None)

def norm_yn(v):
    if pd.isna(v): return None
    return YN_NORM.get(str(v).strip().lower(), None)

def parse_multi(v):
    if pd.isna(v) or str(v).strip() == "": return []
    return [x.strip() for x in str(v).replace(",", " ").split() if x.strip()]

def parse_other_items(v):
  if pd.isna(v) or str(v).strip() == "":
    return []
  raw = str(v).replace("\r", "\n")
  parts = []
  for chunk in raw.split("\n"):
    parts.extend(chunk.split(";"))
  cleaned = []
  for item in parts:
    txt = str(item).strip(" .,-\t")
    if not txt:
      continue
    low = txt.lower()
    if low in {"none", "na", "n/a", "não", "nao", "sem", "nenhum", "nenhuma"}:
      continue
    cleaned.append(txt)
  return cleaned

def safe_int(v):
    try: return int(float(str(v)))
    except: return None


def comp_macro_from_raw(component):
  """Map heterogeneous componente values into a stable macro-group."""
  txt = str(component or "").replace("\n", " ").strip()
  low = txt.lower()
  if not low:
    return "Outros"
  if "planificada" in low or "em curso" in low:
    return "Intervenção"
  if "orma" in low:
    return "Formações/Treinos"
  if "suppli" in low or "insumo" in low:
    return "Supplies/Insumos"
  if "rh" in low:
    return "Recursos Humanos"
  if "qualidade" in low:
    return "Qualidade"
  if "diagn" in low:
    return "Diagnóstico"
  if "ratamento" in low:
    return "Tratamento"
  if "igil" in low:
    return "Vigilância"
  if "ontrolo vectorial" in low:
    return "Controlo vectorial integrado"
  if "ria na gravidez" in low:
    return "Gestão de malária na gravidez"
  if "uimiopreven" in low:
    return "Quimioprevenção"
  if "a social e de comportamento" in low:
    return "Comunicação para Mudança Social e de Comportamento"
  if "supervis" in low:
    return "Supervisão"
  # No keyword matched — use the component name itself so any program's
  # catalog components appear as their own groups rather than all collapsing
  # into a single "Outros" row.
  return txt if txt else "Outros"

# ─────────────────────────────────────────────────────────────────────────────
# CARREGAMENTO DE DADOS
# ─────────────────────────────────────────────────────────────────────────────

def load_data(path, sheet="Responses"):
    df = pd.read_excel(path, sheet_name=sheet, dtype=str)
    df.columns = [c.strip().lower() for c in df.columns]
    missing = {"expert_code", "intervention", "gate"} - set(df.columns)
    if missing:
        sys.exit(f"Erro: colunas obrigatórias em falta: {missing}")
    return df


def load_expected_experts(experts_file=None):
    """Load unique expected experts from experts.txt, ignoring comments/blank lines."""
    if experts_file is not None:
        selected_file = Path(experts_file)
    else:
        # Default to experts.txt in the current working directory only.
        selected_file = Path.cwd() / "experts.txt"

    if selected_file is None:
        return []
    if not selected_file.exists():
        return []

    experts = set()
    with open(selected_file, "r", encoding="utf-8") as f:
        for line in f:
            raw_line = line.strip()
            lower_line = raw_line.lower()
            if (
                not raw_line
                or raw_line.startswith("#")
                or "# test" in lower_line
                or "# ignore" in lower_line
            ):
                continue
            # Support optional inline comments after '#'
            entry = raw_line.split("#", 1)[0].strip().lower()
            if entry:
                experts.add(entry)

    return sorted(experts)


def _split_codes(value):
    """Split exclusion codes by comma/semicolon/whitespace."""
    if not value:
        return []
    import re
    return [tok.strip() for tok in re.split(r"[,;\s]+", str(value)) if tok.strip()]


def _load_simple_config(config_path=None):
    """Load key=value config from config.env-style files."""
    candidates = []
    if config_path:
        candidates.append(Path(config_path))
    else:
        # Prefer cwd config, then repo root (parent of /code)
        candidates.append(Path.cwd() / "config.env")
        candidates.append(Path(__file__).resolve().parents[1] / "config.env")

    cfg_file = next((p for p in candidates if p.exists()), None)
    if not cfg_file:
        return {}

    cfg = {}
    with open(cfg_file, "r", encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if not line or line.startswith("#") or "=" not in line:
                continue
            k, _, v = line.partition("=")
            cfg[k.strip()] = v.strip()
    return cfg


def load_excluded_codes(cli_codes=None, exclude_file=None, config_path=None):
    """Collect exclusion expert codes from CLI, file, and config file."""
    excluded = set()

    # 1) Config values
    cfg = _load_simple_config(config_path)
    cfg_codes = cfg.get("REPORT_EXCLUDE_EXPERTS", cfg.get("EXCLUDE_EXPERTS", ""))
    excluded.update(_split_codes(cfg_codes))

    cfg_file = cfg.get("REPORT_EXCLUDE_FILE", cfg.get("EXCLUDE_FILE", "")).strip()
    if cfg_file and not exclude_file:
        exclude_file = cfg_file

    # 2) Exclude file (if provided by config or CLI)
    if exclude_file:
        if not os.path.exists(exclude_file):
            print(f"Aviso: ficheiro de exclusões não encontrado: {exclude_file}")
        else:
            with open(exclude_file, "r", encoding="utf-8") as f:
                for line in f:
                    raw = line.strip()
                    if not raw or raw.startswith("#"):
                        continue
                    # Inline comments allowed: CODE_X # reason
                    raw = raw.split("#", 1)[0].strip()
                    excluded.update(_split_codes(raw))

    # 3) CLI list has highest precedence (adds more exclusions)
    excluded.update(_split_codes(cli_codes or ""))

    # Normalize and de-duplicate while preserving intent
    return sorted({c.strip() for c in excluded if c and c.strip()})

# ─────────────────────────────────────────────────────────────────────────────
# AGREGAÇÃO
# ─────────────────────────────────────────────────────────────────────────────

def aggregate(df, interventions):
    """Aggregate responses per intervention. Returns dict keyed by code."""
    def is_placeholder_comment(value):
        txt = str(value).strip().lower()
        txt = txt.strip('"\'`.,;:!?()[]{}')
        return txt in {
            "sem comentarios",
            "sem comentário",
            "sem comentario",
            "sem comentários",
            "sem comment",
            "no comments",
        }

    inv_codes  = [r["code"] for r in interventions]
    inv_labels = {r["code"]: r["label"] for r in interventions}
    results = {}

    for inv in interventions:
        code = inv["code"]
        sub  = df[df["intervention"] == code].copy()

        # Per expert: prefer row with gate filled
        deduped = []
        for _, grp in sub.groupby("expert_code", sort=False):
            filled = grp[grp["gate"].notna() & (grp["gate"].str.strip() != "")]
            deduped.append(filled.iloc[0] if len(filled) > 0 else grp.iloc[0])

        n_sim = n_poss = n_nao = n_missing = 0
        n_imp_1 = n_imp_2 = n_imp_3 = 0
        impacts = []
        dup_yes = intg_yes = res_yes = 0
        which_dup_counts  = defaultdict(int)
        which_intg_counts = defaultdict(int)
        which_dup_other_counts = defaultdict(int)
        which_intg_other_counts = defaultdict(int)
        comments = []
        # For scoring diagnostics
        gate_scores = []   # (gate_score, impact, exp_weight) per respondent
        gate_only_scores = []  # gate_score only (unweighted), includes nao=0 when gate answered
        all_exp_vals = []      # experience weight for ALL gate-answered respondents (for S_pond denominator)

        GATE_SCORE = {"sim_def": 1.0, "possivelmente": 0.5, "nao": 0.0}

        for row in deduped:
            g = norm_gate(row.get("gate", None))
            if   g == "sim_def":       n_sim  += 1
            elif g == "possivelmente": n_poss += 1
            elif g == "nao":           n_nao  += 1
            else:                      n_missing += 1

            if g is not None:
                gs  = GATE_SCORE.get(g, 0.0)
                exp = safe_int(row.get("exp", None))
                exp = exp if exp and 1 <= exp <= 3 else 1  # neutral weight if missing
                gate_only_scores.append(gs)
                all_exp_vals.append(exp)

            if g in ("sim_def", "possivelmente"):
                imp = safe_int(row.get("impact", None))
                if imp and 1 <= imp <= 3:
                    impacts.append(imp)
                    gate_scores.append((gs, imp, exp))
                    if imp == 1:   n_imp_1 += 1
                    elif imp == 2: n_imp_2 += 1
                    elif imp == 3: n_imp_3 += 1

                if norm_yn(row.get("dup", None)) == "sim":
                    dup_yes += 1
                    for t in parse_multi(row.get("which_dup", None)):
                        if t in inv_codes and t != code:
                            which_dup_counts[t] += 1
                    for t in parse_other_items(row.get("which_dup_other", None)):
                        which_dup_other_counts[t] += 1

                if norm_yn(row.get("intg", None)) == "sim":
                    intg_yes += 1
                    for t in parse_multi(row.get("which_intg", None)):
                        if t in inv_codes and t != code:
                            which_intg_counts[t] += 1
                    for t in parse_other_items(row.get("which_intg_other", None)):
                        which_intg_other_counts[t] += 1

                if norm_yn(row.get("res", None)) == "sim":
                    res_yes += 1

                cmt = row.get("cmt", None)
                if cmt and not pd.isna(cmt):
                  cmt_txt = str(cmt).strip()
                  if cmt_txt not in ("", ".", "None") and not is_placeholder_comment(cmt_txt):
                    comments.append(cmt_txt)

        n_resp = n_sim + n_poss + n_nao
        n_positive = n_sim + n_poss

        top_dup  = sorted(which_dup_counts,  key=which_dup_counts.get,  reverse=True)[:3]
        top_intg = sorted(which_intg_counts, key=which_intg_counts.get, reverse=True)[:3]

        avg_impact = round(sum(impacts) / len(impacts), 2) if impacts else 0

        # S_base = mean(gate_score over ALL respondents) × mean(impact over POSITIVE respondents)
        # gate_only_scores includes gs=0 for "não", so mean correctly reflects overall support.
        # impact mean is only from those who endorsed optimisation (they're the ones who rated impact).
        if gate_scores:
            imp_vals = [x[1] for x in gate_scores]
            exp_vals = [x[2] for x in gate_scores]
            mean_gs  = sum(gate_only_scores) / len(gate_only_scores)  # over ALL gate-answered
            mean_imp = sum(imp_vals) / len(imp_vals)                  # over positive respondents
            s_base   = round(mean_gs * mean_imp, 3)
            # S_pond = Σ(gs × imp × exp) / Σ(exp over ALL respondents)
            # nao respondents contribute gs=0 → numerator unchanged, but their exp weight
            # belongs in the denominator so high-experience dissenters reduce the score.
            numerator   = sum(gs * imp * exp for gs, imp, exp in gate_scores)
            denominator = sum(all_exp_vals)  # all gate-answered respondents
            s_pond   = round(numerator / denominator, 3) if denominator else 0
            exp_mean = round(sum(exp_vals) / len(exp_vals), 2)
        else:
            s_base = s_pond = exp_mean = 0.0

        gate_mean = round(sum(gate_only_scores) / len(gate_only_scores), 3) if gate_only_scores else 0.0

        pct_dup = round(dup_yes / n_positive * 100) if n_positive else 0
        pct_intg = round(intg_yes / n_positive * 100) if n_positive else 0
        pct_res = round(res_yes / n_positive * 100) if n_positive else 0
        e_score = round((pct_res + pct_dup + pct_intg) / 300, 3) if n_positive else 0.0

        results[code] = {
            "code":      code,
            "label":     inv["label"],
            "url":       inv["url"],
          "programa":  inv.get("programa", ""),
            "component": inv["component"],
          "comp_macro": inv.get("comp_macro", "Outros"),
            "team":      inv.get("team", ""),
            "n_total":   n_resp,
          "n_positive": n_positive,
            "n_missing": n_missing,
            "n_sim":     n_sim,
            "n_poss":    n_poss,
            "n_nao":     n_nao,
            "pct_optimizable": round((n_sim + n_poss) / n_resp * 100) if n_resp else 0,
            "pct_definitely":  round(n_sim / n_resp * 100)            if n_resp else 0,
            "avg_impact": avg_impact,
            "n_imp_1":   n_imp_1,
            "n_imp_2":   n_imp_2,
            "n_imp_3":   n_imp_3,
          "dup_pct":   pct_dup,
          "intg_pct":  pct_intg,
          "res_pct":   pct_res,
          "e_score":   e_score,
            "top_dup":   top_dup,
            "top_intg":  top_intg,
            "dup_counts":  dict(which_dup_counts),
            "intg_counts": dict(which_intg_counts),
            "dup_other_counts": dict(which_dup_other_counts),
            "intg_other_counts": dict(which_intg_other_counts),
            "comments":  comments,
            "composite": round((n_sim / n_resp) * avg_impact, 3) if n_resp and avg_impact else 0,
            "gate_mean": gate_mean,
            "s_base":    s_base,
            "s_pond":    s_pond,
            "exp_mean":  exp_mean,
        }

    return results

# ─────────────────────────────────────────────────────────────────────────────
# ESTATÍSTICAS SUMÁRIAS
# ─────────────────────────────────────────────────────────────────────────────

def summary_stats(results, df, n_experts_expected=None, n_experts_per_team=None):
    n_experts_observed = df["expert_code"].nunique()
    n_experts = n_experts_expected if n_experts_expected and n_experts_expected > 0 else n_experts_observed
    n_inv_80    = sum(1 for r in results.values() if r["pct_optimizable"] >= 80)
    n_imp_high  = sum(1 for r in results.values() if r["avg_impact"] >= 2.5)
    n_unanimous = sum(1 for r in results.values() if r["pct_optimizable"] == 100)

    # Response rate per intervention: n_total / n_experts (% of experts who answered)
    n_experts_int = n_experts if n_experts > 0 else 1
    rr = [r["n_total"] / n_experts_int * 100 for r in results.values() if r["n_total"] > 0]
    rr_median = round(statistics.median(rr), 1) if rr else 0
    rr_min    = round(min(rr), 1)               if rr else 0
    rr_max    = round(max(rr), 1)               if rr else 0

    # Per-intervention breakdown for the response rate table
    rr_detail = sorted(
        [{"label": r["label"], "n": r["n_total"], "pct": round(r["n_total"] / n_experts_int * 100, 1)}
         for r in results.values()],
        key=lambda x: x["pct"], reverse=True
    )

    # Per-team response rate summary (only when Team column is populated in catalog)
    # Denominator: N_EXPERTS_TEAM_X from config when available; otherwise observed count.
    team_codes: dict = {}   # team -> set of intervention codes in that team
    for r in results.values():
        t = r.get("team", "")
        if not t:
            continue
        team_codes.setdefault(t, set()).add(r["code"])

    rr_by_team = []
    _n_per_team = n_experts_per_team or {}
    for t in sorted(team_codes):
        codes = team_codes[t]
        # Prefer configured expert count; fall back to observed unique experts in team
        if t in _n_per_team:
            team_n = _n_per_team[t]
        else:
            team_n = df[df["intervention"].isin(codes)]["expert_code"].nunique() or n_experts_int
        rrs = [results[c]["n_total"] / team_n * 100 for c in codes if c in results]
        if not rrs:
            continue
        rr_by_team.append({
            "team":       t,
            "n_inv":      len(rrs),
            "n_experts":  team_n,
            "rr_median":  round(statistics.median(rrs), 1),
            "rr_min":     round(min(rrs), 1),
            "rr_max":     round(max(rrs), 1),
        })

    return {
        "n_experts":   n_experts,
      "n_experts_observed": n_experts_observed,
        "n_inv":       len(results),
        "n_inv_80":    n_inv_80,
        "n_imp_high":  n_imp_high,
        "n_unanimous": n_unanimous,
        "rr_median":   rr_median,
        "rr_min":      rr_min,
        "rr_max":      rr_max,
        "rr_detail":   rr_detail,
        "rr_by_team":  rr_by_team,
    }

# ─────────────────────────────────────────────────────────────────────────────
# PONTUAÇÃO E RANKS
# ─────────────────────────────────────────────────────────────────────────────

def compute_ranks(results):
    """Add ranks and normalized scores used in scoring diagnostics."""
    items = list(results.values())
    # Normalise to 0-100
    max_gate = max((r.get("gate_mean", 0) for r in items), default=1) or 1
    max_base = max((r["s_base"] for r in items), default=1) or 1
    max_pond = max((r["s_pond"] for r in items), default=1) or 1
    for r in items:
        r["score_gate_n"] = round(r.get("gate_mean", 0) / max_gate * 100, 1)
        r["score_base_n"] = round(r["s_base"] / max_base * 100, 1)
        r["score_wtd_n"]  = round(r["s_pond"] / max_pond * 100, 1)

    sorted_gate = sorted(items, key=lambda r: r.get("gate_mean", 0), reverse=True)
    sorted_base = sorted(items, key=lambda r: r["s_base"], reverse=True)
    sorted_pond = sorted(items, key=lambda r: r["s_pond"], reverse=True)
    rank_gate = {r["code"]: i + 1 for i, r in enumerate(sorted_gate)}
    rank_base = {r["code"]: i + 1 for i, r in enumerate(sorted_base)}
    rank_pond = {r["code"]: i + 1 for i, r in enumerate(sorted_pond)}
    for r in items:
        r["rank_gate"]  = rank_gate[r["code"]]
        r["rank_base"]  = rank_base[r["code"]]
        r["rank_wtd"]   = rank_pond[r["code"]]
        r["rank_delta"] = rank_base[r["code"]] - rank_pond[r["code"]]

# ─────────────────────────────────────────────────────────────────────────────
# SVG HELPERS (sem dependências externas)
# ─────────────────────────────────────────────────────────────────────────────

def svg_hbar_stacked(rows, width=460, bar_h=14, gap=4,
           colors=("#1a6b3a","#d4a017","#e0e0e0"),
           labels=("Sim def.","Possiv.","Não"),
           label_w=220, hrefs=None, sublabels=None, markers=None):
  """Horizontal stacked bar chart. rows = list of (label, [v1,v2,v3]).
  hrefs: optional list of anchor hrefs (one per row) to make labels clickable.
  sublabels: optional list of secondary label strings (e.g. component names).
  markers: optional list of floats in [0,1] — fractional position along bar width
           where a score diamond is overlaid (e.g. gate_mean or (avg_impact-1)/2).
  """
  sub_h = 9  # extra px per row when sublabels present
  has_sub = bool(sublabels)
  row_h = bar_h + (sub_h if has_sub else 0)
  h = len(rows) * (row_h + gap) + 4
  bar_total_w = width - label_w - 40
  out = [f'<svg width="{width}" height="{h}" style="display:block;overflow:visible">']
  for i, (lbl, vals) in enumerate(rows):
    y = i * (row_h + gap)
    total = sum(vals) or 1
    href = hrefs[i] if hrefs and i < len(hrefs) and hrefs[i] else None
    lbl_fill = '#1a5276' if href else '#6b7280'
    lbl_style = ' style="cursor:pointer"' if href else ''
    lbl_text = (f'<text x="{label_w-4}" y="{y+bar_h-2}" text-anchor="end"'
          f' font-size="10" fill="{lbl_fill}"{lbl_style}>{esc(lbl)}</text>')
    out.append(f'<a href="{esc(href)}">{lbl_text}</a>' if href else lbl_text)
    if has_sub:
      sub = (sublabels[i] if i < len(sublabels) else "") or ""
      if sub:
        out.append(f'<text x="{label_w-4}" y="{y+bar_h+sub_h-2}" text-anchor="end"'
                   f' font-size="8" fill="#9ca3af" font-style="italic">{esc(sub)}</text>')
    x = label_w
    for v, col in zip(vals, colors):
      w = round(v / total * bar_total_w)
      if w > 0:
        out.append(f'<rect x="{x}" y="{y}" width="{w}" height="{bar_h}" fill="{col}"/>')
      x += w
    pct = round(sum(vals[:2]) / total * 100) if total else 0
    out.append(f'<text x="{label_w + bar_total_w + 4}" y="{y+bar_h-2}" '
           f'font-size="10" fill="#374151" font-weight="600">{pct}%</text>')
    # Score marker diamond overlaid on bar
    if markers is not None and i < len(markers) and markers[i] is not None:
      frac = max(0.0, min(1.0, float(markers[i])))
      mx = label_w + round(frac * bar_total_w)
      cy_d = y + bar_h // 2
      r_d = 4
      out.append(
        f'<polygon points="{mx},{cy_d-r_d} {mx+r_d},{cy_d} '
        f'{mx},{cy_d+r_d} {mx-r_d},{cy_d}" '
        f'fill="white" stroke="#1e3a5f" stroke-width="1.5" opacity="0.92"/>'
      )
  out.append("</svg>")
  return "".join(out)

def svg_hbar_single(rows, width=460, bar_h=14, gap=4, color="#1a6b3a", label_w=220, fmt=".1f", hrefs=None, sublabels=None):
  """Single horizontal bar per row. rows = list of (label, value, max_value).
  hrefs: optional list of anchor hrefs (one per row) to make labels clickable.
  sublabels: optional list of secondary label strings (e.g. component names).
  """
  sub_h = 9
  has_sub = bool(sublabels)
  row_h = bar_h + (sub_h if has_sub else 0)
  h = len(rows) * (row_h + gap) + 4
  out = [f'<svg width="{width}" height="{h}" style="display:block;overflow:visible">']
  for i, (lbl, val, max_val) in enumerate(rows):
    y = i * (row_h + gap)
    bar_w = width - label_w - 50
    w = round(val / max_val * bar_w) if max_val else 0
    href = hrefs[i] if hrefs and i < len(hrefs) and hrefs[i] else None
    lbl_fill = '#1a5276' if href else '#6b7280'
    lbl_style = ' style="cursor:pointer"' if href else ''
    lbl_text = (f'<text x="{label_w-4}" y="{y+bar_h-2}" text-anchor="end"'
          f' font-size="10" fill="{lbl_fill}"{lbl_style}>{esc(lbl)}</text>')
    out.append(f'<a href="{esc(href)}">{lbl_text}</a>' if href else lbl_text)
    if has_sub:
      sub = (sublabels[i] if i < len(sublabels) else "") or ""
      if sub:
        out.append(f'<text x="{label_w-4}" y="{y+bar_h+sub_h-2}" text-anchor="end"'
                   f' font-size="8" fill="#9ca3af" font-style="italic">{esc(sub)}</text>')
    out.append(f'<rect x="{label_w}" y="{y}" width="{bar_w}" height="{bar_h}" fill="#e5e7eb" rx="1"/>')
    if w > 0:
      out.append(f'<rect x="{label_w}" y="{y}" width="{w}" height="{bar_h}" fill="{color}" rx="1"/>')
    out.append(f'<text x="{label_w + bar_w + 4}" y="{y+bar_h-2}" '
           f'font-size="10" fill="#374151" font-weight="600">'
           f'{val:{fmt}}</text>')
  out.append("</svg>")
  return "".join(out)

def svg_donut(counts_dict, colors, size=100, label=""):
    """Simple donut chart. counts_dict = {label: count}."""
    total = sum(counts_dict.values()) or 1
    cx = cy = size / 2
    r_out, r_in = size * 0.4, size * 0.22
    import math
    angle = -math.pi / 2
    paths = []
    legend = []
    for (lbl, cnt), col in zip(counts_dict.items(), colors):
        sweep = cnt / total * 2 * math.pi
        if sweep < 0.001:
            angle += sweep
            continue
        x1 = cx + r_out * math.cos(angle)
        y1 = cy + r_out * math.sin(angle)
        x2 = cx + r_out * math.cos(angle + sweep)
        y2 = cy + r_out * math.sin(angle + sweep)
        xi1 = cx + r_in * math.cos(angle)
        yi1 = cy + r_in * math.sin(angle)
        xi2 = cx + r_in * math.cos(angle + sweep)
        yi2 = cy + r_in * math.sin(angle + sweep)
        large = 1 if sweep > math.pi else 0
        pct = round(cnt / total * 100)
        paths.append(
            f'<path d="M{xi1:.1f},{yi1:.1f} L{x1:.1f},{y1:.1f} '
            f'A{r_out},{r_out} 0 {large},1 {x2:.1f},{y2:.1f} '
            f'L{xi2:.1f},{yi2:.1f} A{r_in},{r_in} 0 {large},0 {xi1:.1f},{yi1:.1f} Z" '
            f'fill="{col}"><title>{esc(lbl)}: {cnt} ({pct}%)</title></path>'
        )
        legend.append((lbl, cnt, pct, col))
        angle += sweep
    legend_h = len(legend) * 14 + 4
    total_h = max(size, legend_h)
    svg_parts = [f'<svg width="{size + 140}" height="{total_h}" style="display:inline-block;vertical-align:middle">']
    svg_parts += paths
    if label:
        svg_parts.append(f'<text x="{cx}" y="{cy+4}" text-anchor="middle" '
                         f'font-size="9" fill="#6b7280">{esc(label)}</text>')
    for j, (lbl, cnt, pct, col) in enumerate(legend):
        ly = j * 14 + 10
        svg_parts.append(f'<rect x="{size+4}" y="{ly-7}" width="9" height="9" fill="{col}" rx="1"/>')
        svg_parts.append(f'<text x="{size+16}" y="{ly+1}" font-size="10" fill="#374151">'
                         f'{esc(lbl)} <tspan font-weight="600">{pct}%</tspan></text>')
    svg_parts.append("</svg>")
    return "".join(svg_parts)

def svg_scatter_optim_impact_exp(items, width=1060, height=420):
  """
  Single real-scale scatter with overlays:
  - Axes: real x/y values, viewport adjusted to observed ranges
  - Background shaded zones for current dual-cutoff rationale
  Encoding: size = experience mean, color = % reduction resources.
  """
  if not items:
    return '<div style="font-size:12px;color:#6b7280">Sem dados para gráfico de dispersão.</div>'

  # Color/size helpers
  def res_color(res_pct):
    p = max(0.0, min(100.0, float(res_pct or 0.0)))
    t = p / 100.0
    c1 = (226, 232, 240)
    c2 = (91, 94, 166)
    r = round(c1[0] + (c2[0] - c1[0]) * t)
    g = round(c1[1] + (c2[1] - c1[1]) * t)
    b = round(c1[2] + (c2[2] - c1[2]) * t)
    return f"rgb({r},{g},{b})"

  def exp_radius(exp_mean):
    e = max(1.0, min(3.0, float(exp_mean or 1.0)))
    return 4 + (e - 1.0) / 2.0 * 7

  gate_vals = [float(r.get("gate_mean", 0) or 0) for r in items]
  imp_vals = [float(r.get("avg_impact", 1.0) or 1.0) for r in items]
  gmin, gmax = min(gate_vals), max(gate_vals)
  imin, imax = min(imp_vals), max(imp_vals)
  gspan_obs = (gmax - gmin) if (gmax - gmin) > 1e-12 else 1.0
  ispan_obs = (imax - imin) if (imax - imin) > 1e-12 else 1.0

  # Keep plot in real units but zoom to observed range with a small padding.
  x_min = max(0.0, gmin - 0.08 * gspan_obs)
  x_max = min(1.0, gmax + 0.08 * gspan_obs)
  y_min = max(1.0, imin - 0.08 * ispan_obs)
  y_max = min(3.0, imax + 0.08 * ispan_obs)
  x_span = (x_max - x_min) if (x_max - x_min) > 1e-12 else 1.0
  y_span = (y_max - y_min) if (y_max - y_min) > 1e-12 else 1.0

  panel_w = 700
  panel_h = height
  left = 56
  right = 20
  top = 34
  bottom = 50
  w = panel_w - left - right
  h = panel_h - top - bottom

  def x_px(v):
    v = max(x_min, min(x_max, float(v)))
    return left + ((v - x_min) / x_span) * w

  def y_px(v):
    v = max(y_min, min(y_max, float(v)))
    return top + (1.0 - ((v - y_min) / y_span)) * h

  x_cut_actual = statistics.median(gate_vals)
  y_cut_actual = statistics.median(imp_vals)
  x_cut = max(x_min, min(x_max, x_cut_actual))
  y_cut = max(y_min, min(y_max, y_cut_actual))

  out = [f'<svg width="{panel_w}" height="{panel_h}" style="display:block;overflow:visible">']
  out.append(f'<text x="{left}" y="18" font-size="11" fill="#334155" font-weight="600">Escala real (janela min-max observada)</text>')

  # Shaded areas for current dual-cutoff rationale (3 recommendation structures)
  xl = x_px(x_min)
  xm = x_px(x_cut)
  xr = x_px(x_max)
  yb = y_px(y_min)
  ym = y_px(y_cut)
  yt = y_px(y_max)
  out.append(f'<rect x="{xl:.1f}" y="{ym:.1f}" width="{(xm-xl):.1f}" height="{(yb-ym):.1f}" fill="#fee2e2" fill-opacity="0.35"/>')
  out.append(f'<rect x="{xm:.1f}" y="{ym:.1f}" width="{(xr-xm):.1f}" height="{(yb-ym):.1f}" fill="#fef3c7" fill-opacity="0.30"/>')
  out.append(f'<rect x="{xl:.1f}" y="{yt:.1f}" width="{(xm-xl):.1f}" height="{(ym-yt):.1f}" fill="#fef3c7" fill-opacity="0.30"/>')
  out.append(f'<rect x="{xm:.1f}" y="{yt:.1f}" width="{(xr-xm):.1f}" height="{(ym-yt):.1f}" fill="#dcfce7" fill-opacity="0.30"/>')

  # Axes
  out.append(f'<line x1="{left}" y1="{top + h}" x2="{left + w}" y2="{top + h}" stroke="#94a3b8" stroke-width="1"/>')
  out.append(f'<line x1="{left}" y1="{top}" x2="{left}" y2="{top + h}" stroke="#94a3b8" stroke-width="1"/>')

  # Grid and ticks in real-value axes
  x_ticks = [x_min + i * (x_span / 4.0) for i in range(5)]
  y_ticks = [y_min + i * (y_span / 4.0) for i in range(5)]
  for xv in x_ticks:
    xp = x_px(xv)
    out.append(f'<line x1="{xp:.1f}" y1="{top}" x2="{xp:.1f}" y2="{top + h}" stroke="#e5e7eb" stroke-width="1"/>')
    out.append(f'<text x="{xp:.1f}" y="{top + h + 16}" text-anchor="middle" font-size="10" fill="#64748b">{xv:.2f}</text>')
  for yv in y_ticks:
    yp = y_px(yv)
    out.append(f'<line x1="{left}" y1="{yp:.1f}" x2="{left + w}" y2="{yp:.1f}" stroke="#e5e7eb" stroke-width="1"/>')
    out.append(f'<text x="{left - 8}" y="{yp + 3:.1f}" text-anchor="end" font-size="10" fill="#64748b">{yv:.2f}</text>')

  # Current dual cutoffs in real-value axes (clipped to viewport)
  out.append(f'<line x1="{x_px(x_cut):.1f}" y1="{top}" x2="{x_px(x_cut):.1f}" y2="{top + h}" stroke="#64748b" stroke-width="1" stroke-dasharray="4 3"/>')
  out.append(f'<line x1="{left}" y1="{y_px(y_cut):.1f}" x2="{left + w}" y2="{y_px(y_cut):.1f}" stroke="#64748b" stroke-width="1" stroke-dasharray="4 3"/>')
  out.append(f'<text x="{left + 6}" y="{top + 14}" font-size="9" fill="#64748b">cutoff x&gt;mediana={x_cut_actual:.3f}</text>')
  out.append(f'<text x="{left + 6}" y="{top + 26}" font-size="9" fill="#64748b">cutoff y&gt;mediana={y_cut_actual:.3f}</text>')

  # Points — label shows numeric suffix of code (e.g. hiv_01 → 01)
  _code_suffix_re = __import__('re').compile(r'(\d+)$')
  for i, r in enumerate(sorted(items, key=lambda z: z.get("composite", 0), reverse=True), 1):
    code = r.get("code", "")
    _m = _code_suffix_re.search(code)
    dot_label = _m.group(1) if _m else str(i)
    x = x_px(float(r.get("gate_mean", 0) or 0))
    y = y_px(float(r.get("avg_impact", 1.0) or 1.0))
    fill = res_color(r.get("res_pct", 0))
    rad = exp_radius(r.get("exp_mean", 1.0))
    tip = (f'{code} · {esc(r.get("label", ""))} | '
           f'triagem: {r.get("gate_mean", 0):.3f} | impacto: {r.get("avg_impact", 0):.2f} | '
           f'S_base: {r.get("s_base", 0):.3f} | S_pond: {r.get("s_pond", 0):.3f}')
    out.append(f'<circle cx="{x:.1f}" cy="{y:.1f}" r="{rad:.1f}" fill="{fill}" fill-opacity="0.92" stroke="#0f172a" stroke-width="0.7">'
               f'<title>{tip}</title></circle>')
    out.append(f'<text x="{x:.1f}" y="{y + 3:.1f}" text-anchor="middle" font-size="8" fill="#0f172a" font-weight="700" pointer-events="none">{dot_label}</text>')

  # Axis labels
  out.append(f'<text x="{left + w / 2:.1f}" y="{panel_h - 8}" text-anchor="middle" font-size="11" fill="#334155">Score de optimizabilidade</text>')
  out.append(f'<text x="14" y="{top + h / 2:.1f}" transform="rotate(-90 14 {top + h / 2:.1f})" text-anchor="middle" font-size="11" fill="#334155">Impacto esperado (1-3)</text>')

  out.append('</svg>')
  scatter_svg = ''.join(out)

  # Shared legend block
  legend = (
    '<div style="margin-top:6px;display:flex;gap:16px;align-items:flex-start;flex-wrap:wrap">'
    '<div style="font-size:10px;color:#64748b">'
    '<div style="margin-bottom:4px;font-weight:600">Cor = % redução recursos</div>'
    '<div style="width:150px;height:10px;background:linear-gradient(90deg,rgb(226,232,240),rgb(91,94,166));border:1px solid #94a3b8"></div>'
    '<div style="display:flex;justify-content:space-between;width:150px"><span>0%</span><span>50%</span><span>100%</span></div>'
    '</div>'
    '<div style="font-size:10px;color:#64748b">'
    '<div style="margin-bottom:4px;font-weight:600">Tamanho = experiência média</div>'
    '<div style="display:flex;align-items:center;gap:8px"><span style="display:inline-block;width:8px;height:8px;border-radius:50%;background:#cbd5e1;border:1px solid #64748b"></span><span>1.0</span>'
    '<span style="display:inline-block;width:14px;height:14px;border-radius:50%;background:#cbd5e1;border:1px solid #64748b"></span><span>2.0</span>'
    '<span style="display:inline-block;width:20px;height:20px;border-radius:50%;background:#cbd5e1;border:1px solid #64748b"></span><span>3.0</span></div>'
    '</div>'
    '<div style="font-size:10px;color:#64748b">'
    '<div style="margin-bottom:4px;font-weight:600">Overlays de cutoff</div>'
    '<div style="display:flex;gap:10px;align-items:center;flex-wrap:wrap">'
    '<span style="display:inline-block;width:12px;height:12px;background:#dcfce7;border:1px solid #86efac"></span><span>Alta (caixa atual)</span>'
    '<span style="display:inline-block;width:12px;height:12px;background:#fef3c7;border:1px solid #fcd34d"></span><span>Média (caixas atuais)</span>'
    '<span style="display:inline-block;width:12px;height:12px;background:#fee2e2;border:1px solid #fca5a5"></span><span>Baixa (caixa atual)</span>'
    '</div></div>'
    '</div>'
  )

  return (
    f'{scatter_svg}'
    + legend
  )

def svg_alluvial_weighting(items, width=1200, row_h=24, node_w=14, label_w=560):
    """
    Draw an alluvial-like SVG linking rank positions across 2 stages:
    S_base (left) -> S_pond (right).
    Left side: rank numbers only. Right side: full intervention name + delta.
    Labels are on the right so the diagram reads left-to-right with names visible.
    """
    if not items:
      return '<div style="font-size:12px;color:#6b7280">Sem dados para diagrama aluvial.</div>'

    left_num_w = 36   # narrow left margin — just enough for rank numbers
    left_x = left_num_w
    right_x = width - label_w - node_w
    top_pad = 42
    bottom_pad = 24
    n = len(items)
    height = top_pad + bottom_pad + max(1, n) * row_h

    left_order  = sorted(items, key=lambda r: (r.get("rank_base", 10**9), r.get("label", "")))
    right_order = sorted(items, key=lambda r: (r.get("rank_wtd",  10**9), r.get("label", "")))

    left_y  = {r["code"]: top_pad + i * row_h + row_h / 2 for i, r in enumerate(left_order)}
    right_y = {r["code"]: top_pad + i * row_h + row_h / 2 for i, r in enumerate(right_order)}

    def flow_color(delta):
        if delta > 0: return "#1a6b3a"
        if delta < 0: return "#c0392b"
        return "#1a5276"

    out = [
      f'<svg width="{width}" height="{height}" style="display:block;overflow:visible">',
      f'<text x="{left_x + node_w // 2}" y="20" text-anchor="middle" font-size="11" fill="#6b7280" '
      f'style="letter-spacing:.06em;text-transform:uppercase">S_base</text>',
      f'<text x="{right_x + node_w // 2}" y="20" text-anchor="middle" font-size="11" fill="#6b7280" '
      f'style="letter-spacing:.06em;text-transform:uppercase">S_pond</text>',
      f'<line x1="{left_x + node_w}" y1="30" x2="{right_x}" y2="30" stroke="#e5e7eb" stroke-width="1"/>'
    ]

    # Bezier curves
    for r in left_order:
      code = r["code"]
      y1 = left_y[code]
      y2 = right_y[code]
      c1x = left_x + node_w + (right_x - (left_x + node_w)) * 0.35
      c2x = left_x + node_w + (right_x - (left_x + node_w)) * 0.65
      col = flow_color(r.get("rank_delta", 0))
      _comp_ttip = (f' | {esc(r["component"])}' if r.get("component") else "")
      out.append(
        f'<path d="M {left_x + node_w:.1f},{y1:.1f} C {c1x:.1f},{y1:.1f} {c2x:.1f},{y2:.1f} {right_x:.1f},{y2:.1f}" '
        f'stroke="{col}" stroke-opacity="0.45" stroke-width="6" fill="none">'
        f'<title>#{r.get("display_idx", "-")} — {esc(r["label"])}{_comp_ttip} | Rank S_base: {r.get("rank_base", "-")} → Rank S_pond: {r.get("rank_wtd", "-")} '
        f'| Δrank: {r.get("rank_delta", 0)}</title></path>'
      )

    # Left nodes — rank number only
    for r in left_order:
      code = r["code"]
      y = left_y[code] - 6
      out.append(f'<rect x="{left_x}" y="{y:.1f}" width="{node_w}" height="12" fill="#1a6b3a" rx="2"/>')
      out.append(
        f'<text x="{left_x - 4}" y="{y + 9:.1f}" text-anchor="end" font-size="10" fill="#6b7280">'
        f'{r.get("rank_base", "-")}</text>'
      )

    # Right nodes — full name + component + delta
    for r in right_order:
      code = r["code"]
      y = right_y[code] - 6
      out.append(f'<rect x="{right_x}" y="{y:.1f}" width="{node_w}" height="12" fill="#1a5276" rx="2"/>')
      delta = r.get("rank_delta", 0)
      delta_str = f'+{delta}' if delta > 0 else str(delta)
      delta_col = "#1a6b3a" if delta > 0 else ("#c0392b" if delta < 0 else "#6b7280")
      _comp_lbl = (f' — {r["component"]}' if r.get("component") else "")
      out.append(
        f'<text x="{right_x + node_w + 6}" y="{y + 9:.1f}" text-anchor="start" font-size="10" fill="#111827">'
        f'{r.get("display_idx", "-")}. {esc(r["label"])}{esc(_comp_lbl)} '
        f'<tspan fill="{delta_col}" font-weight="600">(Δ{delta_str})</tspan></text>'
      )

    out.append('</svg>')
    return ''.join(out)

# ─────────────────────────────────────────────────────────────────────────────
# ANÁLISE UNIVARIADA
# ─────────────────────────────────────────────────────────────────────────────

def univariate_analysis(results, df):
    """Compute cross-intervention univariate distributions for the report."""
    items = list(results.values())

    # Gate aggregate
    total_sim  = sum(r["n_sim"]  for r in items)
    total_poss = sum(r["n_poss"] for r in items)
    total_nao  = sum(r["n_nao"]  for r in items)

    # Impact distribution
    imp_counts = {1: 0, 2: 0, 3: 0}
    for r in items:
        for v, n in [(1, r["n_nao"]), (2, r["n_poss"]), (3, r["n_sim"])]:
            pass  # we need per-respondent data; use avg_impact × n_total as proxy
    # Use per-row impact from df for accurate distribution
    imp_vals_all = []
    for v in df["impact"].dropna():
        i = safe_int(v)
        if i and 1 <= i <= 3:
            imp_vals_all.append(i)
            imp_counts[i] = imp_counts.get(i, 0) + 1

    # Expertise distribution
    exp_counts = {1: 0, 2: 0, 3: 0}
    for v in df["exp"].dropna() if "exp" in df.columns else []:
        e = safe_int(v)
        if e and 1 <= e <= 3:
            exp_counts[e] = exp_counts.get(e, 0) + 1

    # Per-intervention sorted lists for bar charts
    sorted_gate   = sorted(items, key=lambda r: r["pct_optimizable"], reverse=True)
    sorted_impact = sorted(items, key=lambda r: r["avg_impact"], reverse=True)
    sorted_dup    = sorted(items, key=lambda r: r["dup_pct"],  reverse=True)
    sorted_intg   = sorted(items, key=lambda r: r["intg_pct"], reverse=True)
    sorted_res    = sorted(items, key=lambda r: r["res_pct"],  reverse=True)
    sorted_e      = sorted(items, key=lambda r: r.get("e_score", 0), reverse=True)
    e_vals = [r.get("e_score", 0) for r in items]

    return {
        "gate_agg":      {"Sim def.": total_sim, "Possiv.": total_poss, "Não": total_nao},
        "imp_counts":    imp_counts,
        "exp_counts":    exp_counts,
        "sorted_gate":   sorted_gate,
        "sorted_impact": sorted_impact,
        "sorted_dup":    sorted_dup,
        "sorted_intg":   sorted_intg,
        "sorted_res":    sorted_res,
        "sorted_e":      sorted_e,
        "e_median":      round(statistics.median(e_vals), 3) if e_vals else 0,
        "e_min":         round(min(e_vals), 3) if e_vals else 0,
        "e_max":         round(max(e_vals), 3) if e_vals else 0,
        "imp_vals_all":  imp_vals_all,
    }

# ─────────────────────────────────────────────────────────────────────────────
# GERAÇÃO DE HTML
# ─────────────────────────────────────────────────────────────────────────────

def esc(s):
    return str(s).replace("&","&amp;").replace("<","&lt;").replace(">","&gt;").replace('"','&quot;')

def extract_numeric_code(full_code):
    """Extract leading numeric component from code (e.g., '01' from 'hiv_01'; '1.1' from '1.1.001')."""
    import re
    code_str = str(full_code).strip()
    
    # First, try to extract numeric prefix (digits and dots) at the very start
    match = re.match(r'^([\d.]+)', code_str)
    if match:
        return match.group(1)
    
    # If code starts with letters/underscores, try to find numeric component after last underscore
    # E.g., "hiv_01" -> "01"
    match = re.search(r'_(\d+)$', code_str)
    if match:
        return match.group(1)
    
    # If code contains numeric component anywhere after non-numeric prefix
    # E.g., "HIV01" -> "01"
    match = re.search(r'(\d+)$', code_str)
    if match:
        return match.group(1)
    
    # Fallback: return the whole code if no numeric component found
    return code_str

def format_intervention_label(code, label):
    """Format intervention display as 'numeric_code. label' avoiding duplication."""
    numeric = extract_numeric_code(code)
    label_str = str(label).strip()
    code_str = str(code).strip()
    # If label is the same as code (no metadata), just return numeric code
    if label_str == code_str:
        return numeric
    # Otherwise show numeric code prefix with label name
    return f"{numeric}. {label_str}"

def hiv_report_label(code, label, component=""):
    """HIV report display label with a special-case kludge for generic activity names."""
    import re
    label_str = str(label).strip()
    component_str = str(component).strip()
    label_norm = re.sub(r"\s+", "", label_str.lower())
    if label_norm == "medicamentos/insumos" and component_str:
        label_str = f"{label_str} — {component_str}"
    return format_intervention_label(code, label_str)

def svg_score_strip(all_vals, current_val, color, width=168, height=28):
    """Horizontal strip showing score distribution with current value marked."""
    if not all_vals:
        return ""
    mn = min(all_vals)
    mx = max(all_vals)
    span = (mx - mn) or 1.0
    pad = 6
    track_w = width - 2 * pad
    cy = 10
    def xp(v):
        return pad + (v - mn) / span * track_w
    sv = sorted(all_vals)
    n = len(sv)
    q25 = sv[max(0, n // 4)]
    q75 = sv[min(n - 1, 3 * n // 4)]
    cur_x = xp(current_val)
    q25_x = xp(q25)
    q75_x = xp(q75)
    out = [f'<svg width="{width}" height="{height}" style="display:inline-block;vertical-align:middle;overflow:visible">']
    out.append(f'<rect x="{pad}" y="{cy-2}" width="{track_w}" height="4" rx="2" fill="#e5e7eb"/>')
    iqr_w = max(0.0, q75_x - q25_x)
    out.append(f'<rect x="{q25_x:.1f}" y="{cy-2}" width="{iqr_w:.1f}" height="4" rx="1" fill="#d1d5db"/>')
    for v in sv:
        x = xp(v)
        out.append(f'<line x1="{x:.1f}" y1="{cy-5}" x2="{x:.1f}" y2="{cy+5}" stroke="#9ca3af" stroke-width="0.8" stroke-opacity="0.5"/>')
    out.append(f'<circle cx="{cur_x:.1f}" cy="{cy}" r="5" fill="{color}" stroke="white" stroke-width="1.5"/>')
    out.append(f'<text x="{cur_x:.1f}" y="{cy+18}" text-anchor="middle" font-size="9" fill="{color}" font-weight="600">{current_val:.3f}</text>')
    out.append('</svg>')
    return ''.join(out)


def build_integration_candidates_html(results):
    """Build the 'Integration Candidates' section.

    For each ordered pair (A, B) where A has B in intg_counts, compute
    min_mencoes = min(A→B count, B→A count).  Only pairs where both directions
    have at least one mention are included.  Groups are formed by distinct
    min_mencoes values; the top 3 distinct values are shown.
    """
    import re as _re2
    if not results:
        return ""

    # Collect all bidirectional integration mention counts
    codes = set(results.keys())
    pairs = {}   # (a, b) with a < b  →  {"a_to_b": int, "b_to_a": int}
    for code, r in results.items():
        for target, cnt in r.get("intg_counts", {}).items():
            if target not in codes:
                continue
            key = tuple(sorted([code, target]))
            if key not in pairs:
                pairs[key] = {"a_to_b": 0, "b_to_a": 0}
            if code == key[0]:
                pairs[key]["a_to_b"] += cnt
            else:
                pairs[key]["b_to_a"] += cnt

    # Only keep pairs with at least one mention in each direction
    mutual = {}
    for (a, b), d in pairs.items():
        if d["a_to_b"] > 0 and d["b_to_a"] > 0:
            mutual[(a, b)] = {
                "a_to_b": d["a_to_b"],
                "b_to_a": d["b_to_a"],
                "min_m":  min(d["a_to_b"], d["b_to_a"]),
                "total":  d["a_to_b"] + d["b_to_a"],
            }

    if not mutual:
        return ""

    # Find top 3 distinct min_mencoes levels
    distinct_levels = sorted(set(v["min_m"] for v in mutual.values()), reverse=True)[:3]

    def _lbl(code):
        r = results.get(code, {})
        return r.get("label", code)

    def _comp(code):
        return (results.get(code, {}).get("component") or "").strip()

    def _prog(code):
        return (results.get(code, {}).get("programa") or "").strip()

    def _num(code):
        m = _re2.search(r'(\d+)$', str(code))
        return m.group(1) if m else code

    html = """
  <div class="section-header">
    <h2>Candidatos à Integração</h2>
    <span class="subtitle">Pares com menções mútuas mais frequentes (mín. menções em ambas as direcções)</span>
  </div>
  <div style="font-size:12px;color:var(--muted);margin-bottom:16px">
    Apenas pares em que ambas as intervenções foram identificadas mutuamente pelos especialistas são incluídos.
    <strong>Mín. menções</strong> = mínimo das contagens nas duas direcções (A→B e B→A).
  </div>
"""
    tier_colors = ["#1a6b3a", "#1a5276", "#7a5c00"]
    tier_labels = ["Prioridade 1 — menções mútuas mais altas",
                   "Prioridade 2",
                   "Prioridade 3"]

    for tier_idx, level in enumerate(distinct_levels):
        tier_pairs = sorted(
            [(k, v) for k, v in mutual.items() if v["min_m"] == level],
            key=lambda x: (-x[1]["total"], x[0])
        )
        col = tier_colors[tier_idx % len(tier_colors)]
        html += f"""
  <div style="margin-bottom:28px">
    <div style="font-size:11px;font-weight:700;text-transform:uppercase;letter-spacing:.1em;
                color:{col};border-left:3px solid {col};padding-left:8px;margin-bottom:10px">
      {esc(tier_labels[tier_idx])} &nbsp;·&nbsp; mín. menções = {level}
    </div>
    <div style="display:flex;flex-direction:column;gap:8px">
"""
        for (a, b), d in tier_pairs:
            a_comp = _comp(a)
            b_comp = _comp(b)
            a_prog = _prog(a)
            b_prog = _prog(b)
            comp_note_a = f" <span style='color:var(--muted);font-size:10px'>({esc(a_comp)})</span>" if a_comp else ""
            comp_note_b = f" <span style='color:var(--muted);font-size:10px'>({esc(b_comp)})</span>" if b_comp else ""
            prog_note_a = f"<span style='font-size:9px;color:var(--muted)'>{esc(a_prog)}</span><br>" if a_prog else ""
            prog_note_b = f"<span style='font-size:9px;color:var(--muted)'>{esc(b_prog)}</span><br>" if b_prog else ""
            html += f"""      <div style="background:var(--card-bg);border:1px solid var(--border);border-radius:4px;
                          padding:10px 14px;display:flex;align-items:flex-start;gap:16px">
        <div style="flex:1;min-width:0">
          <div style="font-weight:600;font-size:12px">
            <a href="#card-{_re2.sub(r'[^A-Za-z0-9_-]','-',a)}" class="inv-link">{_num(a)}. {esc(_lbl(a))}</a>{comp_note_a}
          </div>
          <div style="margin-top:1px">{prog_note_a}</div>
        </div>
        <div style="display:flex;flex-direction:column;align-items:center;gap:2px;flex-shrink:0;
                    font-size:10px;color:var(--muted);padding-top:2px">
          <span title="{esc(_lbl(a))} → {esc(_lbl(b))}: {d['a_to_b']} menções">→ {d['a_to_b']}</span>
          <span style="font-size:16px;color:{col};line-height:1">⇄</span>
          <span title="{esc(_lbl(b))} → {esc(_lbl(a))}: {d['b_to_a']} menções">← {d['b_to_a']}</span>
          <span style="font-weight:700;color:{col}">mín={level}</span>
        </div>
        <div style="flex:1;min-width:0;text-align:right">
          <div style="font-weight:600;font-size:12px">
            <a href="#card-{_re2.sub(r'[^A-Za-z0-9_-]','-',b)}" class="inv-link">{_num(b)}. {esc(_lbl(b))}</a>{comp_note_b}
          </div>
          <div style="margin-top:1px;text-align:right">{prog_note_b}</div>
        </div>
      </div>
"""
        html += "    </div>\n  </div>\n"

    return html


def build_network_html(results, topic_label="HIV"):
    """Build an interactive D3 v7 force-directed network of dup/intg relationships.

    Returns an HTML string (with an embedded <script>) ready to drop into the
    report body, or an empty string when there are no relationships to show.
    """
    # ── Build node list ────────────────────────────────────────────────────
    nodes = [
        {
            "id":        r["code"],
            "label":     r["label"],
            "programa":  str(r.get("programa", "") or "Outros").strip() or "Outros",
            "component": str(r.get("component", "") or "").strip(),
            "s_pond":    round(r.get("s_pond", 0), 3),
            "n_total":   r.get("n_total", 0),
        }
        for r in results.values()
    ]

    # ── Symmetrize edges (A→B + B→A → single undirected edge) ─────────────
    edge_map = {}
    for code, r in results.items():
        for target, cnt in r.get("dup_counts", {}).items():
            key = tuple(sorted([code, target]))
            e = edge_map.setdefault(key, {"dup_w": 0, "intg_w": 0})
            e["dup_w"] += cnt
        for target, cnt in r.get("intg_counts", {}).items():
            key = tuple(sorted([code, target]))
            e = edge_map.setdefault(key, {"dup_w": 0, "intg_w": 0})
            e["intg_w"] += cnt

    links = []
    for (src, tgt), w in edge_map.items():
        total = w["dup_w"] + w["intg_w"]
        if total == 0:
            continue
        if w["dup_w"] > 0 and w["intg_w"] == 0:
            etype = "dup"
        elif w["intg_w"] > 0 and w["dup_w"] == 0:
            etype = "intg"
        else:
            etype = "both"
        links.append({
            "source":  src,
            "target":  tgt,
            "dup_w":   w["dup_w"],
            "intg_w":  w["intg_w"],
            "total_w": total,
            "type":    etype,
        })

    if not links:
        return ""

    max_w = max(l["total_w"] for l in links)
    slider_max = min(max_w, 20)

    nodes_json = json.dumps(nodes, ensure_ascii=False)
    links_json = json.dumps(links, ensure_ascii=False)

    return f"""
  <!-- ═══ NETWORK GRAPH ════════════════════════════════════════════════════ -->
  <div class="section-header">
    <h2>Rede de Duplicação / Integração</h2>
    <span class="subtitle">Nós = intervenções · Arestas = relações reportadas pelos especialistas</span>
  </div>
  <div style="font-size:12px;color:var(--muted);margin-bottom:10px">
    Cada aresta liga duas intervenções quando pelo menos um especialista
    identificou duplicação ou potencial de integração entre elas.
    A <strong>espessura</strong> da aresta reflecte o número total de menções;
    a <strong>cor</strong> indica o tipo dominante.
    Os nós são coloridos por programa e o seu tamanho reflecte S<sub>pond</sub>.
    Arraste os nós para reorganizar; utilize a roda do rato para zoom.
    <strong>Clique num nó</strong> para navegar para a ficha detalhada da intervenção.
  </div>
  <div style="display:flex;align-items:center;gap:18px;flex-wrap:wrap;margin-bottom:6px">
    <label style="font-size:12px;display:flex;align-items:center;gap:6px">
      Mín. menções:
      <input id="net-thresh" type="range" min="1" max="{slider_max}" value="2" step="1"
             style="width:100px;accent-color:var(--ink)">
      <span id="net-thresh-val" style="font-weight:600;min-width:18px">2</span>
    </label>
    <label style="font-size:12px;display:flex;align-items:center;gap:6px">
      Relação:
      <select id="net-rel-type" style="font-size:12px;border:1px solid var(--border);border-radius:4px;padding:2px 6px;background:var(--card-bg)">
        <option value="dup">Duplicação</option>
        <option value="intg">Integração</option>
        <option value="both" selected>Ambos</option>
      </select>
    </label>
    <button id="net-recenter"
            style="font-size:11px;padding:3px 10px;border:1px solid var(--border);
                   border-radius:4px;cursor:pointer;background:var(--card-bg)">
      ⤢ Ajustar ao ecrã
    </button>
    <span style="font-size:11px;display:flex;align-items:center;gap:10px">
      <svg width="28" height="10"><line x1="0" y1="5" x2="28" y2="5" stroke="#dc2626" stroke-width="3"/></svg>Duplicação
      <svg width="28" height="10"><line x1="0" y1="5" x2="28" y2="5" stroke="#1d4ed8" stroke-width="3"/></svg>Integração
      <svg width="28" height="10"><line x1="0" y1="5" x2="28" y2="5" stroke="#7c3aed" stroke-width="3"/></svg>Ambos
    </span>
  </div>
  <div style="font-size:11px;color:var(--muted);margin:0 0 8px 0">
    Nota: apenas ligações com outras intervenções de {esc(topic_label)} são incluídas;
    relações reportadas em campo "other" não estão reflectidas neste diagrama.
  </div>
  <div id="net-prog-legend" style="margin-bottom:8px;font-size:11px;display:flex;flex-wrap:wrap;gap:4px"></div>
  <div id="network-chart"
       style="width:100%;height:600px;border:1px solid var(--border);
              border-radius:8px;background:#fff;position:relative;overflow:hidden">
  </div>

  <script src="https://cdn.jsdelivr.net/npm/d3@7/dist/d3.min.js"></script>
  <script>
  (function() {{
    const ALL_NODES = {nodes_json};
    const ALL_LINKS = {links_json};
    const MAX_W = {max_w};

    const TYPE_COLOR = {{dup: '#dc2626', intg: '#1d4ed8', both: '#7c3aed'}};
    const PALETTE = ['#1e7e34','#1a5276','#9b2226','#7c5a00','#005f73',
                     '#732d91','#41690c','#2d5a7b','#724c11','#3d6b61',
                     '#b45309','#0e7490','#6d28d9'];
    const programs = [...new Set(ALL_NODES.map(n => n.programa).filter(Boolean))].sort();
    const progColor = {{}};
    programs.forEach((p, i) => {{ progColor[p] = PALETTE[i % PALETTE.length]; }});

    const container = document.getElementById('network-chart');
    const W = container.offsetWidth || 900;
    const H = 600;

    // Scales
    const strokeW = d3.scaleSqrt().domain([1, MAX_W]).range([1.5, 9]);
    const maxS = Math.max(0.01, d3.max(ALL_NODES, d => d.s_pond));
    const rScale = d3.scaleSqrt().domain([0, maxS]).range([6, 16]);

    // SVG
    const svg = d3.select(container)
      .append('svg')
      .attr('width', W).attr('height', H)
      .style('font-family', "'DM Sans', sans-serif");

    const zoomBeh = d3.zoom().scaleExtent([0.15, 6])
      .on('zoom', ev => gMain.attr('transform', ev.transform));
    svg.call(zoomBeh);

    const gMain = svg.append('g');

    function fitGraph(transitionMs=400) {{
      const node = gMain.node();
      if (!node || !node.childNodes || node.childNodes.length === 0) return;
      let bounds;
      try {{
        bounds = node.getBBox();
      }} catch (_err) {{
        return;
      }}
      if (!bounds || bounds.width <= 0 || bounds.height <= 0) return;

      const pad = 28;
      const scale = Math.max(
        0.15,
        Math.min(6, 0.92 / Math.max(bounds.width / (W - pad * 2), bounds.height / (H - pad * 2)))
      );
      const tx = W / 2 - scale * (bounds.x + bounds.width / 2);
      const ty = H / 2 - scale * (bounds.y + bounds.height / 2);
      const targetTransform = d3.zoomIdentity.translate(tx, ty).scale(scale);

      svg.transition().duration(transitionMs)
         .call(zoomBeh.transform, targetTransform);
    }}

    // Floating tooltip (HTML)
    const tip = d3.select(container)
      .append('div')
      .style('position', 'absolute')
      .style('pointer-events', 'none')
      .style('background', 'rgba(15,25,35,0.92)')
      .style('color', '#f3f4f6')
      .style('border-radius', '6px')
      .style('padding', '8px 12px')
      .style('font-size', '12px')
      .style('line-height', '1.5')
      .style('max-width', '260px')
      .style('display', 'none')
      .style('z-index', '10');

    let simObj = null;

    function rebuildGraph(minW, relType) {{
      const fitLinks = ALL_LINKS
        .filter(d => d.type === relType)
        .filter(d => d.total_w >= minW)
        .map(d => ({{...d}}));
      const visIds = new Set(fitLinks.flatMap(d => [d.source, d.target]));
      const fitNodes = ALL_NODES
        .filter(n => visIds.has(n.id))
        .map(d => ({{...d}}));

      gMain.selectAll('*').remove();
      if (simObj) simObj.stop();

      if (fitLinks.length === 0) {{
        gMain.append('text')
          .attr('x', W / 2).attr('y', H / 2)
          .attr('text-anchor', 'middle')
          .attr('fill', '#9ca3af')
          .text('Sem relações com os filtros seleccionados.');
        return;
      }}

      simObj = d3.forceSimulation(fitNodes)
        .force('link', d3.forceLink(fitLinks).id(d => d.id)
          .distance(d => 60 + 100 / Math.sqrt(d.total_w)))
        .force('charge', d3.forceManyBody().strength(-200))
        .force('center', d3.forceCenter(W / 2, H / 2))
        .force('collide', d3.forceCollide().radius(d => rScale(d.s_pond) + 8));

      // Links
      const linkSels = gMain.append('g')
        .selectAll('line').data(fitLinks).join('line')
        .attr('stroke', d => TYPE_COLOR[d.type])
        .attr('stroke-width', d => strokeW(d.total_w))
        .attr('stroke-opacity', 0.65)
        .on('mouseover', (ev, d) => {{
          const srcId = d.source.id || d.source;
          const tgtId = d.target.id || d.target;
          tip.style('display', 'block').html(
            '<strong>' + srcId + ' ↔ ' + tgtId + '</strong><br>' +
            '<span style="color:#fca5a5">Duplicação: ' + d.dup_w + '</span><br>' +
            '<span style="color:#93c5fd">Integração: ' + d.intg_w + '</span><br>' +
            '<em style="color:#d1d5db">Total menções: ' + d.total_w + '</em>'
          );
        }})
        .on('mousemove', (ev) => {{
          const rect = container.getBoundingClientRect();
          tip.style('left', (ev.clientX - rect.left + 14) + 'px')
             .style('top',  (ev.clientY - rect.top  - 10) + 'px');
        }})
        .on('mouseout', () => tip.style('display', 'none'));

      // Nodes
      const nodeSels = gMain.append('g')
        .selectAll('circle').data(fitNodes).join('circle')
        .attr('r', d => rScale(d.s_pond))
        .attr('fill', d => progColor[d.programa] || '#888')
        .attr('stroke', '#fff').attr('stroke-width', 1.5)
        .style('cursor', 'pointer')
        .call(d3.drag()
          .on('start', (ev, d) => {{
            if (!ev.active) simObj.alphaTarget(0.3).restart();
            d.fx = d.x; d.fy = d.y;
            d._dragMoved = false;
          }})
          .on('drag', (ev, d) => {{
            d.fx = ev.x; d.fy = ev.y;
            d._dragMoved = true;
          }})
          .on('end', (ev, d) => {{
            if (!ev.active) simObj.alphaTarget(0);
            d.fx = null; d.fy = null;
          }}))
        .on('click', (ev, d) => {{
          if (d._dragMoved) {{ d._dragMoved = false; return; }}
          const safeId = d.id.replace(/[^A-Za-z0-9_-]/g, '-');
          const target = document.getElementById('card-' + safeId);
          if (target) {{
            target.scrollIntoView({{behavior: 'smooth', block: 'start'}});
            target.style.outline = '2px solid #1a5276';
            setTimeout(() => {{ target.style.outline = ''; }}, 1800);
          }}
        }})
        .on('mouseover', (ev, d) => {{
          tip.style('display', 'block').html(
            '<strong>' + d.id + '</strong><br>' +
            d.label + '<br>' +
            (d.component ? '<em style="color:#9ca3af;font-size:10px">' + d.component + '</em><br>' : '') +
            '<em style="color:#d1d5db">' + d.programa + '</em><br>' +
            'S\u209a\u2092\u2099\u2091: ' + d.s_pond.toFixed(2) + ' &nbsp;|&nbsp; N resp.: ' + d.n_total
          );
        }})
        .on('mousemove', (ev) => {{
          const rect = container.getBoundingClientRect();
          tip.style('left', (ev.clientX - rect.left + 14) + 'px')
             .style('top',  (ev.clientY - rect.top  - 10) + 'px');
        }})
        .on('mouseout', () => tip.style('display', 'none'));

      // Labels (intervention code below node)
      const labelSels = gMain.append('g')
        .selectAll('text').data(fitNodes).join('text')
        .attr('font-size', '8px')
        .attr('fill', '#111')
        .attr('text-anchor', 'middle')
        .attr('pointer-events', 'none')
        .text(d => d.id);

      simObj.on('tick', () => {{
        linkSels
          .attr('x1', d => d.source.x).attr('y1', d => d.source.y)
          .attr('x2', d => d.target.x).attr('y2', d => d.target.y);
        nodeSels
          .attr('cx', d => d.x).attr('cy', d => d.y);
        labelSels
          .attr('x', d => d.x)
          .attr('y', d => d.y + rScale(d.s_pond) + 10);
      }});

      simObj.on('end', () => fitGraph(500));
      setTimeout(() => fitGraph(500), 700);
    }}

    // Controls
    const slider   = document.getElementById('net-thresh');
    const sliderLbl = document.getElementById('net-thresh-val');
    const relTypeSel = document.getElementById('net-rel-type');
    rebuildGraph(parseInt(slider.value, 10), relTypeSel.value);
    slider.addEventListener('input', function () {{
      sliderLbl.textContent = this.value;
      rebuildGraph(parseInt(this.value, 10), relTypeSel.value);
    }});
    relTypeSel.addEventListener('change', function () {{
      rebuildGraph(parseInt(slider.value, 10), this.value);
    }});

    document.getElementById('net-recenter').addEventListener('click', () => {{
      fitGraph(450);
    }});

    // Programa colour legend
    const legEl = document.getElementById('net-prog-legend');
    programs.forEach(p => {{
      const s = document.createElement('span');
      s.style.cssText = 'display:inline-flex;align-items:center;gap:4px;margin-right:8px';
      s.innerHTML = '<svg width="10" height="10" viewBox="0 0 10 10">' +
        '<circle cx="5" cy="5" r="5" fill="' + progColor[p] + '"/></svg>' + p;
      legEl.appendChild(s);
    }});

  }})();
  </script>
"""


def render_html(results, stats, interventions, source_file, univariate=None,
                include_xyplot=True, include_scoring=True,
                alluvial_top=None, priority_order="s_pond", programa_nome=None,
                topic_label="HIV", simple_table=False, logo_data_uri=None):
    if programa_nome is None:
        programa_nome = "Programa Nacional de Controlo do HIV/SIDA"
    inv_label = {r["code"]: r["label"] for r in interventions}
    # Create display labels with numeric code prefix
    display_label = {r["code"]: hiv_report_label(r["code"], r["label"], r.get("component", "")) 
                     for r in interventions}
    _sort_key = (lambda r: r["s_base"]) if priority_order == "s_base" else (lambda r: r["s_pond"])
    sorted_inv = sorted(results.values(), key=_sort_key, reverse=True)
    all_s_base_vals = sorted(r["s_base"] for r in results.values())
    all_s_pond_vals = sorted(r["s_pond"] for r in results.values())
    _sort_label = "S_base (não-ponderado)" if priority_order == "s_base" else "S_pond (ponderado por experiência)"
    for i, r in enumerate(sorted_inv, 1):
        r["display_idx"] = i
    now = datetime.now().strftime("%d/%m/%Y às %H:%M")
    med_gate = statistics.median([r.get("gate_mean", 0) for r in sorted_inv]) if sorted_inv else 0
    med_imp = statistics.median([r.get("avg_impact", 0) for r in sorted_inv]) if sorted_inv else 0

    def tier(r):
        high_gate = r.get("gate_mean", 0) > med_gate
        high_imp = r.get("avg_impact", 0) > med_imp
        if high_gate and high_imp: return "alta"
        if high_gate or high_imp: return "media"
        return "baixa"

    def _mean(vals):
        vals = [v for v in vals if v is not None]
        return (sum(vals) / len(vals)) if vals else 0.0

    def _score_fill_color(v):
        if v >= 1.0:
            return "#1a6b3a"
        if v >= 0.7:
            return "#d4a017"
        return "#9ca3af"

    def _e_fill_color(p):
        if p >= 75:
            return "#1e7e34"
        if p >= 50:
            return "#f59e0b"
        return "#9ca3af"

    def _safe_anchor(text):
      import re as _re
      return _re.sub(r'[^A-Za-z0-9_-]', '-', str(text or "Outros"))

    # Prepare card grouping: programa → component → [results], preserving catalog order.
    by_prog_comp = OrderedDict()   # {programa: OrderedDict({component: [result]})}
    programa_first_anchor = {}     # {programa: anchor_id}  — for summary table links
    comp_macro_first_anchor = {}   # {comp_macro: anchor_id} — for component table links
    for inv in interventions:
      prog = (inv.get("programa", "") or "Outros").strip() or "Outros"
      comp = (inv.get("component", "") or "Outros").strip() or "Outros"
      if prog not in by_prog_comp:
        by_prog_comp[prog] = OrderedDict()
        prog_anchor = f"prog-section-{_safe_anchor(prog)}"
        programa_first_anchor[prog] = prog_anchor
      if comp not in by_prog_comp[prog]:
        by_prog_comp[prog][comp] = []
      code = inv.get("code")
      if code in results:
        row = results[code]
        by_prog_comp[prog][comp].append(row)
        macro = str(row.get("comp_macro", "") or "Outros").strip() or "Outros"
        if macro not in comp_macro_first_anchor:
          comp_macro_first_anchor[macro] = f"prog-section-{_safe_anchor(prog)}"

    def build_group_summary_rows(group_key, name_links=None):
        buckets = OrderedDict()
        for r in results.values():
            gname = str(r.get(group_key, "") or "Outros").strip() or "Outros"
            buckets.setdefault(gname, []).append(r)

        rows = []
        for gname, items_g in buckets.items():
            s_vals = [x.get("s_pond", 0) for x in items_g]
            s_pond_m = _mean(s_vals)
            s_pond_max = max(s_vals) if s_vals else 0
            impact_m = _mean([x.get("avg_impact", 0) for x in items_g])
            e_m = _mean([x.get("e_score", 0) for x in items_g])
            rows.append({
                "name":      gname,
                "n":         len(items_g),
                "s_pond":    s_pond_m,
                "s_pond_max": s_pond_max,
                "impact":    impact_m,
                "e":         e_m,
            })
        rows.sort(key=lambda x: x["s_pond"], reverse=True)

        html_rows = ""
        for row in rows:
            s_val = row["s_pond"]
            s_max = row["s_pond_max"]
            s_w = max(0, min(100, round(s_val / 5 * 100)))
            s_col = _score_fill_color(s_val)
            e_pct = round(row["e"] * 100)
            e_w = max(0, min(100, e_pct))
            e_col = _e_fill_color(e_pct)
            name_html = esc(row['name'])
            if name_links and row["name"] in name_links:
                name_html = f'<a href="#{name_links[row["name"]]}" class="inv-link">{esc(row["name"])}</a>'
            html_rows += f"""<tr>
              <td style=\"font-weight:500\">{name_html}</td>
              <td style=\"text-align:center;color:var(--muted)\" data-val=\"{row['n']}\">{row['n']}</td>
              <td data-val=\"{s_val:.4f}\">
                <div style=\"display:flex;align-items:center;gap:8px\">
                  <div style=\"flex:1;height:6px;background:#e8edf2;border-radius:3px;overflow:hidden\">
                    <div style=\"width:{s_w}%;height:100%;background:{s_col}\"></div>
                  </div>
                  <span style=\"font-size:11px;font-weight:600;color:#374151\">{s_val:.2f}</span>
                </div>
              </td>
              <td style=\"text-align:center;color:var(--muted);font-size:11px\" data-val=\"{s_max:.4f}\">{s_max:.2f}</td>
              <td style=\"text-align:center\" data-val=\"{row['impact']:.4f}\">{row['impact']:.2f}</td>
              <td data-val=\"{e_pct}\">
                <div style=\"display:flex;align-items:center;gap:8px\">
                  <div style=\"width:50px;height:5px;background:#e8edf2;border-radius:2px;overflow:hidden\">
                    <div style=\"width:{e_w}%;height:100%;background:{e_col}\"></div>
                  </div>
                  <span style=\"font-size:11px;color:#374151\">{e_pct}%</span>
                </div>
              </td>
            </tr>"""
        return html_rows

    programa_rows = build_group_summary_rows("programa", name_links=programa_first_anchor)
    comp_rows = build_group_summary_rows("comp_macro", name_links=comp_macro_first_anchor)
    _sort_th = 'style="cursor:pointer;user-select:none;white-space:nowrap" onclick="sortSummaryTable(this)"'
    _sort_th_c = 'style="cursor:pointer;user-select:none;white-space:nowrap;text-align:center" onclick="sortSummaryTable(this)"'
    summaries_html = f"""
  <div class=\"section-header\">
    <h2>Resumo por Programa</h2>
    <span class=\"subtitle\">Médias de S<sub>pond</sub>, impacto e E por área programática · clique num cabeçalho para ordenar</span>
  </div>
  <table class=\"rank-table\" id=\"tbl-prog\">
    <thead>
      <tr>
        <th {_sort_th}>Programa</th>
        <th {_sort_th_c}>N</th>
        <th {_sort_th}>S<sub>pond</sub> médio ▾</th>
        <th {_sort_th_c}>S<sub>pond</sub> máx</th>
        <th {_sort_th_c}>Impacto pond.</th>
        <th {_sort_th}>E (eficiência)</th>
      </tr>
    </thead>
    <tbody>{programa_rows}</tbody>
  </table>

  <div class=\"section-header\">
    <h2>Resumo por Componente</h2>
    <span class=\"subtitle\">Agrupamento macro de componente com os mesmos indicadores · clique num cabeçalho para ordenar</span>
  </div>
  <table class=\"rank-table\" id=\"tbl-comp\">
    <thead>
      <tr>
        <th {_sort_th}>Componente</th>
        <th {_sort_th_c}>N</th>
        <th {_sort_th}>S<sub>pond</sub> médio ▾</th>
        <th {_sort_th_c}>S<sub>pond</sub> máx</th>
        <th {_sort_th_c}>Impacto pond.</th>
        <th {_sort_th}>E (eficiência)</th>
      </tr>
    </thead>
    <tbody>{comp_rows}</tbody>
  </table>
  <script>
  (function() {{
    function sortSummaryTable(th) {{
      var table = th.closest('table');
      var tbody = table.querySelector('tbody');
      var ths = Array.from(th.parentElement.children);
      var col = ths.indexOf(th);
      var asc = th.dataset.sortDir !== 'asc';
      th.dataset.sortDir = asc ? 'asc' : 'desc';
      ths.forEach(function(h) {{ h.textContent = h.textContent.replace(' ▲','').replace(' ▼',''); }});
      th.textContent = th.textContent + (asc ? ' ▲' : ' ▼');
      var rows = Array.from(tbody.querySelectorAll('tr'));
      rows.sort(function(a, b) {{
        var ca = a.children[col], cb = b.children[col];
        var va = parseFloat(ca.dataset.val !== undefined ? ca.dataset.val : ca.textContent);
        var vb = parseFloat(cb.dataset.val !== undefined ? cb.dataset.val : cb.textContent);
        if (isNaN(va)) va = ca.textContent.trim().toLowerCase();
        if (isNaN(vb)) vb = cb.textContent.trim().toLowerCase();
        if (va < vb) return asc ? -1 : 1;
        if (va > vb) return asc ? 1 : -1;
        return 0;
      }});
      rows.forEach(function(r) {{ tbody.appendChild(r); }});
    }}
    window.sortSummaryTable = sortSummaryTable;
  }})();
  </script>
"""

    alta_rows  = [r for r in sorted_inv if tier(r) == "alta"]
    media_rows = [r for r in sorted_inv if tier(r) == "media"]
    baixa_rows = [r for r in sorted_inv if tier(r) == "baixa"]

    # ── Response rate section ──────────────────────────────────────────────
    rr_rows = ""
    n_experts = stats["n_experts"]
    n_experts_observed = stats.get("n_experts_observed", n_experts)
    for item in stats["rr_detail"]:
        bar_w = round(item["pct"])
        color = "#1a6b3a" if item["pct"] >= 80 else ("#d4a017" if item["pct"] >= 60 else "#e57373")
        rr_rows += f"""<tr>
          <td style="max-width:320px;font-size:12px">{esc(item["label"])}</td>
          <td style="text-align:center;font-weight:600">{item["n"]}/{n_experts}</td>
          <td style="width:200px">
            <div style="display:flex;align-items:center;gap:8px">
              <div style="flex:1;height:8px;background:#e5e7eb;border-radius:2px;overflow:hidden">
                <div style="width:{bar_w}%;height:100%;background:{color}"></div>
              </div>
              <span style="font-size:11px;font-weight:600;color:{color};white-space:nowrap">{item["pct"]}%</span>
            </div>
          </td>
        </tr>"""

    # ── Team response rate section (only when teams present) ───────────────
    _team_colors = {"A": "#1a5c8a", "B": "#2e7d52", "C": "#b45309", "D": "#6b3fa0"}
    rr_by_team_html = ""
    rr_by_team_data = stats.get("rr_by_team", [])
    if rr_by_team_data:
        team_cards = ""
        for td in rr_by_team_data:
            tc = _team_colors.get(td["team"], "#546e7a")
            team_cards += (
                f"<div style='background:#fff;border:1px solid #ECEFF1;border-radius:8px;"
                f"padding:14px 20px;text-align:center;border-top:3px solid {tc};min-width:120px'>"
                f"<div style='font-size:0.7rem;text-transform:uppercase;letter-spacing:0.08em;"
                f"color:#78909C;font-weight:600'>Equipa {esc(td['team'])}</div>"
                f"<div style='font-size:1.6rem;font-weight:700;color:{tc}'>{td['rr_median']}%</div>"
                f"<div style='font-size:0.75rem;color:#78909C'>"
                f"{td['n_inv']} interv. · {td.get('n_experts','?')} esp. · {td['rr_min']}–{td['rr_max']}%</div>"
                f"</div>"
            )
        rr_by_team_html = (
            f"<div style='margin-top:18px'>"
            f"<div style='font-size:12px;font-weight:600;color:#546e7a;margin-bottom:10px;"
            f"text-transform:uppercase;letter-spacing:0.06em'>Taxa de resposta por equipa</div>"
            f"<div style='display:flex;gap:14px;flex-wrap:wrap'>{team_cards}</div>"
            f"</div>"
        )

    # ── Ranking table rows ──────────────────────────────────────────────────
    rank_rows = ""
    for i, r in enumerate(sorted_inv, 1):
        sim_w  = round(r["n_sim"]  / r["n_total"] * 100) if r["n_total"] else 0
        poss_w = round(r["n_poss"] / r["n_total"] * 100) if r["n_total"] else 0
        nao_w  = round(r["n_nao"]  / r["n_total"] * 100) if r["n_total"] else 0
        imp    = r["avg_impact"]
        imp_cls = "imp-high" if imp >= 2.5 else ("imp-med" if imp >= 1.8 else "imp-low")
        imp_lbl = "Alto" if imp >= 2.5 else ("Médio" if imp >= 1.8 else "Baixo")
        label_with_idx = display_label.get(r["code"], r["label"])
        import re as _re
        _safe_id_rank = _re.sub(r'[^A-Za-z0-9_-]', '-', str(r["code"]))
        name_cell = f'<a href="#card-{_safe_id_rank}" class="inv-link">{esc(label_with_idx)}</a>'
        rank_rows += f"""<tr>
          <td class="rank-num">{i}</td>
          <td style="font-weight:500">{name_cell}</td>
          <td><span class="comp-tag">{esc(r["component"])}</span></td>
          <td>
            <div class="mini-bar-wrap">
              <div class="mini-bar-bg">
                <div style="display:flex;height:100%">
                  <div style="width:{sim_w}%;background:var(--sim)"></div>
                  <div style="width:{poss_w}%;background:#d4a017"></div>
                  <div style="width:{nao_w}%;background:#e5e7eb"></div>
                </div>
              </div>
              <span class="bar-counts">{r["n_sim"]}S·{r["n_poss"]}P·{r["n_nao"]}N</span>
            </div>
          </td>
          <td><strong style="font-size:14px">{r["pct_optimizable"]}%</strong></td>
          <td><span class="impact-badge {imp_cls}">{imp:.1f} — {imp_lbl}</span></td>
          <td style="color:var(--muted)">{r["dup_pct"]}%</td>
          <td style="color:var(--muted)">{r["intg_pct"]}%</td>
        </tr>"""

    cards_html = ""
    cmt_toggle_seq = 0
    for prog, comp_dict in by_prog_comp.items():
        prog_anchor = programa_first_anchor.get(prog, f"prog-section-{_safe_anchor(prog)}")
        cards_html += (
            f'<div class="programa-section-header" id="{prog_anchor}" '
            f'style="margin:32px 0 10px;padding:10px 16px;background:var(--accent2);'
            f'color:#fff;border-radius:6px;font-family:\'DM Serif Display\',serif;'
            f'font-size:1.05rem;letter-spacing:0.01em">{esc(prog)}</div>\n'
        )
        for comp, items in comp_dict.items():
          if not items:
              continue
          cards_html += f'<div class="comp-label">{esc(comp)}</div>\n<div class="priority-grid">\n'
          for r in items:
            sim_p  = round(r["n_sim"]  / r["n_total"] * 100) if r["n_total"] else 0
            poss_p = round(r["n_poss"] / r["n_total"] * 100) if r["n_total"] else 0
            nao_p  = round(r["n_nao"]  / r["n_total"] * 100) if r["n_total"] else 0
            imp    = r["avg_impact"]
            imp_cls = "mc-impact-high" if imp >= 2.5 else ("mc-impact-med" if imp >= 1.8 else "mc-impact-low")

            def make_tag(code, cls):
                lbl = display_label.get(code, inv_label.get(code, code))
                import re as _re
                safe_id = _re.sub(r'[^A-Za-z0-9_-]', '-', str(code))
                return f'<a href="#card-{safe_id}" class="{cls}" style="text-decoration:none">{esc(lbl)}</a>'

            dup_tags  = "".join(make_tag(t, "dup-tag")  for t in r["top_dup"])
            intg_tags = "".join(make_tag(t, "intg-tag") for t in r["top_intg"])
            dup_row  = (f'<div class="tag-row"><span class="tag-label dup-label">Duplicação com:</span>'
                        f'{dup_tags}</div>') if r["top_dup"]  else ""
            intg_row = (f'<div class="tag-row"><span class="tag-label">Integrar com:</span>'
                        f'{intg_tags}</div>') if r["top_intg"] else ""

            dup_other_row = ""
            intg_other_row = ""
            dup_other_counts = r.get("dup_other_counts", {}) or {}
            intg_other_counts = r.get("intg_other_counts", {}) or {}
            if dup_other_counts:
              cmt_toggle_seq += 1
              more_id = f'dup-other-{cmt_toggle_seq}'
              sorted_items = sorted(dup_other_counts.items(), key=lambda kv: (-kv[1], kv[0].lower()))
              hidden_items = "".join(
                f'<div class="cmt-item">{esc(txt)} <span style="color:var(--muted)">({cnt})</span></div>'
                for txt, cnt in sorted_items
              )
              dup_other_row = (
                f'<div class="tag-row"><span class="tag-label dup-label">Duplicação (outros):</span>'
                f'<button type="button" class="cmt-toggle" data-target="{more_id}" data-open="false">'
                f'+{len(sorted_items)} item(ns) adicionais</button>'
                f'<div id="{more_id}" class="cmt-hidden">{hidden_items}</div></div>'
              )
            if intg_other_counts:
              cmt_toggle_seq += 1
              more_id = f'intg-other-{cmt_toggle_seq}'
              sorted_items = sorted(intg_other_counts.items(), key=lambda kv: (-kv[1], kv[0].lower()))
              hidden_items = "".join(
                f'<div class="cmt-item">{esc(txt)} <span style="color:var(--muted)">({cnt})</span></div>'
                for txt, cnt in sorted_items
              )
              intg_other_row = (
                f'<div class="tag-row"><span class="tag-label">Integração (outros):</span>'
                f'<button type="button" class="cmt-toggle" data-target="{more_id}" data-open="false">'
                f'+{len(sorted_items)} item(ns) adicionais</button>'
                f'<div id="{more_id}" class="cmt-hidden">{hidden_items}</div></div>'
              )

            cmt_html = ""
            if r["comments"]:
                cmt_html = '<div class="comments-section"><div class="cmt-label">Sugestões dos especialistas (amostra)</div>'
                for c in r["comments"][:2]:
                    cmt_html += f'<div class="cmt-item">"{esc(c)}"</div>'
                if len(r["comments"]) > 2:
                    cmt_toggle_seq += 1
                    more_id = f'cmt-more-{cmt_toggle_seq}'
                    hidden_comments = "".join(
                        f'<div class="cmt-item">"{esc(c)}"</div>' for c in r["comments"][2:]
                    )
                    cmt_html += (
                        f'<button type="button" class="cmt-toggle" '
                        f'data-target="{more_id}" data-open="false">'
                        f'+{len(r["comments"]) - 2} comentário(s) adicionais</button>'
                        f'<div id="{more_id}" class="cmt-hidden">{hidden_comments}</div>'
                    )
                cmt_html += '</div>'

            missing_note = (f'<div class="missing-note">⚠ {r["n_missing"]} resposta(s) em falta / '
                            f'não aplicável</div>') if r["n_missing"] else ""

            label_with_idx = display_label.get(r["code"], r["label"])
            name_cell = (f'<a href="{esc(r["url"])}" target="_blank" rel="noopener noreferrer" class="inv-link">{esc(label_with_idx)}</a>'
                         if r["url"] else esc(label_with_idx))

            import re as _re
            safe_card_id = _re.sub(r'[^A-Za-z0-9_-]', '-', str(r['code']))
            _prog_label = (r.get("programa") or "").strip()
            cards_html += f"""<div class="inv-card" id="card-{safe_card_id}">
              <div class="card-top">
                {f'<span class="prog-tag">{esc(_prog_label)}</span>' if _prog_label else '<span></span>'}
                <span class="comp-tag">{esc(r["component"])}</span>
              </div>
              <h3>{name_cell}</h3>
              <div class="gate-label-sm">Triagem: precisa de optimização? (n={r["n_total"]})</div>
              <div class="gate-bar-wrap">
                <div class="gate-seg seg-sim"  style="width:{sim_p}%"></div>
                <div class="gate-seg seg-poss" style="width:{poss_p}%"></div>
                <div class="gate-seg seg-nao"  style="width:{nao_p}%"></div>
              </div>
              <div class="gate-counts">
                <div class="gc"><div class="dot" style="background:var(--sim)"></div><span class="val">{r["n_sim"]}</span><span class="gc-lbl">Sim def. ({sim_p}%)</span></div>
                <div class="gc"><div class="dot" style="background:#d4a017"></div><span class="val">{r["n_poss"]}</span><span class="gc-lbl">Possiv. ({poss_p}%)</span></div>
                <div class="gc"><div class="dot" style="background:#e0e0e0"></div><span class="val">{r["n_nao"]}</span><span class="gc-lbl">Não ({nao_p}%)</span></div>
              </div>
              <div class="metrics-row">
                <div class="metric-chip"><div class="mc-val" style="color:#7a5c00">{r.get("gate_mean",0):.2f}<span class="mc-denom">/1</span></div><div class="mc-lbl">Score triagem</div></div>
                <div class="metric-chip"><div class="mc-val {imp_cls}">{imp:.1f}<span class="mc-denom">/3</span></div><div class="mc-lbl">Impacto esperado</div></div>
                <div class="metric-chip"><div class="mc-val" style="color:var(--ink)">{r["dup_pct"]}%</div><div class="mc-lbl">Duplicação</div></div>
                <div class="metric-chip"><div class="mc-val" style="color:var(--accent2)">{r["intg_pct"]}%</div><div class="mc-lbl">Integração</div></div>
                <div class="metric-chip"><div class="mc-val" style="color:var(--muted)">{r["res_pct"]}%</div><div class="mc-lbl">↓ Recursos</div></div>
              </div>
              <div class="score-strip-row">
                <div class="score-strip-item"><span class="score-strip-lbl">S<sub>base</sub></span>{svg_score_strip(all_s_base_vals, r["s_base"], "#1a6b3a")}</div>
                <div class="score-strip-item"><span class="score-strip-lbl">S<sub>pond</sub></span>{svg_score_strip(all_s_pond_vals, r["s_pond"], "#1a5276")}</div>
              </div>
              {dup_row}{intg_row}{dup_other_row}{intg_other_row}
              {cmt_html}
              {missing_note}
            </div>\n"""
          cards_html += "</div>\n"  # close priority-grid for this component

    def _tier_links(rows):
      if not rows:
        return "—"
      chunks = []
      for r in rows:
        safe_id = _safe_anchor(r.get("code", ""))
        lbl = display_label.get(r["code"], r["label"])
        chunks.append(f'<a href="#card-{safe_id}" class="inv-link">{esc(lbl)}</a>')
      return " · ".join(chunks)

    alta_str  = _tier_links(alta_rows)
    media_str = _tier_links(media_rows)
    baixa_str = _tier_links(baixa_rows)

    # ── Univariate analysis section ────────────────────────────────────────
    univ_html = ""
    scoring_html = ""
    if univariate:
        u = univariate
        GATE_COLORS = ("#1a6b3a", "#d4a017", "#e0e0e0")
        IMP_COLORS  = ("#1a6b3a", "#fef9e7", "#e0e0e0")  # Alto / Médio / Baixo
        EXP_COLORS  = ("#c7d2fe", "#818cf8", "#4f46e5")

        gate_donut = svg_donut(u["gate_agg"], GATE_COLORS, size=110, label="Triagem")
        imp_donut  = svg_donut(
            {"Alto (3)":  u["imp_counts"].get(3,0),
             "Médio (2)": u["imp_counts"].get(2,0),
             "Baixo (1)": u["imp_counts"].get(1,0)},
            IMP_COLORS, size=110, label="Impacto"
        )
        exp_donut = svg_donut(
            {"Geral": u["exp_counts"].get(1,0),
             "Intermédio": u["exp_counts"].get(2,0),
             "Especialista": u["exp_counts"].get(3,0)},
            EXP_COLORS, size=110, label="Experiência"
        ) if any(u["exp_counts"].values()) else ""

        def _comp(r): return (r.get("component") or "").strip()
        def _hrefs(lst): return [f"#card-{__import__('re').sub(r'[^A-Za-z0-9_-]','-',str(r['code']))}" for r in lst]
        gate_bars = svg_hbar_stacked(
          [(display_label.get(r["code"], r["label"])[:52], [r["n_sim"], r["n_poss"], r["n_nao"]])
             for r in u["sorted_gate"]],
            width=600, bar_h=14, gap=5, colors=GATE_COLORS, label_w=240,
            hrefs=_hrefs(u["sorted_gate"]),
            sublabels=[_comp(r) for r in u["sorted_gate"]],
            markers=[r["gate_mean"] for r in u["sorted_gate"]],
        )
        impact_bars = svg_hbar_stacked(
          [(display_label.get(r["code"], r["label"])[:52], [r["n_imp_3"], r["n_imp_2"], r["n_imp_1"]])
             for r in u["sorted_impact"]],
            width=600, bar_h=14, gap=5, colors=IMP_COLORS, label_w=240,
            hrefs=_hrefs(u["sorted_impact"]),
            sublabels=[_comp(r) for r in u["sorted_impact"]],
            markers=[(r["avg_impact"] - 1) / 2.0 if r["avg_impact"] else None for r in u["sorted_impact"]],
        )
        dup_bars  = svg_hbar_single(
          [(display_label.get(r["code"], r["label"])[:52], r["dup_pct"],  100) for r in u["sorted_dup"]],
            width=600, bar_h=14, gap=5, color="#9b2226", label_w=240, fmt=".0f",
            hrefs=_hrefs(u["sorted_dup"]),
            sublabels=[_comp(r) for r in u["sorted_dup"]],
        )
        intg_bars = svg_hbar_single(
          [(display_label.get(r["code"], r["label"])[:52], r["intg_pct"], 100) for r in u["sorted_intg"]],
            width=600, bar_h=14, gap=5, color="#1a5276", label_w=240, fmt=".0f",
            hrefs=_hrefs(u["sorted_intg"]),
            sublabels=[_comp(r) for r in u["sorted_intg"]],
        )
        res_bars  = svg_hbar_single(
          [(display_label.get(r["code"], r["label"])[:52], r["res_pct"],  100) for r in u["sorted_res"]],
            width=600, bar_h=14, gap=5, color="#5b5ea6", label_w=240, fmt=".0f",
            hrefs=_hrefs(u["sorted_res"]),
            sublabels=[_comp(r) for r in u["sorted_res"]],
        )
        e_bars = svg_hbar_single(
          [(display_label.get(r["code"], r["label"])[:52], (r.get("e_score", 0) * 100), 100) for r in u["sorted_e"]],
            width=600, bar_h=14, gap=5, color="#1e7e34", label_w=240, fmt=".0f",
            hrefs=_hrefs(u["sorted_e"]),
            sublabels=[_comp(r) for r in u["sorted_e"]],
        )
        scatter_oi = svg_scatter_optim_impact_exp(list(results.values())) if include_xyplot else ""
        xy_html = (f'<div class="univ-sub-header" style="margin-top:28px">2h · Dispersão em escala real (janela observada) com zonas da regra actual</div>'
             f'<div class="chart-wrap">{scatter_oi}</div>') if include_xyplot else ""

        univ_html = f"""
  <div class="section-header">
    <h2>Análise Univariada</h2>
    <span class="subtitle">Distribuições globais e por intervenção</span>
  </div>
  <div class="univ-panel">
    <div class="univ-row two-col">
      <div class="univ-col">
        <div class="univ-sub-header">2a · Distribuição da pergunta de triagem</div>
        <div style="font-size:10px;color:var(--muted);margin-bottom:4px">
          <span style="display:inline-block;width:10px;height:9px;background:#1a6b3a;vertical-align:middle;margin-right:2px"></span>Sim def. <strong>(1,0)</strong>
          &nbsp;<span style="display:inline-block;width:10px;height:9px;background:#d4a017;vertical-align:middle;margin-right:2px"></span>Possiv. <strong>(0,5)</strong>
          &nbsp;<span style="display:inline-block;width:10px;height:9px;background:#e0e0e0;vertical-align:middle;margin-right:2px"></span>Não <strong>(0,0)</strong>
          &nbsp;&nbsp;<svg width="10" height="10" style="vertical-align:middle;margin-right:2px"><polygon points="5,1 9,5 5,9 1,5" fill="white" stroke="#1e3a5f" stroke-width="1.5"/></svg>Score triagem médio (◆)
        </div>
        <div class="donut-row">{gate_donut}</div>
        <div class="chart-wrap">{gate_bars}</div>
      </div>
      <div class="univ-col">
        <div class="univ-sub-header">2b · Impacto esperado (escala 1–3)</div>
        <div style="font-size:10px;color:var(--muted);margin-bottom:4px">
          <span style="display:inline-block;width:10px;height:9px;background:#1a6b3a;vertical-align:middle;margin-right:2px"></span>Alto
          &nbsp;<span style="display:inline-block;width:10px;height:9px;background:#fef9e7;border:1px solid #d1d5db;vertical-align:middle;margin-right:2px"></span>Médio
          &nbsp;<span style="display:inline-block;width:10px;height:9px;background:#e0e0e0;vertical-align:middle;margin-right:2px"></span>Baixo
          &nbsp;&nbsp;<svg width="10" height="10" style="vertical-align:middle;margin-right:2px"><polygon points="5,1 9,5 5,9 1,5" fill="white" stroke="#1e3a5f" stroke-width="1.5"/></svg>Impacto médio (◆)
        </div>
        <div class="donut-row">{imp_donut}</div>
        <div class="chart-wrap">{impact_bars}</div>
      </div>
    </div>

    {'<div class="univ-sub-header" style="margin-top:28px">2c · Experiência declarada dos respondentes</div><div style="font-size:10px;color:var(--muted);margin-bottom:4px"><span style="display:inline-block;width:10px;height:9px;background:#c7d2fe;vertical-align:middle;margin-right:2px"></span>Geral <strong>(×1)</strong> &nbsp;<span style="display:inline-block;width:10px;height:9px;background:#818cf8;vertical-align:middle;margin-right:2px"></span>Intermédio <strong>(×2)</strong> &nbsp;<span style="display:inline-block;width:10px;height:9px;background:#4f46e5;vertical-align:middle;margin-right:2px"></span>Especialista <strong>(×3)</strong></div><div class="donut-row">' + exp_donut + '</div>' if exp_donut else ''}

    <div class="univ-row two-col" style="margin-top:28px">
      <div class="univ-col">
        <div class="univ-sub-header">2d · % que identificou duplicação (por intervenção)</div>
        <div class="chart-wrap">{dup_bars}</div>
      </div>
      <div class="univ-col">
        <div class="univ-sub-header">2e · % que identificou potencial de integração</div>
        <div class="chart-wrap">{intg_bars}</div>
      </div>
    </div>

    <div class="univ-sub-header" style="margin-top:28px">2f · % que identificou possibilidade de redução de recursos</div>
    <div class="chart-wrap">{res_bars}</div>

    <div class="univ-sub-header" style="margin-top:28px">2g · Índice de Eficiência (E)</div>
    <div style="display:flex;gap:18px;align-items:center;flex-wrap:wrap;margin-bottom:8px;font-size:11px;color:var(--muted)">
      <span><strong>Mediana:</strong> {u['e_median']*100:.1f}%</span>
      <span><strong>Mínimo:</strong> {u['e_min']*100:.1f}%</span>
      <span><strong>Máximo:</strong> {u['e_max']*100:.1f}%</span>
    </div>
    <div class="chart-wrap">{e_bars}</div>

    {xy_html}
  </div>
"""

        # ── Scoring section ──────────────────────────────────────────────────
        sorted_base = sorted(results.values(), key=lambda r: r["s_base"], reverse=True)
        sorted_pond = sorted(results.values(), key=lambda r: r["s_pond"], reverse=True)
        max_base = max((r["s_base"] for r in results.values()), default=1) or 1
        max_pond = max((r["s_pond"] for r in results.values()), default=1) or 1

        _card_href = lambda r: f"#card-{__import__('re').sub(r'[^A-Za-z0-9_-]','-',str(r['code']))}"
        _num_suffix = lambda code: (__import__('re').search(r'(\d+)$', str(code)) or type('', (), {'group': lambda s, n: code})()).group(1)
        score_base_bars = svg_hbar_single(
            [(f'{_num_suffix(r["code"])} · {r["label"][:34]}', r["s_base"], max_base) for r in sorted_base],
            width=520, bar_h=14, gap=5, color="#1a6b3a", label_w=280, fmt=".3f",
            hrefs=[_card_href(r) for r in sorted_base]
        )
        score_pond_bars = svg_hbar_single(
            [(f'{_num_suffix(r["code"])} · {r["label"][:34]}', r["s_pond"], max_pond) for r in sorted_pond],
            width=520, bar_h=14, gap=5, color="#1a5276", label_w=280, fmt=".3f",
            hrefs=[_card_href(r) for r in sorted_pond]
        )

        _all_by_base = sorted(results.values(), key=lambda r: r.get("rank_base", 10**9))
        alluvial_items = _all_by_base if alluvial_top is None else _all_by_base[:alluvial_top]
        alluvial_svg = svg_alluvial_weighting(alluvial_items)
        if alluvial_top is not None and len(results) > alluvial_top:
            _excl = ", ".join(esc(display_label.get(r["code"], r["label"])) for r in _all_by_base[alluvial_top:])
            alluvial_note = (f'<div style="font-size:11px;color:var(--muted);margin-top:4px">'
                            f'<em>Diagrama mostra as {alluvial_top} intervenções com maior S<sub>base</sub>'
                            f' (de {len(results)} no total). Exclu&#237;das: {_excl}.</em></div>')
        else:
            alluvial_note = ""

        score_table_rows = ""
        if simple_table:
            sorted_simple = sorted(results.values(), key=lambda r: r.get("s_pond", 0), reverse=True)
            for r in sorted_simple:
                grp = tier(r)
                grp_num = 1 if grp == "alta" else (2 if grp == "media" else 3)
                e_pct = r.get("e_score", 0) * 100
                score_table_rows += f"""<tr>
              <td style="font-weight:500;font-size:12px"><a href="#card-{__import__('re').sub(r'[^A-Za-z0-9_-]', '-', str(r['code']))}" class="inv-link">{esc(display_label.get(r["code"], r["label"]))}</a></td>
              <td style="text-align:center;font-weight:600;color:#1a5276">{r.get("s_pond", 0):.3f}</td>
              <td style="text-align:center;color:#1e7e34;font-weight:600">{e_pct:.1f}%</td>
              <td style="text-align:center;font-weight:700">{grp_num}</td>
            </tr>"""
        else:
            for r in sorted_base:
                delta = r.get("rank_delta", 0)
                delta_str = f"+{delta}" if delta > 0 else str(delta)
                delta_col = "#1a6b3a" if delta > 0 else ("#c0392b" if delta < 0 else "#6b7280")
                score_table_rows += f"""<tr>
              <td class="rank-num">{r.get("display_idx","")}</td>
              <td style="font-weight:500;font-size:12px"><a href="#card-{__import__('re').sub(r'[^A-Za-z0-9_-]', '-', str(r['code']))}" class="inv-link">{esc(display_label.get(r["code"], r["label"]))}</a></td>
              <td style="text-align:center;font-weight:600;color:#7a5c00">{r.get("gate_mean", 0):.3f}</td>
              <td style="text-align:center;font-weight:600">{r["s_base"]:.3f}</td>
              <td style="text-align:center;font-weight:600;color:#1a5276">{r["s_pond"]:.3f}</td>
              <td style="text-align:center;font-weight:600;color:{delta_col}">{delta_str}</td>
              <td style="text-align:center;color:var(--muted)">{r["avg_impact"]:.2f}</td>
              <td style="text-align:center;color:var(--muted)">{r["exp_mean"]:.2f}</td>
            </tr>"""

        if include_scoring:
            scoring_html = f"""
  <div class="section-header">
    <h2>Pontuação e Metodologia</h2>
    <span class="subtitle">Duas métricas complementares para priorização</span>
  </div>
  <div class="formula-box">
    <div class="formula-title">Fórmulas de Pontuação</div>
    <div class="formula-grid">
      <div>
        <div class="formula-label">S<sub>base</sub> — Pontuação não ponderada (intervalo: 0–3)</div>
        <div class="formula-code">S<sub>base</sub> = <span class="f-mean">mean</span>(triagem_score<sub>i</sub>) × <span class="f-mean">mean</span>(impact<sub>i</sub>)</div>
        <div class="formula-note">triagem_score: sim_def = 1,0 · possivelmente = 0,5 · não = 0,0</div>
      </div>
      <div>
        <div class="formula-label">S<sub>pond</sub> — Pontuação ponderada por experiência (intervalo: 0–3)</div>
        <div class="formula-code">S<sub>pond</sub> = Σ(triagem<sub>i</sub> × impact<sub>i</sub> × exp<sub>i</sub>) / Σ(exp<sub>i</sub>)</div>
        <div class="formula-note">exp<sub>i</sub> ∈ {{1, 2, 3}} — nível de experiência declarada (omisso = 1)</div>
      </div>
    </div>
  </div>
  <div class="score-charts-row">
    <div class="score-chart-col">
      <div class="score-chart-title" style="color:#1a6b3a">S<sub>base</sub> — não ponderada</div>
      {score_base_bars}
    </div>
    <div class="score-chart-col">
      <div class="score-chart-title" style="color:#1a5276">S<sub>pond</sub> — ponderada por experiência</div>
      {score_pond_bars}
    </div>
  </div>
  <div class="score-chart-col" style="margin-bottom:16px">
    <div class="score-chart-title" style="color:#0f1923">Diagrama Aluvial — S_base -> S_pond</div>
    <div style="font-size:11px;color:var(--muted);margin-bottom:8px">
      1ª coluna: rank por S<sub>base</sub> (não ponderada). 2ª coluna: rank por S<sub>pond</sub>
      (ponderada por experiência declarada).
      Verde/vermelho/azul = subida/descida/sem mudança entre S<sub>base</sub> e S<sub>pond</sub>.
    </div>
    <div class="chart-wrap">{alluvial_svg}</div>
    {alluvial_note}
  </div>
  <table id="score-table" class="rank-table" style="margin-top:16px" {'data-default-sort-col="1" data-default-sort-dir="desc"' if simple_table else ''}>
    <thead>
      <tr>
        {('<th data-sort="text" style="cursor:pointer">Intervenção <span class="sa"></span></th>'
          '<th data-sort="num" style="text-align:center;cursor:pointer">S<sub>pond</sub> <span class="sa"></span></th>'
          '<th data-sort="num" style="text-align:center;cursor:pointer">Índice de eficiência <span class="sa"></span></th>'
          '<th data-sort="num" style="text-align:center;cursor:pointer">Grupo Próximos Passos (1/2/3) <span class="sa"></span></th>') if simple_table else
         ('<th data-sort="num" style="cursor:pointer"># <span class="sa"></span></th>'
          '<th data-sort="text" style="cursor:pointer">Intervenção <span class="sa"></span></th>'
          '<th data-sort="num" style="text-align:center;cursor:pointer">Triagem (média) <span class="sa"></span></th>'
          '<th data-sort="num" style="text-align:center;cursor:pointer">S<sub>base</sub> <span class="sa"></span></th>'
          '<th data-sort="num" style="text-align:center;cursor:pointer">S<sub>pond</sub> <span class="sa"></span></th>'
          '<th data-sort="num" style="text-align:center;cursor:pointer">Δ rank <span class="sa"></span></th>'
          '<th data-sort="num" style="text-align:center;cursor:pointer">Impacto médio <span class="sa"></span></th>'
          '<th data-sort="num" style="text-align:center;cursor:pointer">Exp. média <span class="sa"></span></th>')}
      </tr>
    </thead>
    <tbody>{score_table_rows}</tbody>
  </table>
"""
    else:
      scoring_html = ""

    network_html = build_network_html(results, topic_label=topic_label)
    integration_candidates_html = build_integration_candidates_html(results)

    html = f"""<!DOCTYPE html>
<html lang="pt">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width,initial-scale=1.0">
<title>Delphi W1 — Resultados Agregados</title>
<link href="https://fonts.googleapis.com/css2?family=DM+Serif+Display:ital@0;1&family=DM+Sans:wght@300;400;500;600&display=swap" rel="stylesheet">
<style>
:root {{
  --ink:#0f1923; --paper:#f5f2ed; --accent:#c0392b; --accent2:#1a5276;
  --gold:#b7860b; --muted:#6b7280; --border:#d5cfc6;
  --sim:#1a6b3a; --sim-bg:#e8f5ee; --poss:#7a5c00; --poss-bg:#fef9e7;
  --card-bg:#ffffff;
}}
*{{box-sizing:border-box;margin:0;padding:0}}
body{{font-family:'DM Sans',sans-serif;background:var(--paper);color:var(--ink);line-height:1.5;font-size:14px}}
.site-header{{background:var(--ink);color:white;padding:0 40px;display:flex;align-items:stretch;border-bottom:4px solid var(--accent)}}
.header-brand{{padding:24px 0;border-right:1px solid rgba(255,255,255,.12);padding-right:32px;margin-right:32px}}
.header-brand .eyebrow{{font-size:10px;letter-spacing:.2em;text-transform:uppercase;color:rgba(255,255,255,.5);margin-bottom:4px}}
.header-brand h1{{font-family:'DM Serif Display',serif;font-size:22px;font-weight:400;color:white;line-height:1.2}}
.header-meta{{padding:24px 0;display:flex;flex-direction:column;justify-content:center;gap:4px;flex:1}}
.header-logo{{padding:16px 0 16px 32px;display:flex;align-items:center;margin-left:auto}}
.header-logo img{{max-height:64px;max-width:160px;object-fit:contain}}
.wave-badge{{display:inline-flex;align-items:center;gap:8px;background:var(--accent);color:white;font-size:11px;font-weight:600;letter-spacing:.08em;padding:3px 10px;border-radius:2px;width:fit-content;margin-bottom:6px}}
.header-meta p{{font-size:12px;color:rgba(255,255,255,.6)}}
.header-meta strong{{color:rgba(255,255,255,.9)}}
.intro-band{{background:var(--accent2);color:white;padding:20px 40px;display:flex;align-items:center;gap:20px}}
.intro-band .icon{{font-size:28px;flex-shrink:0}}
.intro-band p{{font-size:13px;line-height:1.6;color:rgba(255,255,255,.88)}}
.intro-band strong{{color:white}}
.summary-bar{{background:white;border-bottom:1px solid var(--border);padding:20px 40px;display:flex;gap:32px;align-items:center;flex-wrap:wrap}}
.stat-pill{{display:flex;flex-direction:column;align-items:center;gap:2px}}
.stat-pill .num{{font-family:'DM Serif Display',serif;font-size:32px;line-height:1;color:var(--accent)}}
.stat-pill .lbl{{font-size:10px;text-transform:uppercase;letter-spacing:.12em;color:var(--muted);text-align:center;max-width:80px}}
.stat-divider{{width:1px;height:40px;background:var(--border)}}
.legend-box{{margin-left:auto;display:flex;gap:16px;align-items:center;flex-wrap:wrap}}
.legend-item{{display:flex;align-items:center;gap:6px;font-size:11px;color:var(--muted)}}
.legend-dot{{width:10px;height:10px;border-radius:50%}}
.dot-sim{{background:var(--sim)}}.dot-poss{{background:#d4a017}}.dot-nao{{background:#d1d5db}}
.main{{padding:32px 40px;max-width:1300px;margin:0 auto}}
.section-header{{display:flex;align-items:baseline;gap:12px;margin-bottom:20px;padding-bottom:10px;border-bottom:2px solid var(--ink)}}
.section-header h2{{font-family:'DM Serif Display',serif;font-size:20px;font-weight:400}}
.section-header .subtitle{{font-size:12px;color:var(--muted);letter-spacing:.04em}}
.note-box{{background:#fffbeb;border:1px solid #fcd34d;border-left:4px solid var(--gold);border-radius:3px;padding:14px 18px;margin-bottom:28px;font-size:12px;color:#78350f;line-height:1.6}}
.note-box strong{{color:#92400e}}
/* ── Response rate section ── */
.rr-panel{{background:white;border:1px solid var(--border);border-radius:4px;margin-bottom:40px;overflow:hidden}}
.rr-summary{{display:flex;gap:0;border-bottom:1px solid var(--border)}}
.rr-stat{{flex:1;padding:16px 20px;text-align:center;border-right:1px solid var(--border)}}
.rr-stat:last-child{{border-right:none}}
.rr-stat .rr-num{{font-family:'DM Serif Display',serif;font-size:28px;line-height:1;color:var(--accent2)}}
.rr-stat .rr-lbl{{font-size:10px;text-transform:uppercase;letter-spacing:.1em;color:var(--muted);margin-top:2px}}
.rr-table{{width:100%;border-collapse:collapse;font-size:12px}}
.rr-table td{{padding:7px 16px;border-bottom:1px solid #f0ede8;vertical-align:middle}}
.rr-table tr:last-child td{{border-bottom:none}}
.rr-table tr:nth-child(even) td{{background:#fafaf8}}
.rr-details{{margin-top:12px}}
.rr-details summary{{cursor:pointer;padding:10px 16px;background:#f7f9fb;border-top:1px solid var(--border);font-size:12px;font-weight:600;color:var(--accent2);user-select:none;display:flex;align-items:center;gap:8px}}
.rr-details summary:hover{{background:#eef2f7}}
.rr-details summary::marker{{content:''}}
.rr-details summary::before{{content:'▶';font-size:10px;transition:transform 0.2s}}
.rr-details[open] summary::before{{transform:rotate(90deg)}}
/* ── Ranking table ── */
.rank-table{{width:100%;border-collapse:collapse;margin-bottom:40px;background:white;border:1px solid var(--border);border-radius:4px;overflow:hidden;font-size:12px}}
.rank-table thead tr{{background:var(--ink);color:white}}
.rank-table th{{padding:10px 14px;font-size:10px;text-transform:uppercase;letter-spacing:.1em;font-weight:600;text-align:left}}
.rank-table td{{padding:9px 14px;border-bottom:1px solid var(--border);vertical-align:middle}}
.rank-table tr:last-child td{{border-bottom:none}}
.rank-table tr:nth-child(even) td{{background:#fafafa}}
.rank-num{{font-family:'DM Serif Display',serif;font-size:18px;color:var(--muted);width:32px}}
.comp-tag{{font-size:10px;background:var(--paper);border:1px solid var(--border);color:var(--muted);padding:2px 6px;border-radius:2px;white-space:nowrap}}
.prog-tag{{font-size:10px;background:var(--accent2);color:#fff;padding:2px 6px;border-radius:2px;white-space:nowrap;font-weight:500}}
.inv-link{{color:inherit;text-decoration:none;border-bottom:1px dotted var(--muted)}}
.inv-link:hover{{color:var(--accent2);border-bottom-color:var(--accent2)}}
.mini-bar-wrap{{display:flex;align-items:center;gap:8px}}
.mini-bar-bg{{flex:1;height:8px;background:#e5e7eb;border-radius:2px;overflow:hidden}}
.bar-counts{{font-size:10px;color:var(--muted);white-space:nowrap}}
.impact-badge{{display:inline-flex;align-items:center;gap:4px;font-size:11px;font-weight:600;padding:3px 8px;border-radius:2px}}
.imp-high{{background:var(--sim-bg);color:var(--sim)}}
.imp-med{{background:var(--poss-bg);color:var(--poss)}}
.imp-low{{background:#f3f4f6;color:var(--muted)}}
/* ── Cards ── */
.comp-label{{font-size:11px;font-weight:700;letter-spacing:.12em;text-transform:uppercase;color:var(--accent2);padding:6px 0;border-bottom:1px solid var(--border);margin-bottom:14px;margin-top:28px}}
.priority-grid{{display:grid;grid-template-columns:1fr 1fr;gap:16px;margin-bottom:16px}}
.inv-card{{background:var(--card-bg);border:1px solid var(--border);border-radius:4px;padding:16px 18px}}
.inv-card:hover{{box-shadow:0 4px 16px rgba(0,0,0,.08)}}
.card-top{{display:flex;align-items:flex-start;justify-content:space-between;gap:8px;margin-bottom:10px}}
.inv-card h3{{font-size:13px;font-weight:600;line-height:1.3;margin-bottom:12px}}
.gate-label-sm{{font-size:10px;text-transform:uppercase;letter-spacing:.1em;color:var(--muted);margin-bottom:5px}}
.gate-bar-wrap{{display:flex;height:12px;border-radius:2px;overflow:hidden;gap:1px;margin-bottom:4px}}
.gate-seg{{height:100%}}
.seg-sim{{background:var(--sim)}}.seg-poss{{background:#d4a017}}.seg-nao{{background:#e0e0e0}}
.gate-counts{{display:flex;gap:12px;margin-bottom:2px}}
.gc{{display:flex;align-items:center;gap:4px;font-size:10px}}
.gc .dot{{width:7px;height:7px;border-radius:50%}}
.gc .val{{font-weight:600}}
.gc-lbl{{color:var(--muted)}}
.metrics-row{{display:flex;gap:8px;margin-top:12px;padding-top:12px;border-top:1px solid var(--border)}}
.metric-chip{{flex:1;background:var(--paper);border-radius:3px;padding:6px 8px;text-align:center}}
.mc-val{{font-family:'DM Serif Display',serif;font-size:16px;line-height:1;margin-bottom:2px}}
.mc-denom{{font-size:10px;font-family:'DM Sans',sans-serif}}
.mc-lbl{{font-size:9px;text-transform:uppercase;letter-spacing:.08em;color:var(--muted)}}
.mc-impact-high{{color:var(--sim)}}.mc-impact-med{{color:var(--poss)}}.mc-impact-low{{color:#9ca3af}}
.score-strip-row{{margin-top:10px;padding-top:8px;border-top:1px solid var(--border);display:flex;flex-direction:column;gap:4px}}
.score-strip-item{{display:flex;align-items:center;gap:6px}}
.score-strip-lbl{{font-size:9px;text-transform:uppercase;letter-spacing:.06em;color:var(--muted);width:38px;flex-shrink:0;text-align:right}}
.tag-row{{margin-top:8px;display:flex;flex-wrap:wrap;gap:4px;align-items:center}}
.tag-label{{font-size:9px;text-transform:uppercase;letter-spacing:.08em;color:var(--muted);margin-right:2px}}
.dup-label{{color:#9b2226}}
.intg-tag{{font-size:10px;font-weight:600;background:#e8f0fe;color:#1a3a8f;padding:2px 7px;border-radius:2px;white-space:normal;line-height:1.4}}
.dup-tag{{font-size:10px;font-weight:600;background:#fde8e8;color:#9b2226;padding:2px 7px;border-radius:2px;white-space:normal;line-height:1.4}}
.comments-section{{margin-top:10px;padding-top:10px;border-top:1px solid var(--border)}}
.cmt-label{{font-size:9px;text-transform:uppercase;letter-spacing:.08em;color:var(--muted);margin-bottom:5px}}
.cmt-item{{font-size:11px;color:var(--ink);font-style:italic;line-height:1.5;margin-bottom:4px;padding-left:8px;border-left:2px solid var(--border)}}
.cmt-toggle{{margin-top:3px;padding:0;border:0;background:none;font-size:10px;color:var(--muted);text-decoration:underline;cursor:pointer}}
.cmt-toggle:hover{{color:var(--ink)}}
.cmt-hidden{{display:none;margin-top:4px}}
.cmt-hidden.is-open{{display:block}}
.missing-note{{font-size:10px;color:#b45309;margin-top:6px;background:#fffbeb;padding:3px 6px;border-radius:2px}}
/* ── Next steps ── */
.next-box{{background:white;border:1px solid var(--border);border-radius:4px;padding:24px;margin-bottom:40px}}
.next-box p{{font-size:13px;line-height:1.8}}
.tier-panels{{margin-top:16px;padding-top:16px;border-top:1px solid var(--border);display:flex;gap:16px;flex-wrap:wrap}}
.tier-panel{{flex:1;min-width:180px;border-radius:3px;padding:12px 16px}}
.tier-alta{{background:var(--sim-bg)}}.tier-media{{background:var(--poss-bg)}}.tier-baixa{{background:#f3f4f6}}
.tier-heading{{font-size:10px;text-transform:uppercase;letter-spacing:.1em;font-weight:700;margin-bottom:6px}}
.tier-alta .tier-heading{{color:var(--sim)}}.tier-media .tier-heading{{color:var(--poss)}}.tier-baixa .tier-heading{{color:var(--muted)}}
.tier-panel p{{font-size:12px;line-height:1.8}}
.report-footer{{background:var(--ink);color:rgba(255,255,255,.5);padding:20px 40px;font-size:11px;display:flex;justify-content:space-between;align-items:center;margin-top:40px;flex-wrap:wrap;gap:8px}}
.report-footer strong{{color:rgba(255,255,255,.8)}}
@media print{{body{{background:white}}.site-header,.intro-band{{-webkit-print-color-adjust:exact;print-color-adjust:exact}}.inv-card{{break-inside:avoid}}}}
/* ── Univariate ── */
.univ-panel{{background:white;border:1px solid var(--border);border-radius:4px;padding:24px;margin-bottom:40px}}
.univ-sub-header{{font-size:11px;font-weight:700;text-transform:uppercase;letter-spacing:.1em;color:var(--accent2);margin-bottom:10px}}
.univ-row.two-col{{display:grid;grid-template-columns:1fr 1fr;gap:24px}}
.univ-col{{min-width:0}}
.donut-row{{margin-bottom:12px}}
.chart-wrap{{overflow-x:auto;padding-bottom:4px}}
/* ── Scoring ── */
.formula-box{{background:#0f1923;color:white;border-radius:4px;padding:20px 24px;margin-bottom:20px}}
.formula-title{{font-size:11px;text-transform:uppercase;letter-spacing:.12em;color:rgba(255,255,255,.5);margin-bottom:12px}}
.formula-grid{{display:grid;grid-template-columns:1fr 1fr;gap:24px}}
.formula-label{{font-size:11px;color:rgba(255,255,255,.6);margin-bottom:6px}}
.formula-code{{font-family:monospace;font-size:14px;color:#86efac;margin-bottom:4px}}
.formula-note{{font-size:10px;color:rgba(255,255,255,.4)}}
.score-charts-row{{display:grid;grid-template-columns:1fr 1fr;gap:24px;margin-bottom:16px}}
.score-chart-col{{background:white;border:1px solid var(--border);border-radius:4px;padding:16px;overflow-x:auto}}
.score-chart-title{{font-size:12px;font-weight:600;margin-bottom:10px}}
@media (max-width: 980px){{
  .univ-row.two-col{{grid-template-columns:1fr}}
}}
</style>
</head>
<body>

<header class="site-header">
  <div class="header-brand">
    <div class="eyebrow">{esc(programa_nome)}</div>
    <h1>Exercício de Optimização<br>de Intervenções</h1>
  </div>
  <div class="header-meta">
    <div class="wave-badge">⬤ &nbsp;ONDA 1 — RESULTADOS AGREGADOS</div>
    <p>Gerado em: <strong>{now}</strong> &nbsp;|&nbsp; Fonte: <strong>{esc(os.path.basename(source_file))}</strong></p>
    <p>Base esperada: <strong>{stats["n_experts"]} especialistas</strong> (experts.txt) &nbsp;|&nbsp; Respondentes observados: <strong>{n_experts_observed}</strong> &nbsp;|&nbsp; Intervenções avaliadas: <strong>{stats["n_inv"]}</strong></p>
    <p>Este documento é <strong>confidencial</strong> — para uso exclusivo dos participantes da oficina Delphi</p>
  </div>
  {f'<div class="header-logo"><img src="{logo_data_uri}" alt="Logótipo"></div>' if logo_data_uri else ''}
</header>

<div class="intro-band">
  <div class="icon">📋</div>
  <p>Este relatório apresenta os <strong>resultados anónimos e agregados</strong> da primeira onda (W1).
  Os dados reflectem as avaliações colectivas dos especialistas sobre o potencial de optimização de cada intervenção.
  <strong>Nenhuma resposta individual é identificável.</strong>
  Com base nestes resultados, o grupo seleccionará as intervenções prioritárias para a Onda 2 (W2).</p>
</div>

<div class="summary-bar">
  <div class="stat-pill"><span class="num">{stats["n_experts"]}</span><span class="lbl">Especialistas esperados</span></div>
  <div class="stat-divider"></div>
  <div class="stat-pill"><span class="num">{n_experts_observed}</span><span class="lbl">Especialistas respondentes</span></div>
  <div class="stat-divider"></div>
  <div class="stat-pill"><span class="num">{stats["n_inv"]}</span><span class="lbl">Intervenções avaliadas</span></div>
  <div class="stat-divider"></div>
  <div class="stat-pill"><span class="num">{stats["n_inv_80"]}</span><span class="lbl">Com ≥80% a favor de optimização</span></div>
  <div class="stat-divider"></div>
  <div class="stat-pill"><span class="num">{stats["n_imp_high"]}</span><span class="lbl">Com impacto esperado alto (≥2,5)</span></div>
  <div class="stat-divider"></div>
  <div class="stat-pill"><span class="num">{stats["n_unanimous"]}</span><span class="lbl">Com consenso unânime (100%)</span></div>
  <div class="legend-box">
    <div class="legend-item"><div class="legend-dot dot-sim"></div>Sim, definitivamente</div>
    <div class="legend-item"><div class="legend-dot dot-poss"></div>Possivelmente</div>
    <div class="legend-item"><div class="legend-dot dot-nao"></div>Não</div>
  </div>
</div>

<div class="main">

  <div class="note-box">
    <strong>Como ler este relatório:</strong> Para cada intervenção são apresentados:
    (1) a distribuição de respostas à pergunta de triagem ("precisa de optimização?"),
    (2) o impacto médio esperado se optimizada (escala 1–3),
    e (3) as percentagens de especialistas que identificaram duplicação, potencial de integração
    e possibilidade de redução de recursos.
    As intervenções estão ordenadas por <em>pontuação composta</em> (% "Sim definitivamente" × impacto médio).
    S = Sim definitivamente &nbsp;|&nbsp; P = Possivelmente &nbsp;|&nbsp; N = Não.
    Os nomes das intervenções são hiperligações para as respectivas fichas de descrição.
  </div>

  <!-- ═══ RESPONSE RATE ═══════════════════════════════════════════════════ -->
  <div class="section-header">
    <h2>Taxa de Resposta por Intervenção</h2>
    <span class="subtitle">Proporção de especialistas que avaliaram cada intervenção</span>
  </div>

  <div class="rr-panel">
    <div class="rr-summary">
      <div class="rr-stat">
        <div class="rr-num">{stats["rr_median"]}%</div>
        <div class="rr-lbl">Mediana</div>
      </div>
      <div class="rr-stat">
        <div class="rr-num">{stats["rr_min"]}%</div>
        <div class="rr-lbl">Mínimo</div>
      </div>
      <div class="rr-stat">
        <div class="rr-num">{stats["rr_max"]}%</div>
        <div class="rr-lbl">Máximo</div>
      </div>
      <div class="rr-stat" style="flex:2;text-align:left;padding-left:24px">
        <div style="font-size:12px;color:var(--muted);line-height:1.6">
          A taxa de resposta por intervenção varia porque o formulário W1 está dividido
          em grupos de intervenções, e nem todos os especialistas completaram todos os grupos.
          Intervenções com taxa abaixo de 60% devem ser interpretadas com cautela.
        </div>
      </div>
    </div>
    <details class="rr-details">
      <summary>Ver detalhes por intervenção ({stats["n_inv"]} intervenções)</summary>
      <table class="rr-table">
        <tbody>{rr_rows}</tbody>
      </table>
    </details>
    {rr_by_team_html}
  </div>

  {summaries_html}

  {network_html}

  {integration_candidates_html}

  {univ_html}
  {scoring_html}

  <!-- ═══ RANKING TABLE ════════════════════════════════════════════════════ -->
  <div class="section-header">
    <h2>Tabela de Prioridade — Todas as Intervenções</h2>
    <span class="subtitle">Ordenadas por {_sort_label}</span>
  </div>

  <table class="rank-table">
    <thead>
      <tr>
        <th>#</th><th>Intervenção</th><th>Componente</th>
        <th style="width:200px">Distribuição de respostas</th>
        <th>% Optimizável</th><th>Impacto esperado</th>
        <th>% Duplicação</th><th>% Integração</th>
      </tr>
    </thead>
    <tbody>{rank_rows}</tbody>
  </table>

  <!-- ═══ DETAIL CARDS ════════════════════════════════════════════════════ -->
  <div class="section-header">
    <h2>Detalhe por Intervenção</h2>
    <span class="subtitle">Agrupado por componente programático</span>
  </div>

  {cards_html}

  <!-- ═══ NEXT STEPS ══════════════════════════════════════════════════════ -->
  <div class="section-header">
    <h2>Próximos Passos — Onda 2</h2>
  </div>
  <div class="next-box">
    <p>Com base nos resultados da W1, o grupo irá <strong>seleccionar colectivamente as intervenções a focar na W2</strong>.
    Recomenda-se priorizar intervenções estritamente acima da mediana simultaneamente em optimizabilidade e impacto esperado
    (cutoffs dinâmicos nesta ronda: triagem_média &gt; {med_gate:.3f}; impacto &gt; {med_imp:.3f}).
    As intervenções seleccionadas serão distribuídas por grupos de trabalho temáticos que irão desenvolver <strong>propostas concretas de optimização</strong>.
    Essas propostas serão submetidas à avaliação anónima colectiva na Onda 3 (W3).</p>
    <div class="tier-panels">
      <div class="tier-panel tier-alta">
        <div class="tier-heading">Candidatas de Alta Prioridade (N={len(alta_rows)})</div>
        <p>{alta_str}</p>
      </div>
      <div class="tier-panel tier-media">
        <div class="tier-heading">Candidatas de Prioridade Média (N={len(media_rows)})</div>
        <p>{media_str}</p>
      </div>
      <div class="tier-panel tier-baixa">
        <div class="tier-heading">Baixa Prioridade / Manter (N={len(baixa_rows)})</div>
        <p>{baixa_str}</p>
      </div>
    </div>
  </div>

</div>

<script>
// Sortable table
(function() {{
  var tbl = document.getElementById('score-table');
  if (!tbl) return;
  var sortCol = -1, sortAsc = true;
  var ths = tbl.querySelectorAll('thead th');
  function applySort(th, idx) {{
    var tbody = tbl.querySelector('tbody');
    var rows = Array.from(tbody.rows);
    rows.sort(function(a, b) {{
      var ta = a.cells[idx].textContent.trim();
      var tb = b.cells[idx].textContent.trim();
      var va, vb;
      if (th.getAttribute('data-sort') === 'text') {{
        va = ta.toLowerCase(); vb = tb.toLowerCase();
      }} else {{
        va = parseFloat(ta.replace(/^\\+/, '')) || 0;
        vb = parseFloat(tb.replace(/^\\+/, '')) || 0;
      }}
      if (va < vb) return sortAsc ? -1 : 1;
      if (va > vb) return sortAsc ? 1 : -1;
      return 0;
    }});
    rows.forEach(function(r) {{ tbody.appendChild(r); }});
    ths.forEach(function(h) {{ var s = h.querySelector('.sa'); if (s) s.textContent = ''; }});
    var sa = th.querySelector('.sa');
    if (sa) sa.textContent = sortAsc ? ' ▲' : ' ▼';
  }}
  ths.forEach(function(th, idx) {{
    th.addEventListener('click', function() {{
      if (sortCol === idx) {{ sortAsc = !sortAsc; }}
      else {{ sortCol = idx; sortAsc = (th.getAttribute('data-sort') !== 'text'); }}
      applySort(th, idx);
    }});
  }});

  var defaultColAttr = tbl.getAttribute('data-default-sort-col');
  if (defaultColAttr !== null) {{
    var dIdx = parseInt(defaultColAttr, 10);
    if (!isNaN(dIdx) && ths[dIdx]) {{
      sortCol = dIdx;
      sortAsc = (tbl.getAttribute('data-default-sort-dir') || 'asc').toLowerCase() !== 'desc';
      applySort(ths[dIdx], dIdx);
    }}
  }}
}})();

document.addEventListener('click', function(evt) {{
  const btn = evt.target.closest('.cmt-toggle');
  if (!btn) return;
  const targetId = btn.getAttribute('data-target');
  if (!targetId) return;
  const panel = document.getElementById(targetId);
  if (!panel) return;

  const isOpen = btn.getAttribute('data-open') === 'true';
  panel.classList.toggle('is-open', !isOpen);
  btn.setAttribute('data-open', isOpen ? 'false' : 'true');

  if (isOpen) {{
    btn.textContent = btn.textContent.replace(/^−[ \t]*/, '+');
  }} else {{
    btn.textContent = btn.textContent.replace(/^[+][ \t]*/, '− ');
  }}
}});
</script>

<div class="report-footer">
  <span>{esc(programa_nome)} — INS &nbsp;|&nbsp; Exercício Delphi 2026 &nbsp;|&nbsp; <strong>CONFIDENCIAL — Uso interno</strong></span>
  <span>Resultados anónimos. Nenhuma resposta individual é identificável. &nbsp;|&nbsp; Gerado em {now}</span>
</div>

</body>
</html>"""
    return html

# ─────────────────────────────────────────────────────────────────────────────
# PPTX GENERATION  (requires: pip install python-pptx cairosvg)
# ─────────────────────────────────────────────────────────────────────────────

# Slide canvas — 16:9 widescreen
_SW = _In(13.33) if _PPTX_OK else 0
_SH = _In(7.5)   if _PPTX_OK else 0

# Brand colours matching the HTML report
_C = {
    "accent":  _RGB(0xC0, 0x39, 0x2B),
    "accent2": _RGB(0x1A, 0x52, 0x76),
    "green":   _RGB(0x2E, 0x7D, 0x52),
    "ink":     _RGB(0x1A, 0x1A, 0x1A),
    "muted":   _RGB(0x78, 0x90, 0x9C),
    "bg":      _RGB(0xF5, 0xF7, 0xFA),
    "white":   _RGB(0xFF, 0xFF, 0xFF),
    "border":  _RGB(0xEC, 0xEF, 0xF1),
} if _PPTX_OK else {}


def _svg_to_png(svg_str, scale=2.0):
    """Convert an SVG string to PNG bytes; returns None if cairosvg unavailable."""
    if not _CAIRO_OK or not svg_str:
        return None
    svg_str = (svg_str
               .replace("DM Sans", "Arial")
               .replace("DM Serif Display", "Georgia"))
    try:
        return _cairosvg.svg2png(bytestring=svg_str.encode("utf-8"), scale=scale)
    except Exception:
        return None


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # fully blank layout


def _rect(slide, l, t, w, h, fill=None, line=None):
    sp = slide.shapes.add_shape(1, l, t, w, h)  # 1 = rectangle autoshape
    if fill:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    else:
        sp.fill.background()
    if line:
        sp.line.color.rgb = line
    else:
        sp.line.fill.background()
    return sp


def _text(slide, text, l, t, w, h, size=13, bold=False,
          color=None, align=None, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = _Pt(size)
    run.font.bold = bold
    run.font.name = "Calibri"
    run.font.color.rgb = color or _C["ink"]
    return tb


def _header_bar(slide, title, subtitle=""):
    _rect(slide, _In(0), _In(0), _SW, _In(0.75), fill=_C["accent2"])
    _text(slide, title, _In(0.28), _In(0.07), _In(11.5), _In(0.62),
          size=18, bold=True, color=_C["white"])
    if subtitle:
        _text(slide, subtitle, _In(9.5), _In(0.07), _In(3.5), _In(0.62),
              size=10, color=_RGB(0xBD, 0xC3, 0xC7), align=_PA.RIGHT)


def _picture(slide, png_bytes, l, t, w):
    if not png_bytes:
        return
    try:
        slide.shapes.add_picture(_io.BytesIO(png_bytes), l, t, width=w)
    except Exception:
        pass


def _slide_title(prs, title, subtitle=""):
    slide = _blank(prs)
    _rect(slide, _In(0), _In(0), _In(0.45), _SH, fill=_C["accent"])
    _rect(slide, _In(0.45), _In(0), _SW - _In(0.45), _SH, fill=_C["accent2"])
    _text(slide, title, _In(0.8), _In(2.1), _In(12), _In(1.9),
          size=34, bold=True, color=_C["white"])
    if subtitle:
        _text(slide, subtitle, _In(0.8), _In(4.0), _In(12), _In(1.0),
              size=14, color=_RGB(0xBD, 0xC3, 0xC7))
    return slide


def _slide_divider(prs, section, note=""):
    slide = _blank(prs)
    _rect(slide, _In(0), _In(0), _SW, _SH, fill=_C["bg"])
    _rect(slide, _In(0), _In(0), _In(0.15), _SH, fill=_C["accent"])
    _text(slide, section, _In(0.5), _In(2.8), _In(12.5), _In(1.4),
          size=30, bold=True, color=_C["accent2"])
    if note:
        _text(slide, note, _In(0.5), _In(4.15), _In(12.5), _In(0.8),
              size=13, color=_C["muted"])
    return slide


def _slide_chart(prs, title, png_bytes, note=""):
    slide = _blank(prs)
    _rect(slide, _In(0), _In(0), _SW, _SH, fill=_C["bg"])
    _header_bar(slide, title)
    _picture(slide, png_bytes, _In(0.2), _In(0.85), _SW - _In(0.4))
    if note:
        _text(slide, note, _In(0.2), _In(7.05), _SW - _In(0.4), _In(0.38),
              size=9, color=_C["muted"])
    return slide


def _slide_two_charts(prs, title, left_png, right_png,
                      left_lbl="", right_lbl=""):
    slide = _blank(prs)
    _rect(slide, _In(0), _In(0), _SW, _SH, fill=_C["bg"])
    _header_bar(slide, title)
    half = (_SW - _In(0.5)) / 2
    for i, (png, lbl) in enumerate([(left_png, left_lbl),
                                     (right_png, right_lbl)]):
        x = _In(0.15) + i * (half + _In(0.2))
        _picture(slide, png, x, _In(0.85), half)
        if lbl:
            _text(slide, lbl, x, _In(7.0), half, _In(0.42),
                  size=9, color=_C["muted"], align=_PA.CENTER)
    return slide


def _slide_table(prs, title, headers, rows, max_rows=20):
    """Paginated table — creates as many slides as needed."""
    total = len(rows)
    for chunk_start in range(0, max(total, 1), max_rows):
        chunk = rows[chunk_start: chunk_start + max_rows]
        slide = _blank(prs)
        _rect(slide, _In(0), _In(0), _SW, _SH, fill=_C["bg"])
        sfx = (f" [{chunk_start+1}–{min(chunk_start+max_rows, total)}/{total}]"
               if total > max_rows else "")
        _header_bar(slide, title + sfx)
        ncols = len(headers)
        nrows = len(chunk) + 1
        tbl = slide.shapes.add_table(
            nrows, ncols,
            _In(0.2), _In(0.85),
            _SW - _In(0.4),
            min(_In(6.3), _In(0.32 * nrows)),
        ).table
        for ci, h in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.text = str(h)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _C["accent2"]
            p = cell.text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = _Pt(9)
            run.font.bold = True
            run.font.color.rgb = _C["white"]
        for ri, row in enumerate(chunk, 1):
            bg = _C["white"] if ri % 2 else _C["bg"]
            for ci, val in enumerate(row):
                cell = tbl.cell(ri, ci)
                cell.text = str(val)
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
                p = cell.text_frame.paragraphs[0]
                run = p.runs[0] if p.runs else p.add_run()
                run.font.size = _Pt(8)
                run.font.color.rgb = _C["ink"]


def _slide_text(prs, title, body_lines):
    """Text-content slide with a title bar and bulleted lines."""
    slide = _blank(prs)
    _rect(slide, _In(0), _In(0), _SW, _SH, fill=_C["bg"])
    _header_bar(slide, title)
    y = _In(0.92)
    for line in body_lines[:18]:
        _text(slide, line, _In(0.5), y, _SW - _In(0.7), _In(0.38),
              size=12, color=_C["ink"])
        y += _In(0.36)
    return slide


def _slide_text_block(prs, title, body_lines, font_size=10):
    """Dense text slide for intervention card-like detail."""
    slide = _blank(prs)
    _rect(slide, _In(0), _In(0), _SW, _SH, fill=_C["bg"])
    _header_bar(slide, title)
    box = slide.shapes.add_textbox(_In(0.4), _In(0.9), _SW - _In(0.8), _SH - _In(1.25))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(body_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(line)
        if p.runs:
            run = p.runs[0]
            run.font.size = _Pt(font_size)
            run.font.name = "Calibri"
            run.font.color.rgb = _C["ink"]
    return slide
def _slide_intervention_card(prs, title, r, tier_num, tier_name, top_dup_txt, top_intg_txt,
                             dup_other_txt, intg_other_txt, comments):
    slide = _blank(prs)
    _rect(slide, _In(0), _In(0), _SW, _SH, fill=_C["bg"])
    _header_bar(slide, title, f"Código: {r.get('code', '')}")

    card_l, card_t, card_w, card_h = _In(0.35), _In(0.95), _In(12.63), _In(6.25)
    _rect(slide, card_l, card_t, card_w, card_h, fill=_C["white"], line=_C["border"])

    comp_txt = f"{r.get('programa', '') or '—'}  •  {r.get('component', '') or '—'}"
    _text(slide, comp_txt, card_l + _In(0.2), card_t + _In(0.12), _In(9.5), _In(0.26),
          size=10, color=_C["muted"])
    _text(slide, f"Grupo: {tier_num} ({tier_name})", card_l + _In(10.1), card_t + _In(0.12),
          _In(2.3), _In(0.26), size=10, bold=True, color=_C["accent2"], align=_PA.RIGHT)

    n_total = r.get("n_total", 0) or 0
    sim_p = round(r.get("n_sim", 0) / n_total * 100) if n_total else 0
    poss_p = round(r.get("n_poss", 0) / n_total * 100) if n_total else 0
    nao_p = max(0, 100 - sim_p - poss_p)

    _text(slide, f"Precisa de optimização? (n={n_total})", card_l + _In(0.2), card_t + _In(0.42),
          _In(5.5), _In(0.24), size=10, bold=True, color=_C["ink"])

    bar_l, bar_t, bar_w, bar_h = card_l + _In(0.2), card_t + _In(0.68), _In(8.6), _In(0.25)
    _rect(slide, bar_l, bar_t, bar_w, bar_h, fill=_RGB(0xEA, 0xEE, 0xF2), line=_RGB(0xEA, 0xEE, 0xF2))
    sim_w = _In(8.6 * sim_p / 100)
    poss_w = _In(8.6 * poss_p / 100)
    nao_w = _In(8.6 * nao_p / 100)
    if sim_w > 0:
        _rect(slide, bar_l, bar_t, sim_w, bar_h, fill=_RGB(0x1A, 0x6B, 0x3A), line=_RGB(0x1A, 0x6B, 0x3A))
    if poss_w > 0:
        _rect(slide, bar_l + sim_w, bar_t, poss_w, bar_h, fill=_RGB(0xD4, 0xA0, 0x17), line=_RGB(0xD4, 0xA0, 0x17))
    if nao_w > 0:
        _rect(slide, bar_l + sim_w + poss_w, bar_t, nao_w, bar_h, fill=_RGB(0xD1, 0xD5, 0xDB), line=_RGB(0xD1, 0xD5, 0xDB))

    _text(slide,
          f"Sim def.: {r.get('n_sim',0)} ({sim_p}%)   |   Possiv.: {r.get('n_poss',0)} ({poss_p}%)   |   Não: {r.get('n_nao',0)} ({nao_p}%)",
          card_l + _In(0.2), card_t + _In(0.96), _In(8.9), _In(0.24), size=9, color=_C["muted"])

    chip_top = card_t + _In(1.28)
    chip_w = _In(2.0)
    gap = _In(0.18)
    chips = [
        ("Impacto", f"{r.get('avg_impact',0):.2f}/3", _RGB(0x1A, 0x6B, 0x3A)),
        ("Duplicação", f"{r.get('dup_pct',0)}%", _C["ink"]),
        ("Integração", f"{r.get('intg_pct',0)}%", _C["accent2"]),
        ("↓ Recursos", f"{r.get('res_pct',0)}%", _C["muted"]),
    ]
    for i, (lbl, val, col) in enumerate(chips):
        x = card_l + _In(0.2) + i * (chip_w + gap)
        _rect(slide, x, chip_top, chip_w, _In(0.86), fill=_RGB(0xF8, 0xFA, 0xFC), line=_C["border"])
        _text(slide, val, x, chip_top + _In(0.13), chip_w, _In(0.34), size=16, bold=True, color=col, align=_PA.CENTER)
        _text(slide, lbl, x, chip_top + _In(0.50), chip_w, _In(0.22), size=9, color=_C["muted"], align=_PA.CENTER)

    score_y = chip_top + _In(1.02)
    _text(
        slide,
        f"S_base: {r.get('s_base',0):.3f}   |   S_pond: {r.get('s_pond',0):.3f}   |   Triagem médio: {r.get('gate_mean',0):.3f}   |   Índice eficiência: {r.get('e_score',0)*100:.1f}%",
        card_l + _In(0.2), score_y, _In(8.9), _In(0.26), size=10, bold=True, color=_C["accent2"]
    )

    _text(slide, f"Top duplicação (catálogo): {top_dup_txt}", card_l + _In(0.2), score_y + _In(0.36),
          _In(12.0), _In(0.32), size=9, color=_C["ink"])
    _text(slide, f"Top integração (catálogo): {top_intg_txt}", card_l + _In(0.2), score_y + _In(0.70),
          _In(12.0), _In(0.32), size=9, color=_C["ink"])
    _text(slide, f"Duplicação (outros): {dup_other_txt}", card_l + _In(0.2), score_y + _In(1.04),
          _In(12.0), _In(0.32), size=9, color=_C["muted"])
    _text(slide, f"Integração (outros): {intg_other_txt}", card_l + _In(0.2), score_y + _In(1.38),
          _In(12.0), _In(0.32), size=9, color=_C["muted"])
    _text(slide, f"URL da ficha: {r.get('url', '') or '—'}", card_l + _In(0.2), score_y + _In(1.72),
          _In(12.0), _In(0.28), size=8, color=_C["muted"])

    cmt_title_y = score_y + _In(2.08)
    _text(slide, f"Sugestões dos especialistas ({len(comments)}):", card_l + _In(0.2), cmt_title_y,
          _In(12.0), _In(0.24), size=9, bold=True, color=_C["ink"])
    shown = comments[:3]
    if shown:
        for i, c in enumerate(shown):
            _text(slide, f"• {c}", card_l + _In(0.25), cmt_title_y + _In(0.26 + 0.22 * i),
                  _In(12.0), _In(0.22), size=8, color=_C["ink"])
        if len(comments) > len(shown):
            _text(slide, f"• ... +{len(comments)-len(shown)} comentário(s)",
                  card_l + _In(0.25), cmt_title_y + _In(0.26 + 0.22 * len(shown)),
                  _In(12.0), _In(0.22), size=8, color=_C["muted"])
    else:
        _text(slide, "• —", card_l + _In(0.25), cmt_title_y + _In(0.26), _In(12.0), _In(0.22),
              size=8, color=_C["muted"])
    return slide


def build_pptx(results, stats, interventions, results_path,
               univariate, include_xyplot=True, include_scoring=True,
               simple_table=False):
    """
    Generate a PowerPoint presentation mirroring the HTML report content.
    Returns a Presentation object, or None if python-pptx is unavailable.
    Chart images require cairosvg; without it slides are created but charts
    will be blank (text and tables are always included).
    """
    if not _PPTX_OK:
        print("  ⚠ python-pptx not installed — skipping PPTX.")
        print("    Run: pip install python-pptx cairosvg")
        return None

    if not _CAIRO_OK:
        print("  ⚠ cairosvg not installed — PPTX will be generated without chart images.")
        print("    Run: pip install cairosvg")

    # Create display labels with numeric code prefix
    display_label = {r["code"]: hiv_report_label(r["code"], r["label"], r.get("component", "")) 
                     for r in interventions}

    prs = _Prs()
    prs.slide_width  = _SW
    prs.slide_height = _SH

    items = sorted(results.values(), key=lambda r: r.get("s_base", 0), reverse=True)
    n_inv = stats["n_inv"]
    n_exp = stats["n_experts"]

    sorted_inv = sorted(results.values(), key=lambda r: r.get("composite", 0), reverse=True)
    for idx, r in enumerate(sorted_inv, 1):
      r["display_idx"] = idx
    med_gate = statistics.median([r.get("gate_mean", 0) for r in sorted_inv]) if sorted_inv else 0
    med_imp  = statistics.median([r.get("avg_impact", 0) for r in sorted_inv]) if sorted_inv else 0

    def _tier(r):
      hg = r.get("gate_mean", 0) > med_gate
      hi = r.get("avg_impact", 0) > med_imp
      if hg and hi:
        return "alta"
      if hg or hi:
        return "media"
      return "baixa"

    def _tier_num(r):
      t = _tier(r)
      return 1 if t == "alta" else (2 if t == "media" else 3)

    def _group_summary_rows(group_key):
      buckets = OrderedDict()
      for r in results.values():
        gname = str(r.get(group_key, "") or "Outros").strip() or "Outros"
        buckets.setdefault(gname, []).append(r)
      rows = []
      for gname, vals in buckets.items():
        n = len(vals)
        spond = sum(x.get("s_pond", 0) for x in vals) / n if n else 0
        imp = sum(x.get("avg_impact", 0) for x in vals) / n if n else 0
        eff = sum(x.get("e_score", 0) for x in vals) / n if n else 0
        g1 = sum(1 for x in vals if _tier(x) == "alta")
        g2 = sum(1 for x in vals if _tier(x) == "media")
        g3 = sum(1 for x in vals if _tier(x) == "baixa")
        rows.append((gname, n, f"{spond:.3f}", f"{imp:.2f}", f"{eff*100:.1f}%", g1 or "—", g2 or "—", g3 or "—"))
      rows.sort(key=lambda x: float(x[2]), reverse=True)
      return rows

    # ── 1: Title ──────────────────────────────────────────────
    _slide_title(
        prs,
        "Delphi W1 — Relatório de Resultados",
        (f"{n_exp} especialistas  ·  {n_inv} intervenções  ·  "
         f"Taxa de resposta mediana: {stats['rr_median']}%  ·  "
         f"{datetime.now().strftime('%d/%m/%Y')}"),
    )

    # ── 2: Executive summary ──────────────────────────────────
    sum_slide = _blank(prs)
    _rect(sum_slide, _In(0), _In(0), _SW, _SH, fill=_C["bg"])
    _header_bar(sum_slide, "Resumo Executivo")
    metrics = [
        ("Especialistas",       str(n_exp)),
        ("Intervenções",        str(n_inv)),
        ("≥80% optimizáveis",   str(stats["n_inv_80"])),
        ("Impacto alto (≥2.5)", str(stats["n_imp_high"])),
        ("Consenso unânime",    str(stats["n_unanimous"])),
        ("Taxa resp. mediana",  f"{stats['rr_median']}%"),
    ]
    bw, bh = _In(1.95), _In(2.4)
    for i, (lbl, val) in enumerate(metrics):
        x = _In(0.27) + i * _In(2.17)
        y = _In(2.0)
        _rect(sum_slide, x, y, bw, bh, fill=_C["white"], line=_C["border"])
        _text(sum_slide, val, x, y + _In(0.25), bw, _In(1.1),
              size=38, bold=True, color=_C["accent2"], align=_PA.CENTER)
        _text(sum_slide, lbl, x, y + _In(1.45), bw, _In(0.75),
              size=10, color=_C["muted"], align=_PA.CENTER)

    # ── 3: Response rate table ────────────────────────────────
    _n_exp_leg = stats.get("n_experts", 1) or 1
    rr_rows = [
        ("", str(d.get("label", "")),
         f"{d.get('pct', 0)}%",
         f"{d.get('n', 0)}/{_n_exp_leg}")
        for d in stats.get("rr_detail", [])
    ]
    if rr_rows:
        _slide_table(prs, "Taxa de Resposta por Intervenção",
                     ["Código", "Intervenção", "Taxa", "Respostas"],
                     rr_rows, max_rows=22)
    rr_team_rows_leg = [
        (f"Equipa {d['team']}", str(d["n_inv"]), str(d.get("n_experts", "—")),
         f"{d['rr_median']}%", f"{d['rr_min']}%", f"{d['rr_max']}%")
        for d in stats.get("rr_by_team", [])
    ]
    if rr_team_rows_leg:
        _slide_table(prs, "Taxa de Resposta por Equipa",
                     ["Equipa", "Intervenções", "Especialistas", "Mediana", "Mín", "Máx"],
                     rr_team_rows_leg, max_rows=22)

    # ── 4: Program/component tabulations ─────────────────────
    programa_rows = _group_summary_rows("programa")
    if programa_rows:
      _slide_table(
        prs,
        "Resumo por Programa",
        ["Programa", "N", "S_pond médio", "Impacto médio", "Índice eficiência", "G1", "G2", "G3"],
        programa_rows,
        max_rows=20,
      )
    componente_rows = _group_summary_rows("comp_macro")
    if componente_rows:
      _slide_table(
        prs,
        "Resumo por Componente",
        ["Componente", "N", "S_pond médio", "Impacto médio", "Índice eficiência", "G1", "G2", "G3"],
        componente_rows,
        max_rows=20,
      )

    # ── Section: Univariate Analysis ──────────────────────────
    _slide_divider(prs, "Análise Univariada",
                   f"{n_inv} intervenções  ·  {n_exp} especialistas")

    gate_agg    = univariate.get("gate_agg", {})
    sorted_gate = univariate.get("sorted_gate",   items)
    sorted_imp  = univariate.get("sorted_impact", items)
    sorted_dup  = univariate.get("sorted_dup",    items)
    sorted_intg = univariate.get("sorted_intg",   items)
    sorted_res  = univariate.get("sorted_res",    items)
    imp_counts  = univariate.get("imp_counts", {})
    exp_counts  = univariate.get("exp_counts", {})

    # 2a: Gate distribution — donut + stacked bars
    gate_stacked_rows = [
      (display_label.get(r["code"], r["label"]), [r["n_sim"], r["n_poss"], r["n_nao"]])
        for r in sorted_gate
    ]
    gate_donut_svg = svg_donut(
        gate_agg or {"Sim def.": 1, "Possiv.": 1, "Não": 1},
        ["#2e7d52", "#d4a017", "#e0e0e0"],
        size=280,
    )
    _slide_two_charts(
        prs, "2a · Distribuição Triagem (Optimizabilidade)",
        _svg_to_png(gate_donut_svg),
        _svg_to_png(svg_hbar_stacked(gate_stacked_rows, width=580)),
        "Distribuição global", "Por intervenção",
    )

    # 2b: Expected impact — donut + bars
    max_imp = max((r["avg_impact"] for r in items), default=3.0) or 3.0
    impact_rows = [(display_label.get(r["code"], r["label"]), r["avg_impact"], max_imp) for r in sorted_imp]
    if imp_counts and any(imp_counts.values()):
        imp_labels = {1: "Baixo", 2: "Médio", 3: "Alto"}
        imp_donut_svg = svg_donut(
            {imp_labels.get(k, str(k)): v
             for k, v in sorted(imp_counts.items()) if v},
            ["#b91c1c", "#b45309", "#2e7d52"], size=280,
        )
        imp_donut_png = _svg_to_png(imp_donut_svg)
    else:
        imp_donut_png = None
    _slide_two_charts(
        prs, "2b · Impacto Esperado",
        imp_donut_png,
        _svg_to_png(svg_hbar_single(impact_rows, width=580)),
        "Distribuição global", "Impacto médio por intervenção",
    )

    # 2c: Expertise distribution
    if exp_counts and any(exp_counts.values()):
        exp_labels = {1: "Baixa", 2: "Média", 3: "Alta"}
        exp_donut_svg = svg_donut(
            {exp_labels.get(k, str(k)): v
             for k, v in sorted(exp_counts.items()) if v},
            ["#f59e0b", "#3b82f6", "#2e7d52"], size=340,
        )
        _slide_chart(prs, "2c · Nível de Experiência dos Especialistas",
                     _svg_to_png(exp_donut_svg))

    # 2d–2f: Duplication, Integration, Resource reduction
    for slide_title, field, src in [
        ("2d · % Sobreposição com Outras Intervenções", "dup_pct",  sorted_dup),
        ("2e · % Potencial de Integração",              "intg_pct", sorted_intg),
        ("2f · % Possibilidade de Redução de Recursos", "res_pct",  sorted_res),
    ]:
        bar_rows = [(display_label.get(r["code"], r["label"]), r[field], 100) for r in src]
        _slide_chart(prs, slide_title,
                     _svg_to_png(svg_hbar_single(bar_rows, width=900)))

    # 2g: Scatter plot
    if include_xyplot:
        _slide_chart(
            prs, "2g · Dispersão: Optimizabilidade × Impacto Esperado",
            _svg_to_png(svg_scatter_optim_impact_exp(items), scale=1.3),
        )

    # ── Section: Scoring ──────────────────────────────────────
    if include_scoring:
        _slide_divider(prs, "Pontuação e Metodologia")

        _slide_text(prs, "Fórmulas de Pontuação", [
            "",
            "S_base  =  média(triagem_score_i)  ×  média(impacto_i)",
            "",
            "S_pond  =  Σ(gate_i × impacto_i × exp_i)  /  Σ(exp_i)",
            "",
            "Onde:",
            "   triagem_score:  Sim definitivamente = 1,  Possivelmente = 0.5,  Não = 0",
            "   impacto:     Baixo = 1,  Médio = 2,  Alto = 3",
            "   exp (ponderação):  Baixa = 1,  Média = 2,  Alta = 3",
        ])

        s_max = max(
            max((r["s_base"] for r in items), default=1.0),
            max((r["s_pond"] for r in items), default=1.0),
        ) or 1.0
        sbase_rows = [(display_label.get(r["code"], r["label"]), r["s_base"], s_max)
                      for r in sorted(items, key=lambda r: r["s_base"], reverse=True)]
        spond_rows = [(display_label.get(r["code"], r["label"]), r["s_pond"], s_max)
                      for r in sorted(items, key=lambda r: r["s_pond"], reverse=True)]
        _slide_two_charts(
            prs, "Pontuações S_base e S_pond",
            _svg_to_png(svg_hbar_single(sbase_rows, width=560, color="#1a5276")),
            _svg_to_png(svg_hbar_single(spond_rows, width=560, color="#2e7d52")),
            "S_base (não ponderado)", "S_pond (ponderado por experiência)",
        )

        _slide_chart(
            prs, "Transição de Rankings: Optimizabilidade → S_base → S_pond",
            _svg_to_png(svg_alluvial_weighting(items), scale=1.2),
        )

        if simple_table:
          score_rows = [
            (
              display_label.get(r["code"], r["label"]),
              f"{r.get('s_pond', 0):.3f}",
              f"{r.get('e_score', 0)*100:.1f}%",
              str(_tier_num(r)),
            )
            for r in sorted(results.values(), key=lambda x: x.get("s_pond", 0), reverse=True)
          ]
          _slide_table(
            prs,
            "Tabela de Pontuações (Simples)",
            ["Intervenção", "S_pond", "Índice de eficiência", "Grupo Próximos Passos (1/2/3)"],
            score_rows,
            max_rows=22,
          )
        else:
          score_rows = [
            (r.get("rank_base", "—"),
             display_label.get(r["code"], r["label"]),
             r.get("component", "")[:22] or "—",
             f"{r['s_base']:.3f}",
             f"{r['s_pond']:.3f}",
             (f"+{r.get('rank_delta',0)}"
              if r.get("rank_delta", 0) > 0 else str(r.get("rank_delta", 0))),
             f"{r.get('avg_impact',0):.2f}",
             f"{r.get('exp_mean',0):.2f}")
            for r in sorted(items, key=lambda r: r.get("rank_base", 999))
          ]
          _slide_table(prs, "Tabela de Pontuações",
                 ["#", "Intervenção", "Componente",
                  "S_base", "S_pond", "Δrank", "Impacto", "Exp."],
                 score_rows, max_rows=22)

    # ── Section: Priority Ranking ─────────────────────────────
    _slide_divider(prs, "Tabela de Priorização",
                   "Ordenado por S_base (pontuação composta)")

    priority_rows = [
        (r.get("rank_base", "—"),
       display_label.get(r["code"], r["label"]),
         r.get("component", "")[:22] or "—",
         f"{r['pct_optimizable']}%",
         f"{r.get('avg_impact',0):.1f}",
         f"{r['dup_pct']}%",
         f"{r['intg_pct']}%")
        for r in sorted(items, key=lambda r: r.get("rank_base", 999))
    ]
    _slide_table(prs, "Ranking de Prioridades",
                 ["#", "Intervenção", "Componente",
                  "% Optim.", "Impacto", "% Dup.", "% Integr."],
                 priority_rows, max_rows=20)

    # ── Section: Next Steps ───────────────────────────────────

    alta  = [display_label.get(r['code'], r['label'])
             for r in sorted_inv if _tier(r) == "alta"]
    media = [display_label.get(r['code'], r['label'])
             for r in sorted_inv if _tier(r) == "media"]
    baixa = [display_label.get(r['code'], r['label'])
             for r in sorted_inv if _tier(r) == "baixa"]

    _slide_divider(prs, "Próximos Passos", "Classificação por prioridade")
    for tier_title, tier_items in [
        ("Alta Prioridade — Optimizar e reforçar",    alta),
        ("Prioridade Média — Avaliar caso a caso",    media),
        ("Baixa Prioridade — Manter ou desinvestir",  baixa),
    ]:
        if tier_items:
            _slide_text(prs, tier_title, tier_items)

    # ── Section: Intervention Details by Component ────────────
    _slide_divider(prs, "Fichas por Componente",
                   "Resumo detalhado por intervenção")

    from collections import OrderedDict as _OD
    by_comp = _OD()
    for intv in interventions:
        comp = intv.get("component", "") or "Outros"
        by_comp.setdefault(comp, []).append(intv)

    def _fmt_tagged_codes(codes):
      if not codes:
        return "—"
      return " · ".join(display_label.get(c, c) for c in codes[:6])

    def _fmt_count_map(count_map, limit=8):
      if not count_map:
        return "—"
      pairs = sorted(count_map.items(), key=lambda kv: (-kv[1], str(kv[0]).lower()))[:limit]
      return "; ".join(f"{k} ({v})" for k, v in pairs)

    for comp, comp_intvs in by_comp.items():
      if not comp_intvs:
        continue
      _slide_divider(prs, f"Componente: {comp}", "Detalhe por intervenção")
      for intv in comp_intvs:
        r = results.get(intv["code"])
        if not r:
          continue
        comments = r.get("comments", []) or []

        _slide_intervention_card(
          prs,
          display_label.get(intv['code'], intv['label']),
          r,
          _tier_num(r),
          _tier(r),
          _fmt_tagged_codes(r.get('top_dup', [])),
          _fmt_tagged_codes(r.get('top_intg', [])),
          _fmt_count_map(r.get('dup_other_counts', {}), limit=8),
          _fmt_count_map(r.get('intg_other_counts', {}), limit=8),
          comments,
        )

    return prs


# ─────────────────────────────────────────────────────────────────────────────
# TEMPLATE-DRIVEN PPTX PIPELINE
# ─────────────────────────────────────────────────────────────────────────────

# ── Style defaults and override machinery ────────────────────────────────────

_PPTX_DEFAULTS = {
    "chart.top_n":                18,
    "chart.label.max_chars":      56,
    "chart.font.family":          "DM Sans",
    "chart.font.size":            11,
    "chart.label.wrap":           False,
    "chart.truncation.notice":    True,
    "table.font.family":          "DM Sans",
    "table.header.font.size":     10,
    "table.body.font.size":       10,
    "table.min.bottom_margin.in": 0.20,
    "card.font.family":           "DM Sans",
    "card.title.font.size":       16,
    "card.body.font.size":        11,
    "card.comments.max_shown":    5,
}

# Maps config/env key names to style dict keys
_PPTX_CFG_KEY_MAP = {
    "PPTX_CHART_TOP_N":              "chart.top_n",
    "PPTX_CHART_LABEL_MAX_CHARS":    "chart.label.max_chars",
    "PPTX_CHART_FONT_FAMILY":        "chart.font.family",
    "PPTX_CHART_FONT_SIZE":          "chart.font.size",
    "PPTX_CHART_LABEL_WRAP":         "chart.label.wrap",
    "PPTX_CHART_TRUNCATION_NOTICE":  "chart.truncation.notice",
    "PPTX_TABLE_FONT_FAMILY":        "table.font.family",
    "PPTX_TABLE_HEADER_FONT_SIZE":   "table.header.font.size",
    "PPTX_TABLE_BODY_FONT_SIZE":     "table.body.font.size",
    "PPTX_CARD_FONT_FAMILY":         "card.font.family",
    "PPTX_CARD_TITLE_FONT_SIZE":     "card.title.font.size",
    "PPTX_CARD_BODY_FONT_SIZE":      "card.body.font.size",
    "PPTX_CARD_COMMENTS_MAX_SHOWN":  "card.comments.max_shown",
}


def parse_pptx_style_tokens(prs):
    """
    Parse PPTX_STYLE_TOKENS blocks from template reference style slides.
    Returns dict mapping token_key -> raw string value.
    Each line after the first in a PPTX_STYLE_TOKENS block must be: key=value.
    """
    tokens = {}
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text.startswith("PPTX_STYLE_TOKENS"):
                continue
            for line in text.splitlines()[1:]:
                line = line.strip()
                if "=" in line:
                    k, _, v = line.partition("=")
                    tokens[k.strip()] = v.strip()
    return tokens


def resolve_pptx_style(cli_overrides=None, cfg=None, token_overrides=None):
    """
    Resolve PPTX style settings applying precedence (highest first):
      1. CLI overrides  2. Config/env  3. Template tokens  4. Hardcoded defaults
    Returns a fully-populated style dict.
    """
    result = dict(_PPTX_DEFAULTS)

    def _cast(target_val, new_val):
        try:
            if isinstance(target_val, bool):
                if isinstance(new_val, bool): return new_val
                return str(new_val).lower() in ("1", "true", "yes")
            if isinstance(target_val, int):   return int(new_val)
            if isinstance(target_val, float): return float(new_val)
            return str(new_val)
        except (ValueError, TypeError):
            return target_val

    # Level 4 → 3 → 2: apply in order so higher precedence wins last
    if token_overrides:
        for k, v in token_overrides.items():
            if k in result:
                result[k] = _cast(result[k], v)
    if cfg:
        for cfg_key, style_key in _PPTX_CFG_KEY_MAP.items():
            if cfg_key in cfg:
                result[style_key] = _cast(result[style_key], cfg[cfg_key])
    if cli_overrides:
        for style_key, v in cli_overrides.items():
            if style_key in result and v is not None:
                result[style_key] = _cast(result[style_key], v)

    return result


def get_layout(prs, preferred_names, fallback_idx=0):
    """Resolve a slide layout by preferred name(s) with index fallback."""
    for name in preferred_names:
        for layout in prs.slide_layouts:
            if layout.name == name:
                return layout
    layouts = prs.slide_layouts
    return layouts[min(fallback_idx, len(layouts) - 1)]


def _get_ph(slide, idx):
    """Return the first placeholder with the given index, or None."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _content_bounds(slide, prs):
    """Return (l, t, w, h) using content placeholder (idx 1) or slide defaults."""
    ph = _get_ph(slide, 1)
    if ph is not None:
        return ph.left, ph.top, ph.width, ph.height
    return _In(0.2), _In(0.85), prs.slide_width - _In(0.4), prs.slide_height - _In(1.1)


def _set_title_ph(slide, text):
    """Write text into the title placeholder (idx 0) if present."""
    ph = _get_ph(slide, 0)
    if ph is not None:
        ph.text = text


def add_title_slide_from_layout(prs, layout, title, subtitle=""):
    slide = prs.slides.add_slide(layout)
    _set_title_ph(slide, title)
    sub_ph = _get_ph(slide, 1) or _get_ph(slide, 2)
    if sub_ph is not None and subtitle:
        sub_ph.text = subtitle
    return slide


def add_section_slide_from_layout(prs, layout, title, note=""):
    slide = prs.slides.add_slide(layout)
    _set_title_ph(slide, title)
    if note:
        sub_ph = _get_ph(slide, 1) or _get_ph(slide, 2)
        if sub_ph is not None:
            sub_ph.text = note
    return slide


def _fit_picture(slide, png_bytes, l, t, max_w, max_h):
    """Insert picture scaled to fit within (max_w × max_h) while preserving aspect ratio."""
    pic = slide.shapes.add_picture(_io.BytesIO(png_bytes), l, t, width=max_w)
    if pic.height > max_h:
        ratio = max_h / pic.height
        pic.height = max_h
        pic.width  = int(pic.width * ratio)
    return pic


def add_visual_slide_from_layout(prs, layout, title, png_bytes, note=""):
    slide = prs.slides.add_slide(layout)
    _set_title_ph(slide, title)
    if not png_bytes:
        return slide
    l, t, w, h = _content_bounds(slide, prs)
    note_h = _In(0.22) if note else 0
    try:
        _fit_picture(slide, png_bytes, l, t, w, h - note_h)
    except Exception:
        pass
    if note:
        _text(slide, note, l, t + h - note_h, w, note_h,
              size=11, color=_C["muted"], align=_PA.LEFT)
    return slide


def add_card_slide_from_layout(prs, layout, png_bytes):
    """Full-page card slide — title placeholder is suppressed, image fills the slide."""
    slide = prs.slides.add_slide(layout)
    # Clear title so it does not overlay the card image
    title_ph = _get_ph(slide, 0)
    if title_ph is not None:
        title_ph.text = ""
    if not png_bytes:
        return slide
    # Fill the entire slide canvas
    try:
        _fit_picture(slide, png_bytes, 0, 0, prs.slide_width, prs.slide_height)
    except Exception:
        pass
    return slide


def add_two_charts_slide_from_layout(prs, layout, title, left_png, right_png,
                                      left_lbl="", right_lbl=""):
    slide = prs.slides.add_slide(layout)
    _set_title_ph(slide, title)
    l, t, w, h = _content_bounds(slide, prs)
    half_w = (w - _In(0.1)) // 2
    for i, (png, lbl) in enumerate([(left_png, left_lbl), (right_png, right_lbl)]):
        if not png:
            continue
        x = l + i * (half_w + _In(0.1))
        try:
            _fit_picture(slide, png, x, t, half_w, h)
        except Exception:
            pass
    return slide


def add_table_slide_from_layout(prs, layout, title, headers, rows,
                                 max_rows=None, style=None, col_widths=None):
    """
    Paginated native table — creates as many slides as needed.
    Row count per slide is computed dynamically from the frame height so that
    the table never overflows the slide bounds.  max_rows acts as an upper cap.
    """
    st = style or _PPTX_DEFAULTS
    hdr_pt  = st["table.header.font.size"]
    body_pt = st["table.body.font.size"]
    font_nm = st["table.font.family"]
    margin  = _In(st["table.min.bottom_margin.in"])

    # Fixed row heights (EMU) used for layout calculation
    _HDR_ROW_H  = _In(0.34)
    _BODY_ROW_H = _In(0.29)

    # First slide is used to read frame bounds for pagination calculation
    def _compute_max_rows(frame_h):
        usable = frame_h - _HDR_ROW_H - margin
        return max(1, int(usable / _BODY_ROW_H))

    total = len(rows)
    chunk_start = 0
    while chunk_start < max(total, 1):
        slide = prs.slides.add_slide(layout)
        l, t, w, h = _content_bounds(slide, prs)

        # Compute capacity for this frame
        computed_cap = _compute_max_rows(h)
        effective_cap = computed_cap if max_rows is None else min(max_rows, computed_cap)

        chunk = rows[chunk_start: chunk_start + effective_cap]
        sfx = (f" [{chunk_start+1}–{min(chunk_start+effective_cap, total)}/{total}]"
               if total > effective_cap else "")
        _set_title_ph(slide, title + sfx)

        ncols = len(headers)
        nrows = len(chunk) + 1
        tbl_h = _HDR_ROW_H + _BODY_ROW_H * len(chunk)
        tbl   = slide.shapes.add_table(nrows, ncols, l, t, w, tbl_h).table

        # Apply deterministic column widths when provided
        if col_widths and len(col_widths) == ncols:
            total_specified = sum(_In(cw) for cw in col_widths)
            scale = w / total_specified if total_specified > 0 else 1.0
            for ci, cw in enumerate(col_widths):
                tbl.columns[ci].width = int(_In(cw) * scale)

        for ci, hdr in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.text = str(hdr)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _C["accent2"]
            p   = cell.text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size  = _Pt(hdr_pt)
            run.font.bold  = True
            run.font.name  = font_nm
            run.font.color.rgb = _C["white"]
        for ri, row in enumerate(chunk, 1):
            bg = _C["white"] if ri % 2 else _C["bg"]
            for ci, val in enumerate(row):
                cell = tbl.cell(ri, ci)
                # Strip embedded newlines to prevent multi-paragraph splits
                # which would inherit the template's default (larger) font.
                clean = str(val).replace("\r\n", " ").replace("\n", " ").replace("\r", " ")
                cell.text = clean
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
                # Apply font to every paragraph and every run so that word-
                # wrapped text does not pick up a different size from the master.
                for p in cell.text_frame.paragraphs:
                    for run in (p.runs if p.runs else [p.add_run()]):
                        run.font.size      = _Pt(body_pt)
                        run.font.name      = font_nm
                        run.font.color.rgb = _C["ink"]

        chunk_start += effective_cap


def add_text_slide_from_layout(prs, layout, title, body_lines, font_size=12):
    slide = prs.slides.add_slide(layout)
    _set_title_ph(slide, title)
    l, t, w, h = _content_bounds(slide, prs)
    body_ph = _get_ph(slide, 1)
    line_h = _In(max(0.32, font_size / 28.0))
    if body_ph is not None:
        tf = body_ph.text_frame
        tf.clear()
        for i, line in enumerate(body_lines[:18]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = str(line)
            if p.runs:
                p.runs[0].font.size = _Pt(font_size)
    else:
        y = t
        for line in body_lines[:18]:
            _text(slide, line, l, y, w, line_h, size=font_size, color=_C["ink"])
            y += line_h
    return slide


def _mini_pct_bar_svg(label, value, max_val, color, width=290, bar_h=14):
    """Single-row horizontal bar SVG suitable for card mini-panels."""
    label_w  = 72
    val_w    = 38
    bar_area = width - label_w - val_w
    fill_w   = int(bar_area * min(value, max_val) / max_val) if max_val else 0
    val_str  = f"{value:.2f}" if isinstance(value, float) and max_val <= 3.0 else (
               f"{value:.0f}%" if max_val == 100 else
               (f"{value:.1f}" if isinstance(value, float) else str(value)))
    return (
        f'<svg width="{width}" height="{bar_h + 6}" xmlns="http://www.w3.org/2000/svg">'
        f'<text x="0" y="{bar_h}" font-size="10" fill="#78909c">{label}</text>'
        f'<rect x="{label_w}" y="2" width="{bar_area}" height="{bar_h - 2}" fill="#eaeff2" rx="2"/>'
        f'<rect x="{label_w}" y="2" width="{fill_w}" height="{bar_h - 2}" fill="{color}" rx="2"/>'
        f'<text x="{label_w + bar_area + 4}" y="{bar_h}" font-size="10" fill="#1a1a1a">{val_str}</text>'
        f'</svg>'
    )


def _mini_gate_strip_svg(n_sim, n_poss, n_nao, width=290, bar_h=16):
    """1-D stacked bar showing Sim/Possiv./Não gate distribution (like HTML report)."""
    total = (n_sim + n_poss + n_nao) or 1
    sim_p  = round(n_sim  / total * 100)
    poss_p = round(n_poss / total * 100)
    nao_p  = max(0, 100 - sim_p - poss_p)
    sim_w  = int(width * sim_p  / 100)
    poss_w = int(width * poss_p / 100)
    nao_w  = max(0, width - sim_w - poss_w)
    h = bar_h + 22
    return (
        f'<svg width="{width}" height="{h}" xmlns="http://www.w3.org/2000/svg">'
        f'<rect x="0"           y="0" width="{sim_w}"  height="{bar_h}" fill="#2e7d52" rx="2"/>'
        f'<rect x="{sim_w}"     y="0" width="{poss_w}" height="{bar_h}" fill="#d4a017"/>'
        f'<rect x="{sim_w+poss_w}" y="0" width="{nao_w}" height="{bar_h}" fill="#e0e0e0" rx="2"/>'
        f'<text x="2" y="{bar_h + 14}" font-size="9" fill="#2e7d52" font-weight="600">Sim {sim_p}%</text>'
        f'<text x="{width//2}" y="{bar_h + 14}" text-anchor="middle" font-size="9" fill="#b07d0b">Poss. {poss_p}%</text>'
        f'<text x="{width - 2}" y="{bar_h + 14}" text-anchor="end" font-size="9" fill="#78909c">Não {nao_p}%</text>'
        f'</svg>'
    )


def _build_card_html_fragment(title, r, tier_num, tier_name,
                               top_dup_txt, top_intg_txt,
                               dup_other_txt, intg_other_txt, comments,
                               style=None, all_items=None):
    """
    Build a self-contained HTML page for a single intervention card.
    Layout mirrors the template card style slide geometry:
      left panel (~65%): main content, metrics, commentary
      right panel (~33%): univariate distribution mini-graphs (1D histograms)
    """
    st       = style or _PPTX_DEFAULTS
    font_fam = st["card.font.family"]
    title_pt = st["card.title.font.size"]
    body_pt  = st["card.body.font.size"]

    n_total = r.get("n_total", 0) or 0
    sim_p   = round(r.get("n_sim",  0) / n_total * 100) if n_total else 0
    poss_p  = round(r.get("n_poss", 0) / n_total * 100) if n_total else 0
    nao_p   = max(0, 100 - sim_p - poss_p)
    comp_txt = f"{r.get('programa', '') or '—'} · {r.get('component', '') or '—'}"
    code     = r.get("code", "")
    tier_col = {"alta": "#1a6b3a", "media": "#d4a017", "baixa": "#78909c"}.get(tier_name, "#78909c")

    # ── Left panel: KPI chips ──
    chips_html = ""
    for lbl, val, col in [
        ("Impacto",    f"{r.get('avg_impact',0):.2f}/3", "#1a6b3a"),
        ("Duplicação", f"{r.get('dup_pct',0)}%",         "#1a1a1a"),
        ("Integração", f"{r.get('intg_pct',0)}%",        "#1a5276"),
        ("↓ Recursos", f"{r.get('res_pct',0)}%",         "#78909c"),
    ]:
        chips_html += (
            f'<div style="flex:1;background:#f8fafc;border:1px solid #eceff1;'
            f'padding:8px 4px;text-align:center;">'
            f'<div style="font-size:20px;font-weight:bold;color:{col};">{val}</div>'
            f'<div style="font-size:{body_pt}px;color:#78909c;margin-top:3px;">{lbl}</div>'
            f'</div>'
        )

    # ── Left panel: comments ──
    max_shown = int(st.get("card.comments.max_shown", 5))
    cmt_html = ""
    shown = comments[:max_shown]
    for c in shown:
        cmt_html += f'<div>• {str(c)[:140]}</div>'
    if len(comments) > len(shown):
        cmt_html += f'<div style="color:#78909c;">• ... +{len(comments)-len(shown)} comentário(s)</div>'
    if not shown:
        cmt_html = '<div style="color:#78909c;">• —</div>'

    # ── Right panel: univariate distributions ──
    # Gate: 1D stacked bar (like HTML report's gate-bar-wrap)
    gate_strip_svg = _mini_gate_strip_svg(
        r.get("n_sim", 0) or 0,
        r.get("n_poss", 0) or 0,
        r.get("n_nao", 0) or 0,
        width=290,
    )

    # Score strips: show position in the full distribution (requires all_items)
    all_s_base  = [x.get("s_base",  0) for x in all_items] if all_items else []
    all_s_pond  = [x.get("s_pond",  0) for x in all_items] if all_items else []
    all_impact  = [x.get("avg_impact", 0) for x in all_items] if all_items else []
    all_dup     = [x.get("dup_pct",  0) for x in all_items] if all_items else []
    all_intg    = [x.get("intg_pct", 0) for x in all_items] if all_items else []
    all_res     = [x.get("res_pct",  0) for x in all_items] if all_items else []

    # svg_score_strip: distribution strip with current value marked
    sbase_strip = svg_score_strip(all_s_base,  r.get("s_base",  0), "#1a5276", width=290, height=30)
    spond_strip = svg_score_strip(all_s_pond,  r.get("s_pond",  0), "#2e7d52", width=290, height=30)
    impact_strip= svg_score_strip(all_impact,  r.get("avg_impact", 0), "#1a6b3a", width=290, height=30)
    dup_strip   = svg_score_strip(all_dup,     r.get("dup_pct",  0), "#1a5276", width=290, height=30)
    intg_strip  = svg_score_strip(all_intg,    r.get("intg_pct", 0), "#d4a017", width=290, height=30)
    res_strip   = svg_score_strip(all_res,     r.get("res_pct",  0), "#78909c", width=290, height=30)

    def _strip_row(label, strip_svg, val_str):
        return (
            f'<div style="margin-bottom:8px;">'
            f'<div style="font-size:9px;color:#78909c;margin-bottom:2px;">'
            f'{label} <span style="font-weight:600;color:#1a1a1a;">{val_str}</span></div>'
            f'{strip_svg}'
            f'</div>'
        )

    return f"""<!DOCTYPE html><html><head><meta charset="utf-8"><style>
*{{box-sizing:border-box;margin:0;padding:0;font-family:'{font_fam}',Arial,sans-serif;}}
body{{background:#1a5276;padding:10px;height:760px;overflow:hidden;}}
.header{{background:#1a5276;color:#fff;padding:8px 14px;display:flex;justify-content:space-between;
  align-items:center;margin-bottom:8px;}}
.header .title{{font-size:{title_pt}px;font-weight:bold;flex:1;margin-right:10px;line-height:1.3;}}
.badge{{background:{tier_col};color:#fff;padding:3px 10px;border-radius:3px;
  font-size:{body_pt}px;white-space:nowrap;}}
.body-row{{display:flex;gap:8px;height:calc(100% - 56px);}}
.left{{flex:2;background:#fff;border-radius:5px;padding:13px;min-width:0;overflow:hidden;}}
.right{{flex:0 0 310px;display:flex;flex-direction:column;gap:8px;overflow:hidden;}}
.rbox{{background:#fff;border-radius:5px;padding:10px 12px;}}
.rbox-hd{{font-size:{body_pt}px;font-weight:bold;color:#1a5276;margin-bottom:8px;
  border-bottom:1px solid #eceff1;padding-bottom:4px;}}
.muted{{color:#78909c;font-size:{body_pt}px;}}
.lbl{{font-size:{body_pt}px;font-weight:bold;color:#1a1a1a;margin:9px 0 4px;}}
.bar-wrap{{background:#eaeff2;height:14px;border-radius:2px;overflow:hidden;display:flex;}}
.chips{{display:flex;gap:5px;margin:9px 0;}}
.score{{font-size:{body_pt}px;font-weight:bold;color:#1a5276;margin:7px 0;line-height:1.5;}}
.det{{font-size:{body_pt}px;color:#1a1a1a;margin:3px 0;}}
.mdet{{font-size:{body_pt}px;color:#78909c;margin:3px 0;}}
.cmt-hd{{font-size:{body_pt}px;font-weight:bold;color:#1a1a1a;margin-top:9px;margin-bottom:3px;}}
.cmt{{font-size:{body_pt}px;color:#1a1a1a;line-height:1.5;}}
</style></head><body>
<div class="header">
  <div class="title">{title}</div>
  <div class="badge">Grupo {tier_num} ({tier_name})</div>
</div>
<div class="body-row">
  <!-- LEFT PANEL -->
  <div class="left">
    <div class="muted">{comp_txt} &nbsp;·&nbsp; {code}</div>
    <div class="lbl">Precisa de optimização? (n={n_total})</div>
    <div class="bar-wrap">
      <div style="width:{sim_p}%;background:#1a6b3a;"></div>
      <div style="width:{poss_p}%;background:#d4a017;"></div>
      <div style="width:{nao_p}%;background:#d1d5db;"></div>
    </div>
    <div class="muted" style="margin-top:4px;">
      Sim def.: {r.get('n_sim',0)} ({sim_p}%) &nbsp;|&nbsp;
      Possiv.: {r.get('n_poss',0)} ({poss_p}%) &nbsp;|&nbsp;
      Não: {r.get('n_nao',0)} ({nao_p}%)
    </div>
    <div class="chips">{chips_html}</div>
    <div class="score">
      S_base: {r.get('s_base',0):.3f} &nbsp;·&nbsp;
      S_pond: {r.get('s_pond',0):.3f} &nbsp;·&nbsp;
      Triagem médio: {r.get('gate_mean',0):.3f} &nbsp;·&nbsp;
      Índice ef.: {r.get('e_score',0)*100:.1f}%
    </div>
    <div class="det">Top dup. (cat.): {top_dup_txt}</div>
    <div class="det">Top intg. (cat.): {top_intg_txt}</div>
    <div class="mdet">Dup. (outros): {dup_other_txt}</div>
    <div class="mdet">Intg. (outros): {intg_other_txt}</div>
    <div class="mdet" style="margin-top:6px;">URL: {r.get('url','') or '—'}</div>
    <div class="cmt-hd">Sugestões ({len(comments)}):</div>
    <div class="cmt">{cmt_html}</div>
  </div>
  <!-- RIGHT PANEL: univariate distributions (1D histograms) -->
  <div class="right">
    <div class="rbox">
      <div class="rbox-hd">Distribuição Triagem</div>
      {gate_strip_svg}
    </div>
    <div class="rbox">
      <div class="rbox-hd">Pontuações (posição relativa)</div>
      {_strip_row("S_base", sbase_strip, f"{r.get('s_base',0):.3f}")}
      {_strip_row("S_pond", spond_strip, f"{r.get('s_pond',0):.3f}")}
    </div>
    <div class="rbox">
      <div class="rbox-hd">Indicadores (posição relativa)</div>
      {_strip_row("Impacto", impact_strip, f"{r.get('avg_impact',0):.2f}/3")}
      {_strip_row("Dup. %",  dup_strip,   f"{r.get('dup_pct',0)}%")}
      {_strip_row("Intg. %", intg_strip,  f"{r.get('intg_pct',0)}%")}
      {_strip_row("Rec. %",  res_strip,   f"{r.get('res_pct',0)}%")}
    </div>
  </div>
</div>
</body></html>"""


def render_card_html_to_png(card_html, width_px=1280, height_px=720):
    """
    Render card HTML to PNG bytes using Playwright.
    Returns None if no browser is available (caller falls back to legacy layout).

    Browser selection order (macOS note: standalone Chromium homebrew is deprecated):
      1. System Google Chrome  (channel="chrome") — preferred; auto-detected from
         /Applications/Google Chrome.app when PLAYWRIGHT_SKIP_BROWSER_DOWNLOAD=1
      2. Playwright Chromium   — fallback if system Chrome not found
      3. Firefox               — last resort
    """
    if not _PLAYWRIGHT_OK:
        return None

    def _screenshot(browser_obj):
        page = browser_obj.new_page(viewport={"width": width_px, "height": height_px})
        page.set_content(card_html, wait_until="domcontentloaded")
        png = page.screenshot(full_page=False)
        browser_obj.close()
        return png

    with _sync_playwright() as p:
        # 1. System Google Chrome
        try:
            return _screenshot(p.chromium.launch(channel="chrome"))
        except Exception:
            pass

        # 2. Playwright Chromium
        try:
            return _screenshot(p.chromium.launch())
        except Exception:
            pass

        # 3. Firefox
        try:
            return _screenshot(p.firefox.launch())
        except Exception as exc:
            print(f"  ⚠ Playwright: no usable browser found — {exc}")
            print("    Install Google Chrome or run: python -m playwright install chromium")
            return None


def build_pptx_template(results, stats, interventions, results_path,
                         univariate, template_path,
                         include_xyplot=True, include_scoring=True,
                         simple_table=False,
                         cli_overrides=None, cfg=None):
    """
    Template-driven PPTX generation.
    Loads a .pptx template for layout and typography; renders figures and
    intervention cards as images; keeps tables as native editable objects.
    Falls back to legacy card slide builder if Playwright is unavailable.

    Style resolution precedence (highest first):
      CLI overrides  >  config/env  >  PPTX_STYLE_TOKENS in template  >  defaults
    """
    if not _PPTX_OK:
        print("  ⚠ python-pptx not installed — skipping PPTX.")
        return None
    if not _CAIRO_OK:
        print("  ⚠ cairosvg not installed — chart images will be blank.")

    prs = _Prs(template_path)

    # Resolve style: parse template tokens, apply precedence chain
    tokens = parse_pptx_style_tokens(prs)
    style  = resolve_pptx_style(cli_overrides=cli_overrides, cfg=cfg,
                                 token_overrides=tokens)
    top_n           = style["chart.top_n"]
    max_ch          = style["chart.label.max_chars"]
    trunc_notice    = style["chart.truncation.notice"]

    def _trunc(s):
        s = str(s)
        return (s[:max_ch] + "…") if len(s) > max_ch else s

    def _trunc_note(n_shown, n_total):
        """Return a truncation notice string when data was capped, else empty."""
        if trunc_notice and n_shown < n_total:
            return f"Top {n_shown} de {n_total} intervenções mostradas"
        return ""

    lay_title   = get_layout(prs, ["Title Slide"],        fallback_idx=0)
    lay_section = get_layout(prs, ["Section Header"],     fallback_idx=2)
    lay_content = get_layout(prs, ["Title and Content"],  fallback_idx=1)

    display_label = {r["code"]: hiv_report_label(r["code"], r["label"], r.get("component", ""))
                     for r in interventions}

    # Helper: truncated display label for chart rows
    def _dlbl(r):
        return _trunc(display_label.get(r["code"], r["label"]))

    items = sorted(results.values(), key=lambda r: r.get("s_base", 0), reverse=True)
    n_inv = stats["n_inv"]
    n_exp = stats["n_experts"]

    sorted_inv = sorted(results.values(), key=lambda r: r.get("composite", 0), reverse=True)
    for idx, r in enumerate(sorted_inv, 1):
        r["display_idx"] = idx
    med_gate = statistics.median([r.get("gate_mean", 0) for r in sorted_inv]) if sorted_inv else 0
    med_imp  = statistics.median([r.get("avg_impact", 0) for r in sorted_inv]) if sorted_inv else 0

    def _tier(r):
        hg = r.get("gate_mean", 0) > med_gate
        hi = r.get("avg_impact", 0) > med_imp
        if hg and hi:  return "alta"
        if hg or hi:   return "media"
        return "baixa"

    def _tier_num(r):
        t = _tier(r)
        return 1 if t == "alta" else (2 if t == "media" else 3)

    def _group_summary_rows(group_key):
        buckets = OrderedDict()
        for r in results.values():
            gname = str(r.get(group_key, "") or "Outros").strip() or "Outros"
            buckets.setdefault(gname, []).append(r)
        rows = []
        for gname, vals in buckets.items():
            n     = len(vals)
            spond = sum(x.get("s_pond", 0)      for x in vals) / n if n else 0
            imp   = sum(x.get("avg_impact", 0)  for x in vals) / n if n else 0
            eff   = sum(x.get("e_score", 0)     for x in vals) / n if n else 0
            g1    = sum(1 for x in vals if _tier(x) == "alta")
            g2    = sum(1 for x in vals if _tier(x) == "media")
            g3    = sum(1 for x in vals if _tier(x) == "baixa")
            rows.append((gname, n, f"{spond:.3f}", f"{imp:.2f}", f"{eff*100:.1f}%",
                         g1 or "—", g2 or "—", g3 or "—"))
        rows.sort(key=lambda x: float(x[2]), reverse=True)
        return rows

    # ── 1: Title ──
    add_title_slide_from_layout(
        prs, lay_title,
        "Delphi W1 — Relatório de Resultados",
        (f"{n_exp} especialistas  ·  {n_inv} intervenções  ·  "
         f"Taxa de resposta mediana: {stats['rr_median']}%  ·  "
         f"{datetime.now().strftime('%d/%m/%Y')}"),
    )

    # ── 2: Executive summary ──
    sum_slide = prs.slides.add_slide(lay_content)
    _set_title_ph(sum_slide, "Resumo Executivo")
    cl, ct, cw, ch = _content_bounds(sum_slide, prs)
    metrics = [
        ("Especialistas",       str(n_exp)),
        ("Intervenções",        str(n_inv)),
        ("≥80% optimizáveis",   str(stats["n_inv_80"])),
        ("Impacto alto (≥2.5)", str(stats["n_imp_high"])),
        ("Consenso unânime",    str(stats["n_unanimous"])),
        ("Taxa resp. mediana",  f"{stats['rr_median']}%"),
    ]
    bw = cw // len(metrics)
    bh = _In(2.4)
    for i, (lbl, val) in enumerate(metrics):
        x = cl + i * bw
        y = ct + (ch - bh) // 2
        _rect(sum_slide, x + _In(0.05), y, bw - _In(0.1), bh,
              fill=_C["white"], line=_C["border"])
        _text(sum_slide, val, x + _In(0.05), y + _In(0.2), bw - _In(0.1), _In(1.1),
              size=32, bold=True, color=_C["accent2"], align=_PA.CENTER)
        _text(sum_slide, lbl, x + _In(0.05), y + _In(1.4), bw - _In(0.1), _In(0.75),
              size=10, color=_C["muted"], align=_PA.CENTER)

    # ── 3: Response rate table ──
    rr_rows = [
        ("", str(d.get("label", "")),
         f"{d.get('pct', 0)}%", f"{d.get('n', 0)}/{n_exp}")
        for d in stats.get("rr_detail", [])
    ]
    if rr_rows:
        add_table_slide_from_layout(
            prs, lay_content, "Taxa de Resposta por Intervenção",
            ["Código", "Intervenção", "Taxa", "Respostas"],
            rr_rows, style=style,
            col_widths=[0.9, 5.5, 0.9, 1.3],
        )

    # ── 3b: Team response rate summary (only when Team column present in catalog) ──
    rr_team_rows = stats.get("rr_by_team", [])
    if rr_team_rows:
        team_tbl_rows = [
            (
                f"Equipa {d['team']}",
                str(d["n_inv"]),
                str(d.get("n_experts", "—")),
                f"{d['rr_median']}%",
                f"{d['rr_min']}%",
                f"{d['rr_max']}%",
            )
            for d in rr_team_rows
        ]
        add_table_slide_from_layout(
            prs, lay_content, "Taxa de Resposta por Equipa",
            ["Equipa", "Intervenções", "Especialistas", "Mediana", "Mín", "Máx"],
            team_tbl_rows, style=style,
            col_widths=[1.2, 1.1, 1.2, 1.1, 0.9, 0.9],
        )

    # ── 4: Program / component summaries ──
    for tbl_title, group_key, col_label in [
        ("Resumo por Programa",   "programa",   "Programa"),
        ("Resumo por Componente", "comp_macro", "Componente"),
    ]:
        group_rows = _group_summary_rows(group_key)
        if group_rows:
            add_table_slide_from_layout(
                prs, lay_content, tbl_title,
                [col_label, "N", "S_pond médio", "Impacto médio",
                 "Índice eficiência", "G1", "G2", "G3"],
                group_rows, style=style,
            )

    # ── Section: Univariate analysis ──
    add_section_slide_from_layout(prs, lay_section, "Análise Univariada",
                                  f"{n_inv} intervenções  ·  {n_exp} especialistas")

    gate_agg    = univariate.get("gate_agg", {})
    # Apply top-N cap to all ranked chart lists
    n_total_inv  = len(items)
    sorted_gate  = univariate.get("sorted_gate",   items)[:top_n]
    sorted_imp   = univariate.get("sorted_impact", items)[:top_n]
    sorted_dup   = univariate.get("sorted_dup",    items)[:top_n]
    sorted_intg  = univariate.get("sorted_intg",   items)[:top_n]
    sorted_res   = univariate.get("sorted_res",    items)[:top_n]
    imp_counts  = univariate.get("imp_counts", {})
    exp_counts  = univariate.get("exp_counts", {})

    # 2a: Gate distribution — donut + stacked bars with truncated labels
    gate_stacked_rows = [
        (_dlbl(r), [r["n_sim"], r["n_poss"], r["n_nao"]])
        for r in sorted_gate
    ]
    gate_donut_svg = svg_donut(
        gate_agg or {"Sim def.": 1, "Possiv.": 1, "Não": 1},
        ["#2e7d52", "#d4a017", "#e0e0e0"], size=280,
    )
    add_two_charts_slide_from_layout(
        prs, lay_content, "2a · Distribuição Triagem (Optimizabilidade)",
        _svg_to_png(gate_donut_svg),
        _svg_to_png(svg_hbar_stacked(gate_stacked_rows, width=580, label_w=280)),
        "Distribuição global",
        "Por intervenção" + (f"  ·  " + _trunc_note(len(sorted_gate), n_total_inv)
                             if _trunc_note(len(sorted_gate), n_total_inv) else ""),
    )

    # 2b: Impact distribution
    max_imp = max((r["avg_impact"] for r in items), default=3.0) or 3.0
    impact_rows = [(_dlbl(r), r["avg_impact"], max_imp) for r in sorted_imp]
    if imp_counts and any(imp_counts.values()):
        imp_labels = {1: "Baixo", 2: "Médio", 3: "Alto"}
        imp_donut_svg = svg_donut(
            {imp_labels.get(k, str(k)): v for k, v in sorted(imp_counts.items()) if v},
            ["#b91c1c", "#b45309", "#2e7d52"], size=280,
        )
        imp_donut_png = _svg_to_png(imp_donut_svg)
    else:
        imp_donut_png = None
    add_two_charts_slide_from_layout(
        prs, lay_content, "2b · Impacto Esperado",
        imp_donut_png,
        _svg_to_png(svg_hbar_single(impact_rows, width=580, label_w=280)),
        "Distribuição global",
        "Impacto médio por intervenção" + (f"  ·  " + _trunc_note(len(sorted_imp), n_total_inv)
                                           if _trunc_note(len(sorted_imp), n_total_inv) else ""),
    )

    # 2c: Expertise distribution
    if exp_counts and any(exp_counts.values()):
        exp_labels = {1: "Baixa", 2: "Média", 3: "Alta"}
        exp_donut_svg = svg_donut(
            {exp_labels.get(k, str(k)): v for k, v in sorted(exp_counts.items()) if v},
            ["#f59e0b", "#3b82f6", "#2e7d52"], size=340,
        )
        add_visual_slide_from_layout(prs, lay_content,
                                     "2c · Nível de Experiência dos Especialistas",
                                     _svg_to_png(exp_donut_svg))

    # 2d–2f: Duplication / Integration / Resource reduction
    for slide_title, field, src in [
        ("2d · % Sobreposição com Outras Intervenções", "dup_pct",  sorted_dup),
        ("2e · % Potencial de Integração",              "intg_pct", sorted_intg),
        ("2f · % Possibilidade de Redução de Recursos", "res_pct",  sorted_res),
    ]:
        bar_rows = [(_dlbl(r), r[field], 100) for r in src]
        add_visual_slide_from_layout(
            prs, lay_content, slide_title,
            _svg_to_png(svg_hbar_single(bar_rows, width=900, label_w=340)),
            note=_trunc_note(len(src), n_total_inv),
        )

    # 2g: Scatter plot — embed SVG legend panel (HTML legend stripped; legend
    #     re-built as native SVG elements to the right of the plot)
    if include_xyplot:
        scatter_raw = svg_scatter_optim_impact_exp(items)
        svg_end = scatter_raw.find("</svg>")
        scatter_body = scatter_raw[:svg_end] if svg_end != -1 else scatter_raw
        # Build compact SVG legend appended inside the <svg> tag
        # The scatter panel_w is 700; total svg width is 1060; legend uses 700..1060
        scatter_legend = (
            '<rect x="714" y="30" width="330" height="130" rx="4" fill="#f8fafc" stroke="#e2e8f0" stroke-width="1"/>'
            '<text x="724" y="50" font-size="10" font-weight="600" fill="#334155">Cor = % redução de recursos</text>'
            '<defs><linearGradient id="lg1" x1="0" x2="1" y1="0" y2="0">'
            '<stop offset="0%" stop-color="rgb(226,232,240)"/>'
            '<stop offset="100%" stop-color="rgb(91,94,166)"/>'
            '</linearGradient></defs>'
            '<rect x="724" y="56" width="150" height="10" rx="2" fill="url(#lg1)" stroke="#94a3b8" stroke-width="0.5"/>'
            '<text x="724" y="80" font-size="9" fill="#64748b">0%</text>'
            '<text x="791" y="80" text-anchor="middle" font-size="9" fill="#64748b">50%</text>'
            '<text x="876" y="80" text-anchor="end" font-size="9" fill="#64748b">100%</text>'
            '<text x="724" y="100" font-size="10" font-weight="600" fill="#334155">Tamanho = experiência média</text>'
            '<circle cx="738" cy="116" r="4" fill="#cbd5e1" stroke="#64748b" stroke-width="0.7"/>'
            '<text x="746" y="120" font-size="9" fill="#374151">1.0</text>'
            '<circle cx="774" cy="116" r="7" fill="#cbd5e1" stroke="#64748b" stroke-width="0.7"/>'
            '<text x="784" y="120" font-size="9" fill="#374151">2.0</text>'
            '<circle cx="816" cy="116" r="10" fill="#cbd5e1" stroke="#64748b" stroke-width="0.7"/>'
            '<text x="828" y="120" font-size="9" fill="#374151">3.0</text>'
            '<text x="724" y="142" font-size="10" font-weight="600" fill="#334155">Zonas de cutoff</text>'
            '<rect x="724" y="148" width="10" height="10" fill="#dcfce7" stroke="#86efac" stroke-width="0.7"/>'
            '<text x="737" y="157" font-size="9" fill="#374151">Alta</text>'
            '<rect x="760" y="148" width="10" height="10" fill="#fef3c7" stroke="#fcd34d" stroke-width="0.7"/>'
            '<text x="773" y="157" font-size="9" fill="#374151">Média</text>'
            '<rect x="808" y="148" width="10" height="10" fill="#fee2e2" stroke="#fca5a5" stroke-width="0.7"/>'
            '<text x="821" y="157" font-size="9" fill="#374151">Baixa</text>'
        )
        scatter_svg_pptx = scatter_body + scatter_legend + "</svg>"
        add_visual_slide_from_layout(
            prs, lay_content,
            "2g · Dispersão: Optimizabilidade × Impacto Esperado",
            _svg_to_png(scatter_svg_pptx, scale=1.3),
        )

    # ── Section: Scoring ──
    if include_scoring:
        add_section_slide_from_layout(prs, lay_section, "Pontuação e Metodologia")
        add_text_slide_from_layout(prs, lay_content, "Fórmulas de Pontuação", [
            "",
            "S_base  =  média(triagem_score_i)  ×  média(impacto_i)",
            "",
            "S_pond  =  Σ( gate_i × impacto_i × exp_i )  /  Σ( exp_i )",
            "",
            "Onde:",
            "   triagem_score :  Sim definitivamente = 1   ·   Possivelmente = 0.5   ·   Não = 0",
            "   impacto    :  Baixo = 1   ·   Médio = 2   ·   Alto = 3",
            "   exp (ponderação) :  Baixa = 1   ·   Média = 2   ·   Alta = 3",
        ], font_size=14)

        s_max = max(
            max((r["s_base"] for r in items), default=1.0),
            max((r["s_pond"] for r in items), default=1.0),
        ) or 1.0
        sbase_rows = [(_dlbl(r), r["s_base"], s_max)
                      for r in sorted(items, key=lambda r: r["s_base"], reverse=True)[:top_n]]
        spond_rows = [(_dlbl(r), r["s_pond"], s_max)
                      for r in sorted(items, key=lambda r: r["s_pond"], reverse=True)[:top_n]]
        add_two_charts_slide_from_layout(
            prs, lay_content, "Pontuações S_base e S_pond",
            _svg_to_png(svg_hbar_single(sbase_rows, width=560, color="#1a5276", label_w=280)),
            _svg_to_png(svg_hbar_single(spond_rows, width=560, color="#2e7d52", label_w=280)),
            "S_base (não ponderado)"
            + (f"  ·  " + _trunc_note(len(sbase_rows), n_total_inv)
               if _trunc_note(len(sbase_rows), n_total_inv) else ""),
            "S_pond (ponderado por experiência)",
        )
        alluvial_items = items[:top_n]
        add_visual_slide_from_layout(
            prs, lay_content,
            "Transição de Rankings: Optimizabilidade → S_base → S_pond",
            _svg_to_png(svg_alluvial_weighting(alluvial_items), scale=1.2),
            note=_trunc_note(len(alluvial_items), n_total_inv),
        )

        if simple_table:
            score_rows = [
                (_dlbl(r),
                 f"{r.get('s_pond', 0):.3f}",
                 f"{r.get('e_score', 0)*100:.1f}%",
                 str(_tier_num(r)))
                for r in sorted(results.values(), key=lambda x: x.get("s_pond", 0), reverse=True)
            ]
            add_table_slide_from_layout(
                prs, lay_content, "Tabela de Pontuações (Simples)",
                ["Intervenção", "S_pond", "Índice de eficiência",
                 "Grupo Próximos Passos (1/2/3)"],
                score_rows, style=style,
                col_widths=[5.0, 1.0, 1.5, 1.0],
            )
        else:
            score_rows = [
                (r.get("rank_base", "—"),
                 _dlbl(r),
                 _trunc(r.get("component", "") or "—"),
                 f"{r['s_base']:.3f}", f"{r['s_pond']:.3f}",
                 (f"+{r.get('rank_delta',0)}"
                  if r.get("rank_delta", 0) > 0 else str(r.get("rank_delta", 0))),
                 f"{r.get('avg_impact',0):.2f}",
                 f"{r.get('exp_mean',0):.2f}")
                for r in sorted(items, key=lambda r: r.get("rank_base", 999))
            ]
            add_table_slide_from_layout(
                prs, lay_content, "Tabela de Pontuações",
                ["#", "Intervenção", "Componente",
                 "S_base", "S_pond", "Δrank", "Impacto", "Exp."],
                score_rows, style=style,
                col_widths=[0.4, 4.0, 1.6, 0.8, 0.8, 0.6, 0.8, 0.7],
            )

    # ── Section: Priority ranking ──
    add_section_slide_from_layout(prs, lay_section, "Tabela de Priorização",
                                  "Ordenado por S_base (pontuação composta)")
    priority_rows = [
        (r.get("rank_base", "—"),
         _dlbl(r),
         _trunc(r.get("component", "") or "—"),
         f"{r['pct_optimizable']}%",
         f"{r.get('avg_impact',0):.1f}",
         f"{r['dup_pct']}%",
         f"{r['intg_pct']}%")
        for r in sorted(items, key=lambda r: r.get("rank_base", 999))
    ]
    add_table_slide_from_layout(
        prs, lay_content, "Ranking de Prioridades",
        ["#", "Intervenção", "Componente",
         "% Optim.", "Impacto", "% Dup.", "% Integr."],
        priority_rows, style=style,
        col_widths=[0.4, 4.2, 1.6, 0.9, 0.8, 0.8, 0.9],
    )

    # ── Section: Next steps ──
    alta  = [display_label.get(r["code"], r["label"])
             for r in sorted_inv if _tier(r) == "alta"]
    media = [display_label.get(r["code"], r["label"])
             for r in sorted_inv if _tier(r) == "media"]
    baixa = [display_label.get(r["code"], r["label"])
             for r in sorted_inv if _tier(r) == "baixa"]

    add_section_slide_from_layout(prs, lay_section, "Próximos Passos",
                                  "Classificação por prioridade")
    for tier_title, tier_items in [
        ("Alta Prioridade — Optimizar e reforçar",   alta),
        ("Prioridade Média — Avaliar caso a caso",    media),
        ("Baixa Prioridade — Manter ou desinvestir",  baixa),
    ]:
        if tier_items:
            add_text_slide_from_layout(prs, lay_content, tier_title, tier_items)

    # ── Section: Intervention cards ──
    add_section_slide_from_layout(prs, lay_section, "Fichas por Componente",
                                  "Resumo detalhado por intervenção")

    from collections import OrderedDict as _OD
    by_comp = _OD()
    for intv in interventions:
        comp = intv.get("component", "") or "Outros"
        by_comp.setdefault(comp, []).append(intv)

    def _fmt_tagged_codes(codes):
        if not codes:
            return "—"
        return " · ".join(display_label.get(c, c) for c in codes[:6])

    def _fmt_count_map(count_map, limit=8):
        if not count_map:
            return "—"
        pairs = sorted(count_map.items(), key=lambda kv: (-kv[1], str(kv[0]).lower()))[:limit]
        return "; ".join(f"{k} ({v})" for k, v in pairs)

    _playwright_warned = False
    for comp, comp_intvs in by_comp.items():
        if not comp_intvs:
            continue
        add_section_slide_from_layout(prs, lay_section,
                                      f"Componente: {comp}", "Detalhe por intervenção")
        for intv in comp_intvs:
            r = results.get(intv["code"])
            if not r:
                continue
            comments   = r.get("comments", []) or []
            card_title = display_label.get(intv["code"], intv["label"])
            card_html  = _build_card_html_fragment(
                card_title, r, _tier_num(r), _tier(r),
                _fmt_tagged_codes(r.get("top_dup", [])),
                _fmt_tagged_codes(r.get("top_intg", [])),
                _fmt_count_map(r.get("dup_other_counts", {}), limit=8),
                _fmt_count_map(r.get("intg_other_counts", {}), limit=8),
                comments,
                style=style,
                all_items=list(results.values()),
            )
            card_png = render_card_html_to_png(card_html, width_px=1280, height_px=720)
            if card_png:
                add_card_slide_from_layout(prs, lay_content, card_png)
            else:
                if not _playwright_warned:
                    if not _PLAYWRIGHT_OK:
                        print("  ⚠ Playwright not installed — using legacy card layout.")
                        print("    Install: pip install playwright")
                        print("    macOS (Chrome preferred): brew install --cask google-chrome")
                        print("    Then:   PLAYWRIGHT_SKIP_BROWSER_DOWNLOAD=1 python -m playwright install")
                        print("    Or:     python -m playwright install chromium")
                    else:
                        print("  ⚠ No usable browser for Playwright — using legacy card layout.")
                    _playwright_warned = True
                _slide_intervention_card(
                    prs, card_title, r, _tier_num(r), _tier(r),
                    _fmt_tagged_codes(r.get("top_dup", [])),
                    _fmt_tagged_codes(r.get("top_intg", [])),
                    _fmt_count_map(r.get("dup_other_counts", {}), limit=8),
                    _fmt_count_map(r.get("intg_other_counts", {}), limit=8),
                    comments,
                )

    return prs


# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────

def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(0)

    results_path = sys.argv[1]
    dict_path    = None
    output_dir   = None
    exclude_codes_arg = None
    exclude_file = None
    config_path = None
    simple_sections = False
    simple_table = False
    generate_pptx = False
    pptx_template     = None
    pptx_engine       = "template"
    pptx_cli_overrides = {}
    alluvial_top = None
    priority_order = "s_pond"
    logo_path = None

    # Parse optional arguments
    i = 2
    while i < len(sys.argv):
      arg = sys.argv[i]

      if arg in ("--output-dir", "-o"):
        if i + 1 >= len(sys.argv):
          sys.exit("Erro: --output-dir requer um valor.")
        output_dir = sys.argv[i + 1]
        i += 2
        continue
      if arg.startswith("--output-dir="):
        output_dir = arg.split("=", 1)[1]
        i += 1
        continue

      if arg in ("--exclude-experts", "-x"):
        if i + 1 >= len(sys.argv):
          sys.exit("Erro: --exclude-experts requer um valor.")
        exclude_codes_arg = sys.argv[i + 1]
        i += 2
        continue
      if arg.startswith("--exclude-experts="):
        exclude_codes_arg = arg.split("=", 1)[1]
        i += 1
        continue

      if arg in ("--exclude-file", "-X"):
        if i + 1 >= len(sys.argv):
          sys.exit("Erro: --exclude-file requer um caminho.")
        exclude_file = sys.argv[i + 1]
        i += 2
        continue
      if arg.startswith("--exclude-file="):
        exclude_file = arg.split("=", 1)[1]
        i += 1
        continue

      if arg == "--config":
        if i + 1 >= len(sys.argv):
          sys.exit("Erro: --config requer um caminho.")
        config_path = sys.argv[i + 1]
        i += 2
        continue
      if arg.startswith("--config="):
        config_path = arg.split("=", 1)[1]
        i += 1
        continue

      if arg in ("--simple-sections", "--compact-report"):
        simple_sections = True
        i += 1
        continue

      if arg == "--simple-table":
        simple_table = True
        i += 1
        continue

      if arg == "--pptx":
        generate_pptx = True
        i += 1
        continue

      if arg in ("--pptx-template",):
        if i + 1 >= len(sys.argv):
          sys.exit("Erro: --pptx-template requer um caminho.")
        pptx_template = sys.argv[i + 1]
        i += 2
        continue
      if arg.startswith("--pptx-template="):
        pptx_template = arg.split("=", 1)[1]
        i += 1
        continue

      if arg in ("--pptx-engine",):
        if i + 1 >= len(sys.argv):
          sys.exit("Erro: --pptx-engine requer um valor (legacy ou template).")
        eng = sys.argv[i + 1].lower()
        if eng not in ("legacy", "template"):
          sys.exit("Erro: --pptx-engine deve ser 'legacy' ou 'template'.")
        pptx_engine = eng
        i += 2
        continue
      if arg.startswith("--pptx-engine="):
        eng = arg.split("=", 1)[1].lower()
        if eng not in ("legacy", "template"):
          sys.exit("Erro: --pptx-engine deve ser 'legacy' ou 'template'.")
        pptx_engine = eng
        i += 1
        continue

      # ── PPTX style overrides ──────────────────────────────────────────────
      _PPTX_CLI_ARGS = {
        "--pptx-chart-top-n":           "chart.top_n",
        "--pptx-chart-label-max-chars": "chart.label.max_chars",
        "--pptx-chart-font-family":     "chart.font.family",
        "--pptx-chart-font-size":       "chart.font.size",
        "--pptx-table-font-family":     "table.font.family",
        "--pptx-table-header-font-size":"table.header.font.size",
        "--pptx-table-body-font-size":  "table.body.font.size",
        "--pptx-card-font-family":      "card.font.family",
        "--pptx-card-title-font-size":  "card.title.font.size",
        "--pptx-card-body-font-size":   "card.body.font.size",
      }
      _matched_pptx = False
      for _cli_flag, _style_key in _PPTX_CLI_ARGS.items():
        if arg == _cli_flag:
          if i + 1 >= len(sys.argv):
            sys.exit(f"Erro: {_cli_flag} requer um valor.")
          pptx_cli_overrides[_style_key] = sys.argv[i + 1]
          i += 2
          _matched_pptx = True
          break
        if arg.startswith(_cli_flag + "="):
          pptx_cli_overrides[_style_key] = arg.split("=", 1)[1]
          i += 1
          _matched_pptx = True
          break
      if _matched_pptx:
        continue

      if arg in ("--logo",):
        if i + 1 >= len(sys.argv):
          sys.exit("Erro: --logo requer um caminho para a imagem.")
        logo_path = sys.argv[i + 1]
        i += 2
        continue
      if arg.startswith("--logo="):
        logo_path = arg.split("=", 1)[1]
        i += 1
        continue

      if arg in ("--alluvial-top",):
        if i + 1 >= len(sys.argv):
          sys.exit("Erro: --alluvial-top requer um número inteiro.")
        try:
            alluvial_top = int(sys.argv[i + 1])
        except ValueError:
            sys.exit("Erro: --alluvial-top requer um número inteiro.")
        i += 2
        continue
      if arg.startswith("--alluvial-top="):
        try:
            alluvial_top = int(arg.split("=", 1)[1])
        except ValueError:
            sys.exit("Erro: --alluvial-top requer um número inteiro.")
        i += 1
        continue

      if arg == "--priority-order":
        if i + 1 >= len(sys.argv):
          sys.exit("Erro: --priority-order requer um valor (s_base ou s_pond).")
        pord = sys.argv[i + 1].lower()
        if pord not in ("s_base", "s_pond"):
            sys.exit("Erro: --priority-order deve ser 's_base' ou 's_pond'.")
        priority_order = pord
        i += 2
        continue
      if arg.startswith("--priority-order="):
        pord = arg.split("=", 1)[1].lower()
        if pord not in ("s_base", "s_pond"):
            sys.exit("Erro: --priority-order deve ser 's_base' ou 's_pond'.")
        priority_order = pord
        i += 1
        continue

      if arg.startswith("-"):
        print(f"Aviso: argumento não reconhecido (ignorado): {arg}")
        i += 1
        continue

      if dict_path is None:
        dict_path = arg
      else:
        print(f"Aviso: argumento posicional extra ignorado: {arg}")
      i += 1

    # Default output directory: ./reports relative to current working directory
    if output_dir is None:
        output_dir = "reports"
    output_dir = os.path.abspath(output_dir)

    if not os.path.exists(results_path):
        sys.exit(f"Erro: ficheiro de resultados não encontrado: {results_path}")
    if dict_path and not os.path.exists(dict_path):
        print(f"Aviso: dicionário não encontrado: {dict_path} — continuando sem ele.")
        dict_path = None

    # Load config to get program name and other settings
    cfg = _load_simple_config(config_path)
    programa_nome = cfg.get("TOPIC_LABEL", "Programa Nacional de Controlo do HIV/SIDA")
    topic_label = cfg.get("TOPIC_LABEL", "HIV")
    if logo_path is None and cfg.get("REPORT_LOGO"):
        logo_path = cfg["REPORT_LOGO"]

    expected_experts = load_expected_experts()
    n_expected_experts = len(expected_experts)

    # Ensure output directory exists
    os.makedirs(output_dir, exist_ok=True)

    print(f"A carregar resultados: '{results_path}'...")
    df = load_data(results_path, "Responses")
    print(f"  {len(df)} linhas · {df['expert_code'].nunique()} especialistas únicos")
    if n_expected_experts > 0:
      print(f"  Base esperada (experts.txt): {n_expected_experts} especialistas")
    else:
      print("  Aviso: experts.txt não encontrado ou vazio; usando especialistas observados.")

    print("A carregar metadados das intervenções...")
    interventions = load_metadata(results_path, dict_path)
    if interventions is None:
        df_codes = sorted(df["intervention"].dropna().unique())
        interventions = [{"code": c, "label": c, "programa": "", "component": "", "comp_macro": "Outros", "url": ""} for c in df_codes]
        print(f"  Metadados: códigos brutos da folha Responses ({len(interventions)} intervenções)")

    programas_found = sorted({
      str(inv.get("programa", "")).strip()
      for inv in interventions
      if str(inv.get("programa", "")).strip()
    })
    componentes_found = sorted({
      str(inv.get("component", "")).strip()
      for inv in interventions
      if str(inv.get("component", "")).strip()
    })

    if programas_found:
      print(f"  Programas encontrados ({len(programas_found)}): " + ", ".join(programas_found))
    else:
      print("  Programas encontrados: nenhum valor preenchido nos metadados.")

    if componentes_found:
      print(f"  Componentes encontrados ({len(componentes_found)}): " + ", ".join(componentes_found))
    else:
      print("  Componentes encontrados: nenhum valor preenchido nos metadados.")

    excluded_codes = load_excluded_codes(
      cli_codes=exclude_codes_arg,
      exclude_file=exclude_file,
      config_path=config_path,
    )
    if excluded_codes:
      excluded_lc = {c.lower() for c in excluded_codes}

      before_rows = len(df)
      before_experts = df["expert_code"].nunique()
      before_expert_codes = {
        str(v).strip().lower()
        for v in df["expert_code"].dropna().astype(str).tolist()
        if str(v).strip()
      }

      df = df[
        ~df["expert_code"].fillna("").astype(str).str.strip().str.lower().isin(excluded_lc)
      ].copy()
      after_experts = df["expert_code"].nunique()

      matched = sorted(c for c in excluded_codes if c.lower() in before_expert_codes)
      unmatched = sorted(c for c in excluded_codes if c.lower() not in before_expert_codes)

      print(f"  Exclusões activas: {len(excluded_codes)} especialista(s)")
      print(f"    Removidas {before_rows - len(df)} linha(s) de respostas")
      print(f"    Removidos {before_experts - after_experts} especialista(s)")
      if matched:
        print("    Especialistas excluídos encontrados: " + ", ".join(matched))
      if unmatched:
        print("    Aviso: especialistas não encontrados (ignorados): " + ", ".join(unmatched))

    if not interventions:
      sys.exit("Erro: não restaram intervenções após aplicar as exclusões.")

    # Read per-team expected expert counts from config (N_EXPERTS_TEAM_A … _D etc.)
    n_experts_per_team = {}
    for k, v in cfg.items():
        if k.upper().startswith("N_EXPERTS_TEAM_") and v.strip().isdigit():
            team_letter = k.upper()[len("N_EXPERTS_TEAM_"):]
            if team_letter:
                n_experts_per_team[team_letter] = int(v.strip())
    if n_experts_per_team:
        print(f"  Especialistas esperados por equipa: " +
              ", ".join(f"{t}={n}" for t, n in sorted(n_experts_per_team.items())))

    print("A agregar respostas...")
    results = aggregate(df, interventions)
    compute_ranks(results)
    stats   = summary_stats(results, df, n_expected_experts,
                            n_experts_per_team=n_experts_per_team or None)

    print(f"  {stats['n_experts']} especialistas · {stats['n_inv']} intervenções")
    print(f"  Taxa de resposta: mediana={stats['rr_median']}% "
          f"intervalo=[{stats['rr_min']}%–{stats['rr_max']}%]")

    # Console ranking table
    sorted_base = sorted(results.values(), key=lambda r: r["s_base"], reverse=True)
    print(f"\n  {'Rank':>4}  {'S_base':>7}  {'S_pond':>7}  {'Δrank':>6}  Intervenção")
    print(f"  {'----':>4}  {'-------':>7}  {'-------':>7}  {'------':>6}  -----------")
    for r in sorted_base:
        delta = r.get("rank_delta", 0)
        delta_str = f"+{delta}" if delta > 0 else str(delta)
        print(f"  {r['rank_base']:>4}  {r['s_base']:>7.3f}  {r['s_pond']:>7.3f}  {delta_str:>6}  {r['label'][:52]}")

    print("\nA calcular estatísticas univariadas...")
    univariate = univariate_analysis(results, df)

    # Load logo as base64 data URI if provided
    logo_data_uri = None
    if logo_path:
        import base64, mimetypes
        lp = os.path.abspath(logo_path)
        if not os.path.isfile(lp):
            print(f"Aviso: ficheiro de logótipo não encontrado: {lp}")
        else:
            mime = mimetypes.guess_type(lp)[0] or "image/png"
            with open(lp, "rb") as _lf:
                logo_data_uri = f"data:{mime};base64,{base64.b64encode(_lf.read()).decode()}"

    print("A gerar relatório HTML...")
    html = render_html(
      results,
      stats,
      interventions,
      results_path,
      univariate,
      include_xyplot=not simple_sections,
      include_scoring=not simple_sections,
      alluvial_top=alluvial_top,
      priority_order=priority_order,
      programa_nome=programa_nome,
      topic_label=topic_label,
      simple_table=simple_table,
      logo_data_uri=logo_data_uri,
    )

    ts       = datetime.now().strftime("%Y%m%d_%H%M")
    out_path = os.path.join(output_dir, f"delphi_w1_hiv_relatorio_{ts}.html")
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(html)

    print(f"\n✓ Relatório HTML guardado em: {out_path}")

    if generate_pptx:
        # Resolve default template path relative to this script
        _default_tpl = os.path.normpath(os.path.join(
            os.path.dirname(os.path.abspath(__file__)),
            "..", "docs", "templates", "INS_template_report_style.pptx",
        ))
        if pptx_template is None and os.path.exists(_default_tpl):
            pptx_template = _default_tpl

        if pptx_engine == "template":
            if pptx_template and not os.path.exists(pptx_template):
                sys.exit(f"Erro: ficheiro de template PPTX não encontrado: {pptx_template}")
            if not pptx_template:
                print("  ⚠ Template PPTX não encontrado; a usar motor legado.")
                pptx_engine = "legacy"

        print(f"\nA gerar apresentação PowerPoint (motor: {pptx_engine})...")
        if pptx_engine == "template":
            prs = build_pptx_template(
                results, stats, interventions, results_path,
                univariate, pptx_template,
                include_xyplot=not simple_sections,
                include_scoring=not simple_sections,
                simple_table=simple_table,
                cli_overrides=pptx_cli_overrides or None,
                cfg=cfg,
            )
        else:
            prs = build_pptx(
                results, stats, interventions, results_path,
                univariate,
                include_xyplot=not simple_sections,
                include_scoring=not simple_sections,
                simple_table=simple_table,
            )
        if prs is not None:
            pptx_path = os.path.join(output_dir, f"delphi_w1_hiv_relatorio_{ts}.pptx")
            prs.save(pptx_path)
            print(f"✓ Apresentação PowerPoint guardada em: {pptx_path}")

if __name__ == "__main__":
    main()
